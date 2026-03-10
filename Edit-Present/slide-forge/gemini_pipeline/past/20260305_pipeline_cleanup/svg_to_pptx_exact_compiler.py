#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
High-fidelity SVG -> PPTX compiler.

This compiler prioritizes visual consistency over editability:
1) Rasterize SVG to PNG with CairoSVG.
2) Create a PPTX slide with the same canvas size.
3) Place the rendered PNG as a full-slide background image.

Result: the PPTX should look as close as possible to the SVG rendering.
"""

from __future__ import annotations

import argparse
import re
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional, Tuple

import cairosvg
from pptx import Presentation
from pptx.util import Inches


DEFAULT_DPI = 96.0


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Compile SVG to visually-identical PPTX.")
    parser.add_argument(
        "--input",
        "-i",
        default="/Users/xiaoxiaobo/Downloads/sjtuwenber/gemini_pipeline/output/svg/image.svg",
        help="Input SVG file path.",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="/Users/xiaoxiaobo/Downloads/sjtuwenber/gemini_pipeline/output/pptx/image_exact.pptx",
        help="Output PPTX file path.",
    )
    parser.add_argument(
        "--dpi",
        type=float,
        default=DEFAULT_DPI,
        help=f"SVG pixel DPI for slide size mapping (default: {DEFAULT_DPI}).",
    )
    parser.add_argument(
        "--render-scale",
        type=float,
        default=1.0,
        help="Raster render scale multiplier (default: 1.0).",
    )
    parser.add_argument(
        "--keep-rendered-png",
        action="store_true",
        help="Keep the intermediate rendered PNG next to output PPTX.",
    )
    return parser.parse_args()


def _parse_svg_length_to_px(raw: Optional[str], dpi: float) -> Optional[float]:
    if not raw:
        return None
    s = raw.strip()
    m = re.match(r"^([+-]?\d*\.?\d+)([a-zA-Z%]*)$", s)
    if not m:
        return None
    value = float(m.group(1))
    unit = m.group(2).lower()

    if unit in ("", "px"):
        return value
    if unit == "in":
        return value * dpi
    if unit == "pt":
        return value * dpi / 72.0
    if unit == "pc":
        return value * dpi / 6.0
    if unit == "cm":
        return value * dpi / 2.54
    if unit == "mm":
        return value * dpi / 25.4
    # Percent/em/ex unsupported without context; caller may fallback to viewBox.
    return None


def read_svg_canvas_size(svg_path: Path, dpi: float) -> Tuple[float, float]:
    tree = ET.parse(svg_path)
    root = tree.getroot()

    width_px = _parse_svg_length_to_px(root.get("width"), dpi)
    height_px = _parse_svg_length_to_px(root.get("height"), dpi)

    if width_px and height_px and width_px > 0 and height_px > 0:
        return width_px, height_px

    view_box = (root.get("viewBox") or "").strip()
    if view_box:
        parts = view_box.replace(",", " ").split()
        if len(parts) == 4:
            try:
                vb_w = float(parts[2])
                vb_h = float(parts[3])
                if vb_w > 0 and vb_h > 0:
                    return vb_w, vb_h
            except ValueError:
                pass

    raise ValueError(f"Cannot determine SVG canvas size: {svg_path}")


def compile_svg_to_pptx_exact(
    input_svg: Path,
    output_pptx: Path,
    dpi: float = DEFAULT_DPI,
    render_scale: float = 1.0,
    keep_rendered_png: bool = False,
) -> Path:
    width_px, height_px = read_svg_canvas_size(input_svg, dpi)
    out_w = max(1, int(round(width_px * render_scale)))
    out_h = max(1, int(round(height_px * render_scale)))

    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    if keep_rendered_png:
        rendered_png = output_pptx.with_suffix(".rendered.png")
    else:
        rendered_png = Path(tempfile.NamedTemporaryFile(delete=False, suffix=".png").name)

    cairosvg.svg2png(
        url=str(input_svg),
        write_to=str(rendered_png),
        output_width=out_w,
        output_height=out_h,
    )

    prs = Presentation()
    prs.slide_width = Inches(width_px / dpi)
    prs.slide_height = Inches(height_px / dpi)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(rendered_png),
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    prs.save(str(output_pptx))

    if not keep_rendered_png and rendered_png.exists():
        rendered_png.unlink(missing_ok=True)

    return output_pptx


def main() -> None:
    args = parse_args()
    input_svg = Path(args.input).expanduser().resolve()
    output_pptx = Path(args.output).expanduser().resolve()

    if not input_svg.exists():
        raise FileNotFoundError(f"Input SVG not found: {input_svg}")

    result = compile_svg_to_pptx_exact(
        input_svg=input_svg,
        output_pptx=output_pptx,
        dpi=args.dpi,
        render_scale=args.render_scale,
        keep_rendered_png=args.keep_rendered_png,
    )
    print(f"Compiled successfully: {result}")


if __name__ == "__main__":
    main()
