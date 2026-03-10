#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SVG to PPTX Pro: Convert SVG elements to editable PPTX shapes.

Unlike svg_to_pptx_slide.py which renders SVG as a background image,
this script converts each SVG element (rect, circle, path, image, text, etc.)
to native PPTX shapes, preserving editability and layer order.
"""

import argparse
import base64
import hashlib
import io
import json
import os
import re
import tempfile
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple, Union

DEFAULT_RUNTIME_CONFIG_PATH = (
    Path(__file__).resolve().parent.parent / "config" / "runtime_api_config.json"
)

try:
    from .constants import (
        CT_NS,
        DEFAULT_DPI,
        FREEFORM_LOCAL_UNITS,
        REL_IMAGE,
        REL_NS,
        SUPPORTED_IMAGE_ASPECT_RATIOS,
        SVG_BLIP_EXT_URI,
        SVG_BLIP_NS,
        SVG_NS,
        TRANSPARENT_PNG_BYTES,
        XLINK_NS,
    )
    from .pptx.coordinates import CoordinateConverter
    from .pptx.chart_codegen import (
        CHART_CODE_PROMPT as _CHART_CODE_PROMPT_IMPL,
        execute_chart_code as _execute_chart_code_impl,
        generate_chart_code as _generate_chart_code_impl,
    )
    from .io.config_loader import (
        load_config as _load_config_impl,
        load_placeholders as _load_placeholders_impl,
    )
    from .io.profile_loader import (
        DEFAULT_PROFILE as _DEFAULT_PROFILE_IMPL,
        SUPPORTED_PROFILES as _SUPPORTED_PROFILES_IMPL,
        apply_profile_overrides as _apply_profile_overrides_impl,
        load_profile_spec as _load_profile_spec_impl,
    )
    from .pptx.image_placeholder import (
        add_image_placeholder as _add_image_placeholder_impl,
        apply_transform_chain as _apply_transform_chain_impl,
        classify_image_placeholder_group as _classify_image_placeholder_group_impl,
        clip_path_is_rect as _clip_path_is_rect_impl,
        extract_image_placeholders as _extract_image_placeholders_impl,
        parse_matrix_simple as _parse_matrix_simple_impl,
        parse_scale_simple as _parse_scale_simple_impl,
        rasterize_clipped_placeholder as _rasterize_clipped_placeholder_impl,
    )
    from .pptx.grouping import extract_visual_group_meta as _extract_visual_group_meta_impl
    from .pptx.path_renderer import add_svg_path as render_svg_path
    from .pptx.svg_blip import (
        add_svgblip_region_picture as _add_svgblip_region_picture_impl,
        build_svg_region_snippet as _build_svg_region_snippet_impl,
        count_group_graphics as _count_group_graphics_impl,
        elem_is_hidden_bbox_rect as _elem_is_hidden_bbox_rect_impl,
        group_local_bbox_from_hidden_rect as _group_local_bbox_from_hidden_rect_impl,
        inject_svg_blips_into_pptx as _inject_svg_blips_into_pptx_impl,
        should_render_group_as_svgblip as _should_render_group_as_svgblip_impl,
        transform_bbox as _transform_bbox_impl,
    )
    from .pptx.image_refill import (
        _iter_genai_parts as _iter_genai_parts_impl,
        apply_background_alpha as _apply_background_alpha_impl,
        crop_placeholder_from_source as _crop_placeholder_from_source_impl,
        estimate_border_bg_color as _estimate_border_bg_color_impl,
        generate_placeholder_image as _generate_placeholder_image_impl,
        guess_image_mime as _guess_image_mime_impl,
        parse_aspect_ratio as _parse_aspect_ratio_impl,
        pick_supported_aspect_ratio as _pick_supported_aspect_ratio_impl,
        rect_intersection_area as _rect_intersection_area_impl,
        redraw_placeholder_crop_without_text as _redraw_placeholder_crop_without_text_impl,
        resolve_source_image_for_svg as _resolve_source_image_for_svg_impl,
    )
    from .pptx.semantic_text import (
        add_semantic_text_items_absolute as _add_semantic_text_items_absolute_impl,
        add_semantic_textbox as _add_semantic_textbox_impl,
        extract_semantic_textboxes as _extract_semantic_textboxes_impl,
    )
    from .pptx.shape_style import (
        apply_fill_to_shape as _apply_fill_to_shape_impl,
        apply_stroke_to_shape as _apply_stroke_to_shape_impl,
    )
    from .pptx.text_helpers import (
        assemble_line_text as _assemble_line_text_impl,
        baseline_to_top_offset_px as _baseline_to_top_offset_px_impl,
        estimate_text_width_px as _estimate_text_width_px_impl,
        group_text_lines as _group_text_lines_impl,
        read_text_content as _read_text_content_impl,
        should_insert_space as _should_insert_space_impl,
    )
    from .pptx.text_style import (
        apply_text_run_style,
        font_family_is_theme,
        pick_font_name,
        set_run_ea_font,
        set_run_font_size_from_px,
    )
    from .utils.colors import parse_color
    from .utils.effects import (
        apply_svg_filter_shadow_if_needed as _apply_svg_filter_shadow_if_needed_impl,
        extract_simple_drop_shadow_filters as _extract_simple_drop_shadow_filters_impl,
    )
    from .utils.lengths import parse_length, parse_opacity
    from .utils.style_context import StyleContext
    from .utils.svg_helpers import natural_sort_key, tag_name
    from .utils.text import CJK_RE, GENERIC_FONTS
    from .utils.transform_extract import (
        parse_transform_rotation as _parse_transform_rotation_impl,
        parse_transform_xy as _parse_transform_xy_impl,
    )
    from .utils.transforms import TransformMatrix, parse_transform
except ImportError:
    from compiler.constants import (
        CT_NS,
        DEFAULT_DPI,
        FREEFORM_LOCAL_UNITS,
        REL_IMAGE,
        REL_NS,
        SUPPORTED_IMAGE_ASPECT_RATIOS,
        SVG_BLIP_EXT_URI,
        SVG_BLIP_NS,
        SVG_NS,
        TRANSPARENT_PNG_BYTES,
        XLINK_NS,
    )
    from compiler.pptx.coordinates import CoordinateConverter
    from compiler.pptx.chart_codegen import (
        CHART_CODE_PROMPT as _CHART_CODE_PROMPT_IMPL,
        execute_chart_code as _execute_chart_code_impl,
        generate_chart_code as _generate_chart_code_impl,
    )
    from compiler.io.config_loader import (
        load_config as _load_config_impl,
        load_placeholders as _load_placeholders_impl,
    )
    from compiler.io.profile_loader import (
        DEFAULT_PROFILE as _DEFAULT_PROFILE_IMPL,
        SUPPORTED_PROFILES as _SUPPORTED_PROFILES_IMPL,
        apply_profile_overrides as _apply_profile_overrides_impl,
        load_profile_spec as _load_profile_spec_impl,
    )
    from compiler.pptx.image_placeholder import (
        add_image_placeholder as _add_image_placeholder_impl,
        apply_transform_chain as _apply_transform_chain_impl,
        classify_image_placeholder_group as _classify_image_placeholder_group_impl,
        clip_path_is_rect as _clip_path_is_rect_impl,
        extract_image_placeholders as _extract_image_placeholders_impl,
        parse_matrix_simple as _parse_matrix_simple_impl,
        parse_scale_simple as _parse_scale_simple_impl,
        rasterize_clipped_placeholder as _rasterize_clipped_placeholder_impl,
    )
    from compiler.pptx.grouping import extract_visual_group_meta as _extract_visual_group_meta_impl
    from compiler.pptx.path_renderer import add_svg_path as render_svg_path
    from compiler.pptx.svg_blip import (
        add_svgblip_region_picture as _add_svgblip_region_picture_impl,
        build_svg_region_snippet as _build_svg_region_snippet_impl,
        count_group_graphics as _count_group_graphics_impl,
        elem_is_hidden_bbox_rect as _elem_is_hidden_bbox_rect_impl,
        group_local_bbox_from_hidden_rect as _group_local_bbox_from_hidden_rect_impl,
        inject_svg_blips_into_pptx as _inject_svg_blips_into_pptx_impl,
        should_render_group_as_svgblip as _should_render_group_as_svgblip_impl,
        transform_bbox as _transform_bbox_impl,
    )
    from compiler.pptx.image_refill import (
        _iter_genai_parts as _iter_genai_parts_impl,
        apply_background_alpha as _apply_background_alpha_impl,
        crop_placeholder_from_source as _crop_placeholder_from_source_impl,
        estimate_border_bg_color as _estimate_border_bg_color_impl,
        generate_placeholder_image as _generate_placeholder_image_impl,
        guess_image_mime as _guess_image_mime_impl,
        parse_aspect_ratio as _parse_aspect_ratio_impl,
        pick_supported_aspect_ratio as _pick_supported_aspect_ratio_impl,
        rect_intersection_area as _rect_intersection_area_impl,
        redraw_placeholder_crop_without_text as _redraw_placeholder_crop_without_text_impl,
        resolve_source_image_for_svg as _resolve_source_image_for_svg_impl,
    )
    from compiler.pptx.semantic_text import (
        add_semantic_text_items_absolute as _add_semantic_text_items_absolute_impl,
        add_semantic_textbox as _add_semantic_textbox_impl,
        extract_semantic_textboxes as _extract_semantic_textboxes_impl,
    )
    from compiler.pptx.shape_style import (
        apply_fill_to_shape as _apply_fill_to_shape_impl,
        apply_stroke_to_shape as _apply_stroke_to_shape_impl,
    )
    from compiler.pptx.text_helpers import (
        assemble_line_text as _assemble_line_text_impl,
        baseline_to_top_offset_px as _baseline_to_top_offset_px_impl,
        estimate_text_width_px as _estimate_text_width_px_impl,
        group_text_lines as _group_text_lines_impl,
        read_text_content as _read_text_content_impl,
        should_insert_space as _should_insert_space_impl,
    )
    from compiler.pptx.text_style import (
        apply_text_run_style,
        font_family_is_theme,
        pick_font_name,
        set_run_ea_font,
        set_run_font_size_from_px,
    )
    from compiler.utils.colors import parse_color
    from compiler.utils.effects import (
        apply_svg_filter_shadow_if_needed as _apply_svg_filter_shadow_if_needed_impl,
        extract_simple_drop_shadow_filters as _extract_simple_drop_shadow_filters_impl,
    )
    from compiler.utils.lengths import parse_length, parse_opacity
    from compiler.utils.style_context import StyleContext
    from compiler.utils.svg_helpers import natural_sort_key, tag_name
    from compiler.utils.text import CJK_RE, GENERIC_FONTS
    from compiler.utils.transform_extract import (
        parse_transform_rotation as _parse_transform_rotation_impl,
        parse_transform_xy as _parse_transform_xy_impl,
    )
    from compiler.utils.transforms import TransformMatrix, parse_transform

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

# Optional: openai for chart generation
try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

# Chart data types
try:
    from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
except ImportError:
    CategoryChartData = None
    XyChartData = None
    BubbleChartData = None

try:
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    XL_CHART_TYPE = None

# Optional: PIL for image cropping
try:
    from PIL import Image, ImageChops, ImageDraw, ImageFilter, ImageStat
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# Optional: Google GenAI for image placeholder refill
try:
    from google import genai as google_genai
    from google.genai import types as google_genai_types
    HAS_GOOGLE_GENAI = True
except Exception:
    google_genai = None
    google_genai_types = None
    HAS_GOOGLE_GENAI = False

# Optional: CairoSVG for rasterizing non-rect clip paths
try:
    import cairosvg
    HAS_CAIROSVG = True
except Exception:
    # If CairoSVG or its native cairo dependency is missing/not loadable,
    # gracefully disable rasterization instead of crashing.
    HAS_CAIROSVG = False

# ============================================================================
# Constants and Regex Patterns
# ============================================================================

# Regex patterns

# Active per-slide filter mapping (filter_id -> simple drop-shadow spec).
_ACTIVE_DROP_SHADOW_FILTERS: Dict[str, Dict[str, float]] = {}


# ============================================================================
# Utility Functions
# ============================================================================


def extract_simple_drop_shadow_filters(root: ET.Element) -> Dict[str, Dict[str, float]]:
    return _extract_simple_drop_shadow_filters_impl(root)


def apply_svg_filter_shadow_if_needed(shape: Any, elem: ET.Element, converter: "CoordinateConverter") -> None:
    _apply_svg_filter_shadow_if_needed_impl(shape, elem, converter, _ACTIVE_DROP_SHADOW_FILTERS)


def _to_local_coord(value: float, min_val: float, span: float, local_units: int) -> int:
    """Map absolute coordinate to freeform local coordinate with rounding."""
    if span <= 0:
        return 0
    return int(round((value - min_val) / span * local_units))


def _elem_is_hidden_bbox_rect(elem: ET.Element) -> bool:
    return _elem_is_hidden_bbox_rect_impl(elem)


def _group_local_bbox_from_hidden_rect(elem: ET.Element) -> Optional[Tuple[float, float, float, float]]:
    return _group_local_bbox_from_hidden_rect_impl(elem)


def _transform_bbox(mat: "TransformMatrix", x: float, y: float, w: float, h: float) -> Tuple[float, float, float, float]:
    return _transform_bbox_impl(mat, x, y, w, h)


def _count_group_graphics(elem: ET.Element) -> Dict[str, int]:
    return _count_group_graphics_impl(elem)


def should_render_group_as_svgblip(elem: ET.Element) -> bool:
    return _should_render_group_as_svgblip_impl(elem)


def build_svg_region_snippet(root: ET.Element, x: float, y: float, w: float, h: float) -> bytes:
    return _build_svg_region_snippet_impl(root, x, y, w, h)


def add_svgblip_region_picture(
    slide: Any,
    left_px: float,
    top_px: float,
    width_px: float,
    height_px: float,
    converter: "CoordinateConverter",
    shape_name: str,
) -> Optional[Any]:
    return _add_svgblip_region_picture_impl(
        slide=slide,
        left_px=left_px,
        top_px=top_px,
        width_px=width_px,
        height_px=height_px,
        converter=converter,
        shape_name=shape_name,
    )


def parse_transform_rotation(transform: Optional[str]) -> Optional[float]:
    return _parse_transform_rotation_impl(transform)


def parse_transform_xy(transform: Optional[str]) -> Tuple[float, float]:
    return _parse_transform_xy_impl(transform)


def read_text_content(elem: ET.Element) -> str:
    return _read_text_content_impl(elem)


def group_text_lines(items: List[Dict[str, object]], line_tol: float) -> List[Dict[str, object]]:
    return _group_text_lines_impl(items, line_tol)


def should_insert_space(prev_item: Dict[str, object], curr_item: Dict[str, object]) -> bool:
    return _should_insert_space_impl(prev_item, curr_item)


def assemble_line_text(items: List[Dict[str, object]]) -> str:
    return _assemble_line_text_impl(items)


def estimate_text_width_px(text: str, font_size: float, letter_spacing: float = 0.0) -> float:
    return _estimate_text_width_px_impl(text, font_size, letter_spacing)


def baseline_to_top_offset_px(font_size: float, text: str) -> float:
    return _baseline_to_top_offset_px_impl(font_size, text)


def parse_aspect_ratio(aspect_ratio: str) -> Optional[float]:
    return _parse_aspect_ratio_impl(aspect_ratio)


def pick_supported_aspect_ratio(width: float, height: float) -> str:
    return _pick_supported_aspect_ratio_impl(width, height)


def rect_intersection_area(a: Dict[str, Any], b: Dict[str, Any]) -> float:
    return _rect_intersection_area_impl(a, b)


def _iter_genai_parts(response: Any) -> List[Any]:
    return _iter_genai_parts_impl(response)


def generate_placeholder_image(
    caption: str,
    width: float,
    height: float,
    image_refill_config: Dict[str, Any],
    remove_bg: Optional[bool] = None,
    remove_bg_mode: Optional[str] = None,
) -> Tuple[Optional[Path], bool, Optional[str]]:
    return _generate_placeholder_image_impl(
        caption,
        width,
        height,
        image_refill_config,
        remove_bg=remove_bg,
        remove_bg_mode=remove_bg_mode,
    )


def guess_image_mime(path: Path) -> str:
    return _guess_image_mime_impl(path)


def redraw_placeholder_crop_without_text(
    seed_image: Path,
    caption: str,
    width: float,
    height: float,
    image_refill_config: Dict[str, Any],
) -> Tuple[Optional[Path], bool, Optional[str]]:
    return _redraw_placeholder_crop_without_text_impl(
        seed_image, caption, width, height, image_refill_config
    )


def resolve_source_image_for_svg(svg_path: Path, image_refill_config: Dict[str, Any]) -> Optional[Path]:
    return _resolve_source_image_for_svg_impl(svg_path, image_refill_config)


def estimate_border_bg_color(img: Image.Image, border: int = 2) -> Tuple[int, int, int]:
    return _estimate_border_bg_color_impl(img, border=border)


def apply_background_alpha(
    img_rgba: Image.Image,
    bg_color: Tuple[int, int, int],
    threshold: float,
    feather: float,
) -> Image.Image:
    return _apply_background_alpha_impl(img_rgba, bg_color, threshold, feather)


def crop_placeholder_from_source(
    source_image: Path,
    placeholder: Dict[str, Any],
    svg_width: float,
    svg_height: float,
    image_refill_config: Dict[str, Any],
    semantic_textboxes: Optional[List[Dict[str, Any]]] = None,
) -> Tuple[Optional[Path], bool, Optional[str]]:
    return _crop_placeholder_from_source_impl(
        source_image,
        placeholder,
        svg_width,
        svg_height,
        image_refill_config,
        semantic_textboxes,
    )


# ============================================================================
# Image Placeholder Extraction (Correct Transform Chain Handling)
# ============================================================================

def parse_matrix_simple(transform: str) -> Tuple[float, float, float, float, float, float]:
    return _parse_matrix_simple_impl(transform)


def apply_transform_chain(
    x: float,
    y: float,
    w: float,
    h: float,
    transforms: List[str],
) -> Tuple[float, float, float, float]:
    return _apply_transform_chain_impl(x, y, w, h, transforms)


def clip_path_is_rect(clip_elem: ET.Element) -> bool:
    return _clip_path_is_rect_impl(clip_elem)


def rasterize_clipped_placeholder(
    placeholder: Dict[str, Any],
    converter: "CoordinateConverter",
) -> Optional[Tuple[str, Tuple[int, int, int, int]]]:
    return _rasterize_clipped_placeholder_impl(placeholder, converter)


def parse_scale_simple(transform: str) -> Tuple[float, float]:
    return _parse_scale_simple_impl(transform)


def classify_image_placeholder_group(elem: ET.Element) -> str:
    return _classify_image_placeholder_group_impl(elem)


def extract_image_placeholders(svg_path: Path) -> List[Dict[str, Any]]:
    return _extract_image_placeholders_impl(svg_path)


# ============================================================================
# Config and Placeholders Loading
# ============================================================================

def load_config(config_path: Optional[Path]) -> Dict[str, Any]:
    return _load_config_impl(config_path)


def load_placeholders(json_path: Optional[Path]) -> Dict[Tuple[str, str], Dict[str, Any]]:
    return _load_placeholders_impl(json_path)


def load_profile_spec(profile: str, profile_dir: Optional[Path]) -> Dict[str, Any]:
    return _load_profile_spec_impl(profile, profile_dir=profile_dir)


def apply_profile_overrides(
    base_config: Dict[str, Any],
    profile_spec: Dict[str, Any],
    section: str,
) -> Dict[str, Any]:
    return _apply_profile_overrides_impl(base_config, profile_spec, section)


def _cli_flag_present(flag: str) -> bool:
    return any(arg == flag or arg.startswith(f"{flag}=") for arg in os.sys.argv[1:])


# ============================================================================
# Chart Generation (from LLM)
# ============================================================================

CHART_CODE_PROMPT = _CHART_CODE_PROMPT_IMPL


def generate_chart_code(
    caption: str,
    chart_spec: Optional[Any] = None,
    api_key: Optional[str] = None,
    base_url: Optional[str] = None,
    model: Optional[str] = None,
    max_tokens: int = 5000,
    temperature: float = 0.2,
    error_hint: Optional[str] = None,
) -> Optional[str]:
    return _generate_chart_code_impl(
        caption=caption,
        chart_spec=chart_spec,
        api_key=api_key,
        base_url=base_url,
        model=model,
        max_tokens=max_tokens,
        temperature=temperature,
        error_hint=error_hint,
    )


def execute_chart_code(
    code: str,
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
) -> Tuple[bool, Optional[str]]:
    return _execute_chart_code_impl(
        code=code,
        slide=slide,
        left=left,
        top=top,
        width=width,
        height=height,
    )


# ============================================================================
# Shape Addition Functions
# ============================================================================

def apply_fill_to_shape(shape, fill_color: Optional[str], opacity: float = 1.0) -> None:
    _apply_fill_to_shape_impl(shape=shape, fill_color=fill_color, opacity=opacity)


def apply_stroke_to_shape(shape, stroke_color: Optional[str], stroke_width: Optional[str],
                          converter: CoordinateConverter, opacity: float = 1.0,
                          stroke_linecap: Optional[str] = None,
                          stroke_linejoin: Optional[str] = None) -> None:
    _apply_stroke_to_shape_impl(
        shape=shape,
        stroke_color=stroke_color,
        stroke_width=stroke_width,
        converter=converter,
        opacity=opacity,
        stroke_linecap=stroke_linecap,
        stroke_linejoin=stroke_linejoin,
    )


def add_svg_rect(slide, elem: ET.Element, transform: TransformMatrix,
                 style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG rect element as PPTX rectangle shape."""
    x = parse_length(elem.get("x"), 0.0)
    y = parse_length(elem.get("y"), 0.0)
    width = parse_length(elem.get("width"), 0.0)
    height = parse_length(elem.get("height"), 0.0)
    rx = parse_length(elem.get("rx"), 0.0)
    ry = parse_length(elem.get("ry"), 0.0)
    
    if width <= 0 or height <= 0:
        return None
    
    # Apply transform to corners
    x1, y1 = transform.transform_point(x, y)
    x2, y2 = transform.transform_point(x + width, y + height)
    
    # Get final bounding box
    left = min(x1, x2)
    top = min(y1, y2)
    final_width = abs(x2 - x1)
    final_height = abs(y2 - y1)
    
    # Convert to EMU
    left_emu = converter.to_emu_x(left)
    top_emu = converter.to_emu_y(top)
    width_emu = converter.to_emu_width(final_width)
    height_emu = converter.to_emu_height(final_height)
    
    # Choose shape type based on rounded corners
    if rx > 0 or ry > 0:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE
    
    shape = slide.shapes.add_shape(
        shape_type, left_emu, top_emu, width_emu, height_emu
    )
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(
        shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity,
        stroke_linecap=(style.stroke_linecap or elem.get("stroke-linecap")),
        stroke_linejoin=(style.stroke_linejoin or elem.get("stroke-linejoin")),
    )
    apply_svg_filter_shadow_if_needed(shape, elem, converter)
    
    # Apply rotation if present
    rotation = transform.get_rotation_degrees()
    if abs(rotation) > 0.1:
        shape.rotation = rotation
    
    return shape


def add_svg_circle(slide, elem: ET.Element, transform: TransformMatrix,
                   style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG circle element as PPTX oval shape."""
    cx = parse_length(elem.get("cx"), 0.0)
    cy = parse_length(elem.get("cy"), 0.0)
    r = parse_length(elem.get("r"), 0.0)
    
    if r <= 0:
        return None
    
    # Transform center and radius
    center_x, center_y = transform.transform_point(cx, cy)
    sx, sy = transform.get_scale()
    final_rx = r * sx
    final_ry = r * sy
    
    # Calculate bounding box
    left = center_x - final_rx
    top = center_y - final_ry
    width = final_rx * 2
    height = final_ry * 2
    
    # Convert to EMU
    left_emu = converter.to_emu_x(left)
    top_emu = converter.to_emu_y(top)
    width_emu = converter.to_emu_width(width)
    height_emu = converter.to_emu_height(height)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left_emu, top_emu, width_emu, height_emu
    )
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(
        shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity,
        stroke_linecap=(style.stroke_linecap or elem.get("stroke-linecap")),
        stroke_linejoin=(style.stroke_linejoin or elem.get("stroke-linejoin")),
    )
    apply_svg_filter_shadow_if_needed(shape, elem, converter)
    
    return shape


def add_svg_ellipse(slide, elem: ET.Element, transform: TransformMatrix,
                    style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG ellipse element as PPTX oval shape."""
    cx = parse_length(elem.get("cx"), 0.0)
    cy = parse_length(elem.get("cy"), 0.0)
    rx = parse_length(elem.get("rx"), 0.0)
    ry = parse_length(elem.get("ry"), 0.0)
    
    if rx <= 0 or ry <= 0:
        return None
    
    # Transform center and radii
    center_x, center_y = transform.transform_point(cx, cy)
    sx, sy = transform.get_scale()
    final_rx = rx * sx
    final_ry = ry * sy
    
    # Calculate bounding box
    left = center_x - final_rx
    top = center_y - final_ry
    width = final_rx * 2
    height = final_ry * 2
    
    # Convert to EMU
    left_emu = converter.to_emu_x(left)
    top_emu = converter.to_emu_y(top)
    width_emu = converter.to_emu_width(width)
    height_emu = converter.to_emu_height(height)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left_emu, top_emu, width_emu, height_emu
    )
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(
        shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity,
        stroke_linecap=(style.stroke_linecap or elem.get("stroke-linecap")),
        stroke_linejoin=(style.stroke_linejoin or elem.get("stroke-linejoin")),
    )
    apply_svg_filter_shadow_if_needed(shape, elem, converter)
    
    # Apply rotation if present
    rotation = transform.get_rotation_degrees()
    if abs(rotation) > 0.1:
        shape.rotation = rotation
    
    return shape


def add_svg_line(slide, elem: ET.Element, transform: TransformMatrix,
                 style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG line element as PPTX connector."""
    x1 = parse_length(elem.get("x1"), 0.0)
    y1 = parse_length(elem.get("y1"), 0.0)
    x2 = parse_length(elem.get("x2"), 0.0)
    y2 = parse_length(elem.get("y2"), 0.0)
    
    # Transform endpoints
    tx1, ty1 = transform.transform_point(x1, y1)
    tx2, ty2 = transform.transform_point(x2, y2)
    
    # Convert to EMU
    x1_emu = converter.to_emu_x(tx1)
    y1_emu = converter.to_emu_y(ty1)
    x2_emu = converter.to_emu_x(tx2)
    y2_emu = converter.to_emu_y(ty2)
    
    # Add as connector shape
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        x1_emu, y1_emu, x2_emu, y2_emu
    )
    
    # Apply stroke style
    stroke = style.stroke or elem.get("stroke") or "#000000"
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    color = parse_color(stroke)
    if color:
        connector.line.color.rgb = color
    if stroke_width:
        width_px = parse_length(stroke_width, 1.0)
        connector.line.width = Pt(width_px * 0.75)
    
    return connector


def add_svg_polygon(slide, elem: ET.Element, transform: TransformMatrix,
                    style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG polygon element as PPTX freeform shape."""
    points_str = elem.get("points", "")
    if not points_str:
        return None
    
    # Parse points
    points = []
    coords = re.findall(r"[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?", points_str)
    for i in range(0, len(coords) - 1, 2):
        try:
            x = float(coords[i])
            y = float(coords[i + 1])
            points.append((x, y))
        except ValueError:
            continue
    
    if len(points) < 3:
        return None
    
    # Transform all points
    transformed = [transform.transform_point(x, y) for x, y in points]
    
    # Find bounding box
    xs = [p[0] for p in transformed]
    ys = [p[1] for p in transformed]
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    width = max_x - min_x
    height = max_y - min_y
    
    if width <= 0 or height <= 0:
        return None
    
    # Convert bounding box to EMU
    left_emu = converter.to_emu_x(min_x)
    top_emu = converter.to_emu_y(min_y)
    
    # Calculate scale: local units to EMU (higher precision to reduce jagged edges)
    local_units = FREEFORM_LOCAL_UNITS
    scale_x = converter.to_emu_width(width) / local_units
    scale_y = converter.to_emu_height(height) / local_units
    
    # Create freeform shape using FreeformBuilder
    # Normalize points to local coordinate system (0 to local_units)
    start_local_x = _to_local_coord(transformed[0][0], min_x, width, local_units)
    start_local_y = _to_local_coord(transformed[0][1], min_y, height, local_units)
    
    builder = slide.shapes.build_freeform(start_local_x, start_local_y, scale=(scale_x, scale_y))
    
    # Draw lines to other points
    line_segments = []
    for px, py in transformed[1:]:
        local_x = _to_local_coord(px, min_x, width, local_units)
        local_y = _to_local_coord(py, min_y, height, local_units)
        line_segments.append((local_x, local_y))
    
    if line_segments:
        builder.add_line_segments(line_segments, close=True)
    
    shape = builder.convert_to_shape(left_emu, top_emu)
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(
        shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity,
        stroke_linecap=(style.stroke_linecap or elem.get("stroke-linecap")),
        stroke_linejoin=(style.stroke_linejoin or elem.get("stroke-linejoin")),
    )
    apply_svg_filter_shadow_if_needed(shape, elem, converter)
    
    return shape


def add_svg_polyline(slide, elem: ET.Element, transform: TransformMatrix,
                     style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG polyline element as PPTX freeform shape."""
    points_str = elem.get("points", "")
    if not points_str:
        return None
    
    # Parse points
    points = []
    coords = re.findall(r"[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?", points_str)
    for i in range(0, len(coords) - 1, 2):
        try:
            x = float(coords[i])
            y = float(coords[i + 1])
            points.append((x, y))
        except ValueError:
            continue
    
    if len(points) < 2:
        return None
    
    # Transform all points
    transformed = [transform.transform_point(x, y) for x, y in points]
    
    # Find bounding box
    xs = [p[0] for p in transformed]
    ys = [p[1] for p in transformed]
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    width = max_x - min_x or 1
    height = max_y - min_y or 1
    
    # Convert bounding box to EMU
    left_emu = converter.to_emu_x(min_x)
    top_emu = converter.to_emu_y(min_y)
    
    # Calculate scale: local units to EMU (higher precision to reduce jagged edges)
    local_units = FREEFORM_LOCAL_UNITS
    scale_x = converter.to_emu_width(width) / local_units
    scale_y = converter.to_emu_height(height) / local_units
    
    # Normalize points to local coordinate system
    start_local_x = _to_local_coord(transformed[0][0], min_x, width, local_units)
    start_local_y = _to_local_coord(transformed[0][1], min_y, height, local_units)
    
    builder = slide.shapes.build_freeform(start_local_x, start_local_y, scale=(scale_x, scale_y))
    
    # Draw lines to other points (not closed)
    line_segments = []
    for px, py in transformed[1:]:
        local_x = _to_local_coord(px, min_x, width, local_units)
        local_y = _to_local_coord(py, min_y, height, local_units)
        line_segments.append((local_x, local_y))
    
    if line_segments:
        builder.add_line_segments(line_segments, close=False)
    
    shape = builder.convert_to_shape(left_emu, top_emu)
    
    # Polyline typically has no fill
    fill = style.fill or elem.get("fill") or "none"
    stroke = style.stroke or elem.get("stroke") or "#000000"
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(
        shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity,
        stroke_linecap=(style.stroke_linecap or elem.get("stroke-linecap")),
        stroke_linejoin=(style.stroke_linejoin or elem.get("stroke-linejoin")),
    )
    apply_svg_filter_shadow_if_needed(shape, elem, converter)
    
    return shape


def add_svg_path(slide, elem: ET.Element, transform: TransformMatrix,
                 style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    return render_svg_path(
        slide=slide,
        elem=elem,
        transform=transform,
        style=style,
        converter=converter,
        apply_fill_to_shape=apply_fill_to_shape,
        apply_stroke_to_shape=apply_stroke_to_shape,
        apply_svg_filter_shadow_if_needed=apply_svg_filter_shadow_if_needed,
    )

def add_svg_image(slide, elem: ET.Element, transform: TransformMatrix,
                  style: StyleContext, converter: CoordinateConverter,
                  svg_path: Optional[Path] = None) -> Optional[Any]:
    """Add SVG image element as PPTX picture."""
    x = parse_length(elem.get("x"), 0.0)
    y = parse_length(elem.get("y"), 0.0)
    width = parse_length(elem.get("width"), 0.0)
    height = parse_length(elem.get("height"), 0.0)
    
    # Get href (try both with and without namespace)
    href = elem.get(f"{{{XLINK_NS}}}href") or elem.get("href", "")
    
    if not href or width <= 0 or height <= 0:
        return None
    
    # Transform position
    tx, ty = transform.transform_point(x, y)
    sx, sy = transform.get_scale()
    final_width = width * sx
    final_height = height * sy
    
    # Convert to EMU
    left_emu = converter.to_emu_x(tx)
    top_emu = converter.to_emu_y(ty)
    width_emu = converter.to_emu_width(final_width)
    height_emu = converter.to_emu_height(final_height)
    
    # Handle data URI
    if href.startswith("data:"):
        # Parse data URI: data:[<mediatype>][;base64],<data>
        match = re.match(r"data:([^;,]+)?(?:;base64)?,(.+)", href, re.DOTALL)
        if not match:
            return None
        
        media_type = match.group(1) or "image/png"
        data = match.group(2)
        
        try:
            image_data = base64.b64decode(data)
        except Exception:
            return None
        
        # Create temp file with correct extension
        ext = ".png"
        if "jpeg" in media_type or "jpg" in media_type:
            ext = ".jpg"
        elif "gif" in media_type:
            ext = ".gif"
        elif "webp" in media_type:
            ext = ".webp"
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(image_data)
            tmp_path = tmp.name
        
        try:
            shape = slide.shapes.add_picture(
                tmp_path, left_emu, top_emu, width_emu, height_emu
            )
        finally:
            Path(tmp_path).unlink(missing_ok=True)
        
        return shape
    
    # Handle external file reference
    if svg_path:
        image_path = svg_path.parent / href
        if image_path.exists():
            shape = slide.shapes.add_picture(
                str(image_path), left_emu, top_emu, width_emu, height_emu
            )
            return shape
    
    return None


def add_svg_text(slide, elem: ET.Element, transform: TransformMatrix,
                 style: StyleContext, converter: CoordinateConverter,
                 cjk_font: str = "PingFang SC") -> Optional[Any]:
    """Add SVG text element as PPTX textbox."""
    # Get text content
    text_content = ""
    if elem.text:
        text_content += elem.text
    for child in elem:
        if tag_name(child) == "tspan":
            if child.text:
                text_content += child.text
            if child.tail:
                text_content += child.tail
    
    text_content = text_content.replace("\u00a0", " ").strip()
    if not text_content:
        return None
    
    # Get position
    x = parse_length(elem.get("x"), 0.0)
    y = parse_length(elem.get("y"), 0.0)
    
    # Transform position
    tx, ty = transform.transform_point(x, y)
    
    # Get font size
    font_size_str = style.font_size or elem.get("font-size", "16")
    font_size = parse_length(font_size_str, 16.0)
    sx, sy = transform.get_scale()
    final_font_size = font_size * (sx + sy) / 2
    
    # Estimate text dimensions (rough approximation)
    char_width = final_font_size * 0.6
    text_width = len(text_content) * char_width
    text_height = final_font_size * 1.5
    
    # Adjust position based on text-anchor
    anchor = style.text_anchor or elem.get("text-anchor", "start")
    if anchor == "middle":
        tx -= text_width / 2
    elif anchor == "end":
        tx -= text_width
    
    # Adjust y position (SVG text y is baseline, PPTX is top)
    ty -= final_font_size * 0.8
    
    # Convert to EMU
    left_emu = converter.to_emu_x(tx)
    top_emu = converter.to_emu_y(ty)
    width_emu = converter.to_emu_width(text_width * 1.2)  # Add some padding
    height_emu = converter.to_emu_height(text_height)
    
    # Create textbox
    shape = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text_content
    
    # Set alignment
    if anchor == "middle":
        p.alignment = PP_ALIGN.CENTER
    elif anchor == "end":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    
    # Apply font style
    if p.runs:
        run = p.runs[0]
        set_run_font_size_from_px(run, final_font_size, scale=converter.font_scale_factor())
        
        # Set font family
        font_family = style.font_family or elem.get("font-family")
        use_theme_font = font_family_is_theme(font_family)
        if font_family:
            # Clean up font family string
            fonts = [f.strip().strip("\"'") for f in font_family.split(",")]
            for font in fonts:
                if font.lower() not in GENERIC_FONTS and "msfontservice" not in font.lower():
                    run.font.name = font
                    break
            else:
                # Use CJK font if text contains CJK characters
                if CJK_RE.search(text_content) and cjk_font and not use_theme_font:
                    run.font.name = cjk_font
        if CJK_RE.search(text_content) and not use_theme_font:
            ea_font = run.font.name or cjk_font
            if ea_font:
                set_run_ea_font(run, ea_font)
        
        # Set color
        fill = style.fill or elem.get("fill")
        color = parse_color(fill)
        if color:
            run.font.color.rgb = color
        
        # Set bold
        font_weight = style.font_weight or elem.get("font-weight", "")
        if font_weight in ("bold", "700", "800", "900"):
            run.font.bold = True
    
    return shape


# ============================================================================
# Semantic Layer Textbox Extraction and Addition
# ============================================================================

def extract_semantic_textboxes(
    svg_path: Path,
) -> Tuple[Dict[str, float], List[Dict[str, object]], Dict[str, Dict[str, object]]]:
    return _extract_semantic_textboxes_impl(svg_path)


def add_semantic_textbox(
    slide,
    tb: Dict[str, object],
    converter: CoordinateConverter,
    line_tol: float,
    box_pad: float,
    cjk_font: str,
    width_expand: float = 1.0,
) -> Optional[Any]:
    return _add_semantic_textbox_impl(
        slide=slide,
        tb=tb,
        converter=converter,
        line_tol=line_tol,
        box_pad=box_pad,
        cjk_font=cjk_font,
        width_expand=width_expand,
    )


def add_semantic_text_items_absolute(
    slide,
    tb: Dict[str, object],
    converter: CoordinateConverter,
    cjk_font: str,
    text_pad: float = 1.0,
) -> List[Any]:
    return _add_semantic_text_items_absolute_impl(
        slide=slide,
        tb=tb,
        converter=converter,
        cjk_font=cjk_font,
        text_pad=text_pad,
    )


# ============================================================================
# Image Placeholder Processing (with Correct Position Calculation)
# ============================================================================

def add_image_placeholder(
    slide,
    placeholder: Dict[str, Any],
    converter: CoordinateConverter,
    svg_path: Optional[Path] = None,
    clip_to_canvas: bool = True,
) -> Optional[Any]:
    return _add_image_placeholder_impl(
        slide=slide,
        placeholder=placeholder,
        converter=converter,
        svg_path=svg_path,
        clip_to_canvas=clip_to_canvas,
    )


# ============================================================================
# Main SVG Processing
# ============================================================================

def process_svg_element(slide, elem: ET.Element, parent_transform: TransformMatrix,
                        parent_style: StyleContext, converter: CoordinateConverter,
                        svg_path: Optional[Path] = None, cjk_font: str = "PingFang SC",
                        skip_elements: Optional[set] = None,
                        skip_image_placeholders: bool = True,
                        current_visual_group: Optional[str] = None,
                        visual_group_shapes: Optional[Dict[str, List[Any]]] = None,
                        visual_group_order: Optional[Dict[str, float]] = None,
                        svg_blip_context: Optional[Dict[str, Any]] = None) -> List[Any]:
    """Process a single SVG element and its children, return list of shapes added.
    
    Args:
        skip_image_placeholders: If True, skip elements with data-role="image-placeholder"
                                 (they are processed separately with correct transform handling)
    """
    shapes = []
    tag = tag_name(elem)
    
    # Skip certain elements
    skip_elements = skip_elements or set()
    elem_id = elem.get("id", "")
    if elem_id in skip_elements:
        return shapes
    
    # Skip defs, clipPath, mask, etc.
    if tag in ("defs", "clipPath", "mask", "symbol", "use", "metadata", "title", "desc"):
        return shapes

    # Skip helper bbox rects used only for semantic/visual grouping metadata.
    if tag == "rect" and (elem.get("class") or "").strip() in ("tb-bbox", "vg-bbox"):
        return shapes
    
    # Skip semantic layer (handled separately if needed)
    if elem.get("data-type") == "semantic-layer" or elem_id == "semantic-layer":
        return shapes
    
    # Skip image placeholders if requested (they are processed separately)
    if skip_image_placeholders and (
        elem.get("data-role") == "image-placeholder"
        or elem.get("data-type") == "image-placeholder"
    ):
        return shapes

    data_visual_group = (elem.get("data-visual-group") or "").strip()
    next_visual_group = current_visual_group
    if data_visual_group:
        next_visual_group = data_visual_group
    elif tag == "g" and elem.get("data-type") == "visual-group":
        next_visual_group = (elem.get("id") or elem.get("data-visual-group") or "").strip() or current_visual_group
        if visual_group_order is not None and next_visual_group:
            raw_order = (elem.get("data-order") or "").strip()
            order_val: Optional[float] = None
            if raw_order:
                try:
                    order_val = float(raw_order)
                except Exception:
                    order_val = None
            if next_visual_group not in visual_group_order:
                # Keep discovery order for groups without explicit data-order.
                visual_group_order[next_visual_group] = (
                    order_val if order_val is not None else float(1000000 + len(visual_group_order))
                )
            elif order_val is not None:
                visual_group_order[next_visual_group] = min(visual_group_order[next_visual_group], order_val)
    
    # Update transform
    elem_transform_str = elem.get("transform", "")
    elem_transform = parse_transform(elem_transform_str)
    current_transform = parent_transform.multiply(elem_transform)
    
    # Update style
    current_style = parent_style.update_from_element(elem)

    # Hybrid fidelity mode: render complex icon groups as svgBlip pictures.
    if (
        tag == "g"
        and svg_blip_context
        and bool(svg_blip_context.get("enabled"))
        and should_render_group_as_svgblip(elem)
    ):
        local_bbox = _group_local_bbox_from_hidden_rect(elem)
        if local_bbox is not None:
            gx, gy, gw, gh = _transform_bbox(current_transform, *local_bbox)
            if gw > 0 and gh > 0:
                counter = int(svg_blip_context.get("counter", 0)) + 1
                svg_blip_context["counter"] = counter
                shape_name = f"svgblip-{svg_blip_context.get('slide_index', 1)}-{counter}"
                shape = add_svgblip_region_picture(
                    slide, gx, gy, gw, gh, converter, shape_name
                )
                if shape is not None:
                    root = svg_blip_context.get("root")
                    jobs = svg_blip_context.get("jobs")
                    if isinstance(root, ET.Element) and isinstance(jobs, list):
                        snippet = build_svg_region_snippet(root, gx, gy, gw, gh)
                        jobs.append(
                            {
                                "slide_index": int(svg_blip_context.get("slide_index", 1)),
                                "shape_name": shape_name,
                                "svg_bytes": snippet,
                            }
                        )
                    shapes.append(shape)
                    group_for_shape = data_visual_group or next_visual_group
                    if visual_group_shapes is not None and group_for_shape:
                        visual_group_shapes.setdefault(group_for_shape, []).append(shape)
                    return shapes
    
    # Process based on element type
    shape = None
    
    if tag == "rect":
        shape = add_svg_rect(slide, elem, current_transform, current_style, converter)
    
    elif tag == "circle":
        shape = add_svg_circle(slide, elem, current_transform, current_style, converter)
    
    elif tag == "ellipse":
        shape = add_svg_ellipse(slide, elem, current_transform, current_style, converter)
    
    elif tag == "line":
        shape = add_svg_line(slide, elem, current_transform, current_style, converter)
    
    elif tag == "polygon":
        shape = add_svg_polygon(slide, elem, current_transform, current_style, converter)
    
    elif tag == "polyline":
        shape = add_svg_polyline(slide, elem, current_transform, current_style, converter)
    
    elif tag == "path":
        shape = add_svg_path(slide, elem, current_transform, current_style, converter)
    
    elif tag == "image":
        shape = add_svg_image(slide, elem, current_transform, current_style, converter, svg_path)
    
    elif tag == "text":
        shape = add_svg_text(slide, elem, current_transform, current_style, converter, cjk_font)
    
    elif tag == "g":
        # Process group children recursively
        for child in elem:
            child_shapes = process_svg_element(
                slide, child, current_transform, current_style, converter,
                svg_path, cjk_font, skip_elements, skip_image_placeholders,
                next_visual_group, visual_group_shapes, visual_group_order, svg_blip_context
            )
            shapes.extend(child_shapes)
    
    elif tag == "svg":
        # Nested SVG - process children
        for child in elem:
            child_shapes = process_svg_element(
                slide, child, current_transform, current_style, converter,
                svg_path, cjk_font, skip_elements, skip_image_placeholders,
                next_visual_group, visual_group_shapes, visual_group_order, svg_blip_context
            )
            shapes.extend(child_shapes)
    
    if shape:
        shapes.append(shape)
        group_for_shape = data_visual_group or next_visual_group
        if visual_group_shapes is not None and group_for_shape:
            visual_group_shapes.setdefault(group_for_shape, []).append(shape)
    
    return shapes


def extract_visual_group_meta(root: ET.Element) -> Dict[str, Dict[str, object]]:
    return _extract_visual_group_meta_impl(root)


def _sanitize_svg_amp(text: str) -> str:
    """将未转义的 & 替换为 &amp;，避免 Gemini 生成 SVG 中 '&' 导致 ParseError。"""
    return re.sub(
        r"&(?!(?:amp|lt|gt|quot|apos|#\d+|#x[0-9a-fA-F]+);)",
        "&amp;",
        text,
    )


def convert_svg_to_slide(
    prs: Presentation,
    svg_path: Path,
    dpi: float = 96.0,
    cjk_font: str = "PingFang SC",
    skip_elements: Optional[set] = None,
    line_tol: float = 2.0,
    box_pad: float = 0.0,
    semantic_mode: str = "textbox",
    placeholders_map: Optional[Dict[Tuple[str, str], Dict[str, Any]]] = None,
    chart_config: Optional[Dict[str, Any]] = None,
    skip_charts: bool = False,
    image_refill_config: Optional[Dict[str, Any]] = None,
    svg_blip_jobs: Optional[List[Dict[str, Any]]] = None,
) -> None:
    """Convert a single SVG file to a PPTX slide.
    
    Args:
        placeholders_map: Mapping from (svg_file, placeholder_id) to placeholder info
        chart_config: Configuration for chart generation (API key, model, etc.)
        skip_charts: If True, skip chart generation
        image_refill_config: Optional config to generate and backfill image placeholders
    """
    global _ACTIVE_DROP_SHADOW_FILTERS

    # Parse SVG（首次失败时尝试修复常见 XML 问题后重试）
    try:
        tree = ET.parse(svg_path)
        root = tree.getroot()
    except ET.ParseError as e:
        try:
            raw = svg_path.read_text(encoding="utf-8", errors="replace")
            raw = _sanitize_svg_amp(raw)
            root = ET.fromstring(raw)
            tree = ET.ElementTree(root)
        except Exception as exc:
            # 最终仍无法解析时，记录并跳过当前 SVG，而不是中断整个批次
            print(
                f"  !! 跳过无效 SVG {svg_path.name}: {e} | after sanitize error={exc}",
                flush=True,
            )
            return
    _ACTIVE_DROP_SHADOW_FILTERS = extract_simple_drop_shadow_filters(root)
    
    # Get SVG dimensions
    svg_width = parse_length(root.get("width"), 1920.0)
    svg_height = parse_length(root.get("height"), 1080.0)
    
    # Handle viewBox
    viewbox = root.get("viewBox", "")
    if viewbox:
        vb_parts = viewbox.split()
        if len(vb_parts) == 4:
            try:
                svg_width = float(vb_parts[2])
                svg_height = float(vb_parts[3])
            except ValueError:
                pass
    
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Create coordinate converter
    converter = CoordinateConverter(
        svg_width, svg_height,
        prs.slide_width, prs.slide_height,
        dpi
    )
    
    # Initialize transform and style
    base_transform = TransformMatrix.identity()
    base_style = StyleContext()
    
    # Build skip set - always skip semantic-layer for regular SVG element processing
    # (textboxes from semantic-layer are handled separately)
    actual_skip = set(skip_elements) if skip_elements else set()
    actual_skip.add("semantic-layer")
    
    # Check if we should skip textboxes entirely
    skip_textboxes = skip_elements and "semantic-layer" in skip_elements

    # Visual-group containers (from visual-layer/data-visual-group) for drag-together grouping.
    visual_group_meta = extract_visual_group_meta(root)
    visual_group_shapes: Dict[str, List[Any]] = {}
    visual_group_order: Dict[str, float] = {
        gid: float(meta.get("order") or 1e9) for gid, meta in visual_group_meta.items()
    }
    
    # Extract image placeholders with correct position calculation
    image_placeholders = extract_image_placeholders(svg_path)

    # Pre-extract semantic textboxes once:
    # - used for final editable text reconstruction
    # - used to erase overlapped text from source-cropped images
    textboxes: List[Dict[str, Any]] = []
    semantic_group_meta: Dict[str, Dict[str, object]] = {}
    if not skip_textboxes:
        _, textboxes, semantic_group_meta = extract_semantic_textboxes(svg_path)
    
    image_refill_config = image_refill_config or {}
    slide_index = len(prs.slides)
    svg_blip_context = {
        "enabled": bool(image_refill_config.get("hybrid_svg_blip_icons", False)),
        "root": root,
        "jobs": svg_blip_jobs if isinstance(svg_blip_jobs, list) else [],
        "slide_index": slide_index,
        "counter": 0,
    }

    # Process all children of root (except semantic layer and image placeholders)
    for child in root:
        process_svg_element(
            slide, child, base_transform, base_style, converter,
            svg_path, cjk_font, actual_skip, skip_image_placeholders=True,
            current_visual_group=None,
            visual_group_shapes=visual_group_shapes,
            visual_group_order=visual_group_order,
            svg_blip_context=svg_blip_context,
        )
    
    # Add image placeholders with correct positioning
    svg_name = svg_path.name
    svg_abs = str(svg_path.resolve())
    placeholders_map = placeholders_map or {}
    chart_config = chart_config or {}
    image_count = 0
    chart_count = 0
    generated_image_count = 0
    source_cropped_count = 0
    source_redrawn_count = 0
    suppressed_textbox_count = 0
    source_crop_regions: List[Dict[str, float]] = []
    raster_text_regions: List[Dict[str, float]] = []
    raster_text_visual_group_ids: Set[str] = set()
    chart_shapes_to_front: List[Any] = []
    prefetched_generation: Dict[str, Tuple[Optional[Path], bool, Optional[str]]] = {}

    refill_mode = str(image_refill_config.get("mode") or "gemini").strip().lower()
    if refill_mode not in ("gemini", "source-crop", "auto"):
        refill_mode = "gemini"
    source_image_for_svg = None
    if image_refill_config.get("enabled") and refill_mode in ("source-crop", "auto"):
        source_image_for_svg = resolve_source_image_for_svg(svg_path, image_refill_config)

    # Prefetch placeholder image generation in parallel (gemini mode only).
    # This speeds up API-bound image refill while keeping slide compilation sequential/safe.
    if image_refill_config.get("enabled") and refill_mode == "gemini" and image_placeholders:
        max_workers = int(image_refill_config.get("max_concurrent_requests", 8) or 8)
        if max_workers > 1:
            prefetch_jobs: List[Dict[str, Any]] = []
            for ph in image_placeholders:
                placeholder_id = ph["id"]
                entry = (
                    placeholders_map.get((svg_abs, placeholder_id))
                    or placeholders_map.get((svg_name, placeholder_id))
                )
                is_chart = bool(entry.get("is_chart", False)) if entry is not None else bool(ph.get("is_chart", False))
                if is_chart:
                    continue
                if not str(ph.get("caption") or "").strip():
                    continue
                image_elem = ph.get("image_elem")
                mapped_image_path = str(entry.get("image_path", "")).strip() if entry else ""
                href = ""
                if image_elem is not None:
                    href = image_elem.get(f"{{{XLINK_NS}}}href") or image_elem.get("href") or ""
                if href or mapped_image_path:
                    continue

                remove_bg_pref = ph.get("remove_bg")
                if remove_bg_pref is None and entry:
                    raw_remove_bg = str(entry.get("remove_bg", "")).strip().lower()
                    if raw_remove_bg in ("1", "true", "yes", "on"):
                        remove_bg_pref = True
                    elif raw_remove_bg in ("0", "false", "no", "off"):
                        remove_bg_pref = False
                remove_bg_mode_pref = str(ph.get("remove_bg_mode") or "").strip().lower()
                if (not remove_bg_mode_pref) and entry:
                    remove_bg_mode_pref = str(entry.get("remove_bg_mode") or "").strip().lower()
                if remove_bg_mode_pref not in ("flat", "photo", "auto", "rembg", "chroma", "key", "model"):
                    remove_bg_mode_pref = ""

                prefetch_jobs.append(
                    {
                        "placeholder_id": placeholder_id,
                        "caption": str(ph.get("caption") or ""),
                        "w": float(ph.get("w") or 0.0),
                        "h": float(ph.get("h") or 0.0),
                        "remove_bg": (None if remove_bg_pref is None else bool(remove_bg_pref)),
                        "remove_bg_mode": (remove_bg_mode_pref or None),
                    }
                )

            if prefetch_jobs:
                print(
                    f"    Info: prefetching {len(prefetch_jobs)} placeholder images "
                    f"(concurrency={max_workers})"
                )
                with ThreadPoolExecutor(max_workers=max_workers) as pool:
                    fut_map = {}
                    for job in prefetch_jobs:
                        job_config = dict(image_refill_config)
                        job_config["_memo"] = {}
                        fut = pool.submit(
                            generate_placeholder_image,
                            caption=job["caption"],
                            width=job["w"],
                            height=job["h"],
                            image_refill_config=job_config,
                            remove_bg=job["remove_bg"],
                            remove_bg_mode=job["remove_bg_mode"],
                        )
                        fut_map[fut] = str(job["placeholder_id"])

                    for fut in as_completed(fut_map):
                        pid = fut_map[fut]
                        try:
                            prefetched_generation[pid] = fut.result()
                        except Exception as exc:
                            prefetched_generation[pid] = (None, False, f"prefetch crashed: {exc}")
    
    for ph in image_placeholders:
        placeholder_id = ph["id"]
        entry = (
            placeholders_map.get((svg_abs, placeholder_id))
            or placeholders_map.get((svg_name, placeholder_id))
        )
        if entry is not None:
            text_policy = str(entry.get("text_policy", ph.get("text_policy", "editable"))).strip().lower()
        else:
            text_policy = str(ph.get("text_policy", "editable")).strip().lower()
        if text_policy not in ("editable", "raster"):
            text_policy = "editable"
        
        # Check if this is a chart
        if entry is not None:
            is_chart = bool(entry.get("is_chart", False))
        else:
            is_chart = bool(ph.get("is_chart", False))
        
        if is_chart and not skip_charts:
            chart_built = False
            chart_fallback_image = bool(chart_config.get("chart_fallback_image", False))
            # This is a chart placeholder - must generate chart, not image
            if entry is not None:
                caption = entry.get("caption", "") or ph.get("caption", "")
                chart_spec = entry.get("chart_spec", ph.get("chart_spec"))
            else:
                caption = ph.get("caption", "")
                chart_spec = ph.get("chart_spec")
            if isinstance(chart_spec, str):
                raw_spec = chart_spec.strip()
                if raw_spec:
                    try:
                        chart_spec = json.loads(raw_spec)
                    except Exception:
                        chart_spec = raw_spec
                else:
                    chart_spec = None

            if (caption or chart_spec):
                # Generate chart code
                code = generate_chart_code(
                    caption,
                    chart_spec=chart_spec,
                    api_key=chart_config.get("api_key"),
                    base_url=chart_config.get("base_url"),
                    model=chart_config.get("chart_model"),
                    max_tokens=chart_config.get("chart_max_tokens", 800),
                    temperature=chart_config.get("chart_temperature", 0.2),
                )
                if code:
                    # Execute chart code
                    left_inches = Inches(ph["x"] / dpi)
                    top_inches = Inches(ph["y"] / dpi)
                    width_inches = Inches(ph["w"] / dpi)
                    height_inches = Inches(ph["h"] / dpi)
                    before_shapes = len(slide.shapes)
                    success, error = execute_chart_code(
                        code, slide, left_inches, top_inches, width_inches, height_inches
                    )
                    if success:
                        chart_count += 1
                        chart_built = True
                        after_shapes = len(slide.shapes)
                        if after_shapes > before_shapes:
                            for idx in range(before_shapes, after_shapes):
                                try:
                                    chart_shapes_to_front.append(slide.shapes[idx])
                                except Exception:
                                    pass
                    else:
                        print(f"    Warning: Failed to create chart '{placeholder_id}': {error}")
                else:
                    print(f"    Warning: Failed to generate chart code for '{placeholder_id}'")
            else:
                print(f"    Warning: Chart '{placeholder_id}' has no caption or chart-spec")
            if chart_built:
                continue
            if not chart_fallback_image:
                # Default behavior: do not fallback chart placeholders to images.
                print(f"    Warning: Chart '{placeholder_id}' not rendered as native chart; image fallback disabled.")
                continue
            if not image_refill_config.get("enabled"):
                print(f"    Warning: Chart '{placeholder_id}' image fallback enabled but refill is disabled.")
                continue
            print(f"    Info: Falling back to image generation for chart placeholder '{placeholder_id}'")
        
        # Add as image (only for non-chart placeholders)
        ph_for_insert = dict(ph)
        if entry:
            ph_for_insert["entry"] = entry
        ph_for_insert["text_policy"] = text_policy
        if text_policy == "raster" and ph_for_insert.get("remove_bg") is None:
            # Strongly bound text+graphic regions should keep rectangular crop by default.
            ph_for_insert["remove_bg"] = False

        remove_bg_pref = ph_for_insert.get("remove_bg")
        if remove_bg_pref is None and entry:
            raw_remove_bg = str(entry.get("remove_bg", "")).strip().lower()
            if raw_remove_bg in ("1", "true", "yes", "on"):
                remove_bg_pref = True
            elif raw_remove_bg in ("0", "false", "no", "off"):
                remove_bg_pref = False
        remove_bg_mode_pref = str(ph_for_insert.get("remove_bg_mode") or "").strip().lower()
        if (not remove_bg_mode_pref) and entry:
            remove_bg_mode_pref = str(entry.get("remove_bg_mode") or "").strip().lower()
        if remove_bg_mode_pref not in ("flat", "photo", "auto", "rembg", "chroma", "key", "model"):
            remove_bg_mode_pref = ""
        prefer_generation = bool(
            image_refill_config.get("prefer_generate_remove_bg", True)
        ) and (remove_bg_pref is True)

        shape = add_image_placeholder(slide, ph_for_insert, converter, svg_path)
        generated_new = False
        cropped_new = False
        if (
            shape is None
            and image_refill_config.get("enabled")
            and str(ph.get("caption") or "").strip()
        ):
            tried_generation_first = False

            # 0) for decorative cutout placeholders, prefer text-to-image first.
            if prefer_generation and refill_mode in ("gemini", "source-crop", "auto"):
                tried_generation_first = True
                prefetched = prefetched_generation.pop(placeholder_id, None)
                if prefetched is not None:
                    generated_path, created_new, gen_err = prefetched
                else:
                    generated_path, created_new, gen_err = generate_placeholder_image(
                        caption=str(ph.get("caption") or ""),
                        width=float(ph.get("w") or 0.0),
                        height=float(ph.get("h") or 0.0),
                        image_refill_config=image_refill_config,
                        remove_bg=(
                            None if remove_bg_pref is None else bool(remove_bg_pref)
                        ),
                        remove_bg_mode=(remove_bg_mode_pref or None),
                    )
                if generated_path is not None:
                    ph_for_insert = dict(ph_for_insert)
                    merged_entry = dict(ph_for_insert.get("entry") or {})
                    merged_entry["image_path"] = str(generated_path)
                    ph_for_insert["entry"] = merged_entry
                    shape = add_image_placeholder(slide, ph_for_insert, converter, svg_path)
                    generated_new = created_new
                elif gen_err:
                    print(f"    Warning: image generation failed for '{placeholder_id}': {gen_err}")

            # 1) source-crop path
            if refill_mode in ("source-crop", "auto"):
                if source_image_for_svg is None:
                    if refill_mode == "source-crop":
                        print(f"    Warning: source-crop mode enabled but source image not found for {svg_path.name}")
                else:
                    crop_path, crop_created, crop_err = crop_placeholder_from_source(
                        source_image_for_svg,
                        ph_for_insert,
                        svg_width=float(svg_width),
                        svg_height=float(svg_height),
                        image_refill_config=image_refill_config,
                        semantic_textboxes=([] if text_policy == "raster" else textboxes),
                    )
                    if crop_path is not None:
                        final_image_path = crop_path
                        redraw_enabled = bool(
                            image_refill_config.get("source_crop_redraw_no_text", False)
                        )
                        if redraw_enabled:
                            redraw_path, redraw_created, redraw_err = redraw_placeholder_crop_without_text(
                                seed_image=crop_path,
                                caption=str(ph.get("caption") or ""),
                                width=float(ph.get("w") or 0.0),
                                height=float(ph.get("h") or 0.0),
                                image_refill_config=image_refill_config,
                            )
                            if redraw_path is not None:
                                final_image_path = redraw_path
                                if redraw_created:
                                    source_redrawn_count += 1
                            elif redraw_err:
                                print(f"    Warning: source-crop redraw failed for '{placeholder_id}': {redraw_err}")
                        ph_for_insert = dict(ph_for_insert)
                        merged_entry = dict(ph_for_insert.get("entry") or {})
                        merged_entry["image_path"] = str(final_image_path)
                        merged_entry["fit_mode"] = str(
                            image_refill_config.get("source_crop_fit") or "contain"
                        ).strip().lower()
                        ph_for_insert["entry"] = merged_entry
                        shape = add_image_placeholder(slide, ph_for_insert, converter, svg_path)
                        cropped_new = crop_created
                        source_crop_regions.append(
                            {
                                "x": float(ph.get("x") or 0.0),
                                "y": float(ph.get("y") or 0.0),
                                "w": float(ph.get("w") or 0.0),
                                "h": float(ph.get("h") or 0.0),
                            }
                        )
                    elif crop_err and refill_mode == "source-crop":
                        print(f"    Warning: source-crop failed for '{placeholder_id}': {crop_err}")

            # 2) gemini image generation path
            if shape is None and (not tried_generation_first) and refill_mode in ("gemini", "auto"):
                prefetched = prefetched_generation.pop(placeholder_id, None)
                if prefetched is not None:
                    generated_path, created_new, gen_err = prefetched
                else:
                    generated_path, created_new, gen_err = generate_placeholder_image(
                        caption=str(ph.get("caption") or ""),
                        width=float(ph.get("w") or 0.0),
                        height=float(ph.get("h") or 0.0),
                        image_refill_config=image_refill_config,
                        remove_bg=(
                            None if remove_bg_pref is None else bool(remove_bg_pref)
                        ),
                        remove_bg_mode=(remove_bg_mode_pref or None),
                    )
                if generated_path is not None:
                    ph_for_insert = dict(ph_for_insert)
                    merged_entry = dict(ph_for_insert.get("entry") or {})
                    merged_entry["image_path"] = str(generated_path)
                    ph_for_insert["entry"] = merged_entry
                    shape = add_image_placeholder(slide, ph_for_insert, converter, svg_path)
                    generated_new = created_new
                elif gen_err:
                    print(f"    Warning: image generation failed for '{placeholder_id}': {gen_err}")

        if shape:
            image_count += 1
            if generated_new:
                generated_image_count += 1
            if cropped_new:
                source_cropped_count += 1
            visual_gid = str(ph.get("visual_group_id") or "").strip()
            if visual_gid:
                visual_group_shapes.setdefault(visual_gid, []).append(shape)
            if text_policy == "raster":
                raster_text_regions.append(
                    {
                        "x": float(ph.get("x") or 0.0),
                        "y": float(ph.get("y") or 0.0),
                        "w": float(ph.get("w") or 0.0),
                        "h": float(ph.get("h") or 0.0),
                    }
                )
                if visual_gid:
                    raster_text_visual_group_ids.add(visual_gid)

    # Now extract and add semantic textboxes (unless skipped)
    text_shape_count = 0
    grouped_textbox_count = 0
    grouped_visual_count = 0
    semantic_group_shapes: Dict[str, List[Any]] = {}
    semantic_group_order: Dict[str, float] = {}
    semantic_group_to_visual: Dict[str, str] = {}
    semantic_group_objects: Dict[str, Any] = {}
    semantic_group_fallback_members: Dict[str, List[Any]] = {}
    semantic_group_standalone_members: Dict[str, List[Any]] = {}
    flatten_semantic_for_visual = bool(
        image_refill_config.get("flatten_semantic_groups_for_visual", True)
    )
    if not skip_textboxes:
        for sgid, meta in semantic_group_meta.items():
            visual_gid = str(meta.get("visual_group_id") or "").strip()
            if visual_gid:
                semantic_group_to_visual[sgid] = visual_gid
        suppress_overlap_text = bool(image_refill_config.get("source_crop_skip_overlapped_text", True))
        if bool(image_refill_config.get("source_crop_erase_overlapped_text_in_image", True)):
            suppress_overlap_text = False
        if bool(image_refill_config.get("source_crop_redraw_no_text", False)):
            suppress_overlap_text = False
        overlap_threshold = float(image_refill_config.get("source_crop_text_overlap_threshold", 0.55) or 0.55)
        raster_overlap_threshold = float(
            image_refill_config.get("raster_text_overlap_threshold", 0.2) or 0.2
        )
        for tb in textboxes:
            if tb.get("w", 0) <= 0 or tb.get("h", 0) <= 0:
                continue
            tb_visual_gid = str(tb.get("visual_group_id") or "").strip()
            if tb_visual_gid and tb_visual_gid in raster_text_visual_group_ids:
                # Strong rule: a raster-policy visual group owns its in-graphic labels.
                suppressed_textbox_count += 1
                continue
            if raster_text_regions:
                tb_area = float(tb.get("w") or 0.0) * float(tb.get("h") or 0.0)
                if tb_area > 0:
                    max_raster_overlap_ratio = 0.0
                    for region in raster_text_regions:
                        inter = rect_intersection_area(tb, region)
                        if inter <= 0:
                            continue
                        ratio = inter / tb_area
                        if ratio > max_raster_overlap_ratio:
                            max_raster_overlap_ratio = ratio
                    if max_raster_overlap_ratio >= raster_overlap_threshold:
                        suppressed_textbox_count += 1
                        continue
            if suppress_overlap_text and source_crop_regions:
                tb_area = float(tb.get("w") or 0.0) * float(tb.get("h") or 0.0)
                if tb_area > 0:
                    max_overlap_ratio = 0.0
                    for region in source_crop_regions:
                        inter = rect_intersection_area(tb, region)
                        if inter <= 0:
                            continue
                        ratio = inter / tb_area
                        if ratio > max_overlap_ratio:
                            max_overlap_ratio = ratio
                    if max_overlap_ratio >= overlap_threshold:
                        suppressed_textbox_count += 1
                        continue
            created_shapes: List[Any] = []
            if semantic_mode == "absolute":
                created_shapes = add_semantic_text_items_absolute(
                    slide, tb, converter=converter, cjk_font=cjk_font
                )
                if created_shapes:
                    text_shape_count += len(created_shapes)
                else:
                    shape = add_semantic_textbox(
                        slide, tb, converter=converter, line_tol=line_tol, box_pad=box_pad, cjk_font=cjk_font
                    )
                    if shape:
                        created_shapes = [shape]
                        text_shape_count += 1
            else:
                shape = add_semantic_textbox(
                    slide, tb, converter=converter, line_tol=line_tol, box_pad=box_pad, cjk_font=cjk_font
                )
                if shape:
                    created_shapes = [shape]
                    text_shape_count += 1

            group_id = str(tb.get("group_id") or "").strip()
            if group_id and created_shapes:
                semantic_group_shapes.setdefault(group_id, []).extend(created_shapes)
                visual_gid = str(tb.get("visual_group_id") or "").strip()
                if visual_gid and group_id not in semantic_group_to_visual:
                    semantic_group_to_visual[group_id] = visual_gid
                try:
                    order_val = float(tb.get("group_order") or 1e9)
                except Exception:
                    order_val = 1e9
                if group_id not in semantic_group_order:
                    semantic_group_order[group_id] = order_val
                else:
                    semantic_group_order[group_id] = min(semantic_group_order[group_id], order_val)

    # Group textboxes by semantic text-group so they can be dragged together in PPT.
    if semantic_group_shapes:
        for gid in sorted(
            semantic_group_shapes.keys(),
            key=lambda k: (semantic_group_order.get(k, 1e9), k),
        ):
            members = [s for s in semantic_group_shapes.get(gid, []) if s is not None]
            if not members:
                continue
            # Avoid nested group-shapes in python-pptx: they are prone to coordinate drift.
            # If this semantic group belongs to a visual-group, keep members flat and let
            # visual-group do the single grouping pass.
            if flatten_semantic_for_visual and semantic_group_to_visual.get(gid):
                semantic_group_fallback_members[gid] = members
                continue
            if not semantic_group_to_visual.get(gid):
                # Delay standalone semantic grouping until after visual-groups are built,
                # so they are never hidden under background visual layers.
                semantic_group_standalone_members[gid] = members
                continue
            try:
                gshape = slide.shapes.add_group_shape(members)
                gshape.name = gid
                grouped_textbox_count += 1
                semantic_group_objects[gid] = gshape
            except Exception as exc:
                print(f"    Warning: failed to group semantic text group '{gid}': {exc}")
                semantic_group_fallback_members[gid] = members

    # Group by visual-group (decorations + semantic text groups) so the whole block can be dragged together.
    for sgid, vgid in semantic_group_to_visual.items():
        if not vgid:
            continue
        if sgid in semantic_group_objects:
            visual_group_shapes.setdefault(vgid, []).append(semantic_group_objects[sgid])
        elif sgid in semantic_group_fallback_members:
            visual_group_shapes.setdefault(vgid, []).extend(semantic_group_fallback_members[sgid])
        elif sgid in semantic_group_shapes:
            visual_group_shapes.setdefault(vgid, []).extend(semantic_group_shapes[sgid])

    if visual_group_shapes:
        for vgid in sorted(
            visual_group_shapes.keys(),
            key=lambda k: (visual_group_order.get(k, 1e9), k),
        ):
            seen_members: set = set()
            members: List[Any] = []
            for s in visual_group_shapes.get(vgid, []):
                if s is None:
                    continue
                sid = id(s)
                if sid in seen_members:
                    continue
                seen_members.add(sid)
                members.append(s)
            if not members:
                continue
            try:
                vgshape = slide.shapes.add_group_shape(members)
                vgshape.name = vgid
                grouped_visual_count += 1
            except Exception as exc:
                print(f"    Warning: failed to group visual group '{vgid}': {exc}")

    # Add standalone semantic text-groups after visual groups to keep them visible/editable.
    if semantic_group_standalone_members:
        for gid in sorted(
            semantic_group_standalone_members.keys(),
            key=lambda k: (semantic_group_order.get(k, 1e9), k),
        ):
            members = [s for s in semantic_group_standalone_members.get(gid, []) if s is not None]
            if not members:
                continue
            try:
                gshape = slide.shapes.add_group_shape(members)
                gshape.name = gid
                grouped_textbox_count += 1
                semantic_group_objects[gid] = gshape
            except Exception as exc:
                print(f"    Warning: failed to group standalone semantic group '{gid}': {exc}")

    # Keep chart frames above grouped background/card shapes.
    if chart_shapes_to_front:
        for shape in chart_shapes_to_front:
            try:
                elem = shape._element
                parent = elem.getparent()
                if parent is not None:
                    parent.remove(elem)
                    parent.append(elem)
            except Exception:
                pass

    total_shapes = len(slide.shapes)
    parts = [f"{total_shapes} shapes"]
    if image_count:
        parts.append(f"{image_count} images")
    if generated_image_count:
        parts.append(f"{generated_image_count} generated")
    if source_cropped_count:
        parts.append(f"{source_cropped_count} source-cropped")
    if source_redrawn_count:
        parts.append(f"{source_redrawn_count} redrawn-no-text")
    if suppressed_textbox_count:
        parts.append(f"{suppressed_textbox_count} text-suppressed")
    if chart_count:
        parts.append(f"{chart_count} charts")
    if text_shape_count:
        parts.append(f"{text_shape_count} text-shapes")
    if grouped_textbox_count:
        parts.append(f"{grouped_textbox_count} text-groups")
    if grouped_visual_count:
        parts.append(f"{grouped_visual_count} visual-groups")
    print(f"  Converted: {svg_path.name} ({', '.join(parts)})")


def inject_svg_blips_into_pptx(pptx_path: Path, jobs: List[Dict[str, Any]]) -> None:
    _inject_svg_blips_into_pptx_impl(pptx_path, jobs)


def build_pptx_pro(
    svg_paths: List[Path],
    out_pptx: Path,
    dpi: float = 96.0,
    cjk_font: str = "PingFang SC",
    skip_elements: Optional[set] = None,
    line_tol: float = 2.0,
    box_pad: float = 0.0,
    semantic_mode: str = "textbox",
    placeholders_map: Optional[Dict[Tuple[str, str], Dict[str, Any]]] = None,
    chart_config: Optional[Dict[str, Any]] = None,
    skip_charts: bool = False,
    image_refill_config: Optional[Dict[str, Any]] = None,
) -> None:
    """Build PPTX from multiple SVG files with editable shapes.
    
    Args:
        placeholders_map: Mapping from (svg_file, placeholder_id) to placeholder info
        chart_config: Configuration for chart generation (API key, model, etc.)
        skip_charts: If True, skip chart generation
        image_refill_config: Optional config for image-placeholder generation/refill
    """
    if not svg_paths:
        print("No SVG files to convert.")
        return
    
    # Parse first SVG to get dimensions
    first_tree = ET.parse(svg_paths[0])
    first_root = first_tree.getroot()
    svg_width = parse_length(first_root.get("width"), 1920.0)
    svg_height = parse_length(first_root.get("height"), 1080.0)
    
    # Handle viewBox
    viewbox = first_root.get("viewBox", "")
    if viewbox:
        vb_parts = viewbox.split()
        if len(vb_parts) == 4:
            try:
                svg_width = float(vb_parts[2])
                svg_height = float(vb_parts[3])
            except ValueError:
                pass
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(svg_width / dpi)
    prs.slide_height = Inches(svg_height / dpi)
    
    print(f"Creating PPTX: {svg_width}x{svg_height} px @ {dpi} DPI")
    print(f"Slide size: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    
    # Convert each SVG
    svg_blip_jobs: List[Dict[str, Any]] = []
    for svg_path in svg_paths:
        convert_svg_to_slide(
            prs, svg_path, dpi, cjk_font, skip_elements, line_tol, box_pad, semantic_mode,
            placeholders_map, chart_config, skip_charts, image_refill_config, svg_blip_jobs
        )
    
    # Save
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    if svg_blip_jobs:
        inject_svg_blips_into_pptx(out_pptx, svg_blip_jobs)
    print(f"\nSaved: {out_pptx}")


# ============================================================================
# Command Line Interface
# ============================================================================

def main() -> None:
    parser = argparse.ArgumentParser(
        description="SVG to PPTX Pro: Convert SVG to editable PPTX shapes."
    )
    parser.add_argument(
        "--input", "-i", required=True,
        help="Input SVG file or directory containing SVG files."
    )
    parser.add_argument(
        "--output", "-o", required=True,
        help="Output PPTX file path."
    )
    parser.add_argument(
        "--dpi", type=float, default=96.0,
        help="SVG pixel DPI (default: 96.0)."
    )
    parser.add_argument(
        "--cjk-font", default="PingFang SC",
        help="Font for CJK text (default: PingFang SC)."
    )
    parser.add_argument(
        "--line-tol", type=float, default=2.0,
        help="Line grouping tolerance in px for textboxes (default: 2.0)."
    )
    parser.add_argument(
        "--box-pad", type=float, default=0.0,
        help="Textbox padding in px (default: 0.0)."
    )
    parser.add_argument(
        "--semantic-mode", choices=("absolute", "textbox"), default="textbox",
        help="Semantic text reconstruction mode: absolute (per text item) or textbox (grouped reflow)."
    )
    parser.add_argument(
        "--skip-textboxes", action="store_true",
        help="Skip semantic layer textboxes (don't add text)."
    )
    parser.add_argument(
        "--placeholders", type=str,
        help="Path to image_placeholders.json for chart detection."
    )
    parser.add_argument(
        "--config", type=str,
        help=(
            "Path to config.json for API settings (chart generation). "
            f"If omitted, auto-load {DEFAULT_RUNTIME_CONFIG_PATH} when present."
        ),
    )
    parser.add_argument(
        "--profile",
        choices=_SUPPORTED_PROFILES_IMPL,
        default=_DEFAULT_PROFILE_IMPL,
        help=f"Task profile (default: {_DEFAULT_PROFILE_IMPL}).",
    )
    parser.add_argument(
        "--profile-dir",
        default="",
        help="Optional profile config directory (default: gemini_pipeline/profiles).",
    )
    parser.add_argument(
        "--skip-charts", action="store_true",
        help="Skip chart generation (add images instead)."
    )
    parser.add_argument(
        "--chart-fallback-image",
        action=argparse.BooleanOptionalAction,
        default=None,
        help="Allow chart placeholders to fallback to image refill when native chart generation fails (default: off).",
    )
    parser.add_argument(
        "--refill-placeholders", action="store_true",
        help="Refill image-placeholder blocks using source-crop/gemini strategy."
    )
    parser.add_argument(
        "--refill-mode", choices=("gemini", "source-crop", "auto"), default="gemini",
        help="Placeholder refill mode: gemini, source-crop, or auto (source-crop then gemini)."
    )
    parser.add_argument(
        "--image-api-key",
        default="",
        help="API key for image generation (default: from runtime config; env as fallback)."
    )
    parser.add_argument(
        "--image-api-base", default="",
        help="Optional API base for image generation (default: from runtime config)."
    )
    parser.add_argument(
        "--image-model", default="gemini-3.1-flash-image-preview",
        help="Gemini image model for placeholder refill."
    )
    parser.add_argument(
        "--image-size", default="2K",
        help="Generated image size for image preview models (e.g., 1K/2K/4K)."
    )
    parser.add_argument(
        "--image-aspect-ratio", default="",
        help="Optional fixed aspect ratio like 16:9. If empty, infer from placeholder box."
    )
    parser.add_argument(
        "--image-cache-dir", default="",
        help="Cache directory for generated placeholder images."
    )
    parser.add_argument(
        "--image-max-concurrent", type=int, default=8,
        help="Max concurrent API requests for image-placeholder generation prefetch (default: 8)."
    )
    parser.add_argument(
        "--image-style-prompt", default="",
        help="Optional extra style constraint appended to each placeholder caption."
    )
    parser.add_argument(
        "--source-image", default="",
        help="Source raster image path for source-crop mode (single SVG conversion)."
    )
    parser.add_argument(
        "--source-image-dir", default="",
        help="Source image directory for source-crop mode (match by SVG stem)."
    )
    parser.add_argument(
        "--source-crop-expand", type=float, default=0.0,
        help="Expand source-crop bbox by this many SVG px on each side (default: 0)."
    )
    parser.add_argument(
        "--source-crop-fit", choices=("contain", "cover", "stretch"), default="contain",
        help="Source-crop image fit mode into placeholder box (default: contain)."
    )
    parser.add_argument(
        "--source-crop-bg-threshold", type=float, default=14.0,
        help="Background removal threshold for source-crop RGBA cutout (default: 14)."
    )
    parser.add_argument(
        "--source-crop-feather", type=float, default=18.0,
        help="Background removal feather radius proxy for source-crop (default: 18)."
    )
    parser.add_argument(
        "--source-crop-no-remove-bg", action="store_true",
        help="Disable source-crop background alpha removal."
    )
    parser.add_argument(
        "--source-crop-no-suppress-overlap-text", action="store_true",
        help="Do not suppress semantic textboxes that overlap source-cropped image regions."
    )
    parser.add_argument(
        "--source-crop-text-overlap-threshold", type=float, default=0.55,
        help="Suppress semantic textbox when overlap_area/textbox_area >= threshold (default: 0.55)."
    )
    parser.add_argument(
        "--source-crop-no-erase-overlap-text-in-image", action="store_true",
        help="Do not erase overlapped semantic text from source-cropped image content."
    )
    parser.add_argument(
        "--source-crop-text-erase-overlap-threshold", type=float, default=0.2,
        help="Erase text in source-crop when overlap_area/textbox_area >= threshold (default: 0.2)."
    )
    parser.add_argument(
        "--source-crop-text-erase-overlap-threshold-remove-bg", type=float, default=0.01,
        help="Stricter erase threshold used when placeholder remove_bg=true (default: 0.01)."
    )
    parser.add_argument(
        "--source-crop-text-erase-pad", type=float, default=2.0,
        help="Expand text erase box by this many SVG px on each side (default: 2.0)."
    )
    parser.add_argument(
        "--source-crop-text-erase-feather", type=float, default=1.0,
        help="Feather radius for alpha erase mask in source-crop (default: 1.0)."
    )
    parser.add_argument(
        "--source-crop-text-erase-mode", choices=("alpha", "fill"), default="alpha",
        help="How to remove overlapped text from source-cropped image: alpha or fill."
    )
    parser.add_argument(
        "--source-crop-redraw-no-text", action="store_true",
        help="After source-crop, use image-to-image generation to redraw crop without text."
    )
    parser.add_argument(
        "--source-crop-redraw-model", default="",
        help="Optional model override for source-crop redraw step (default: --image-model)."
    )
    parser.add_argument(
        "--source-crop-redraw-prompt", default="",
        help="Optional prompt override for source-crop redraw step."
    )
    parser.add_argument(
        "--source-crop-use-rembg",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Use rembg first for source-crop remove_bg placeholders (default: on).",
    )
    parser.add_argument(
        "--source-crop-rembg-min-nonopaque-ratio",
        type=float,
        default=0.02,
        help="Min non-opaque ratio required after rembg for source-crop (default: 0.02).",
    )
    parser.add_argument(
        "--source-crop-rembg-max-nonopaque-ratio",
        type=float,
        default=0.995,
        help="Max non-opaque ratio allowed after rembg for source-crop (default: 0.995).",
    )
    parser.add_argument(
        "--source-crop-require-rembg",
        action="store_true",
        help="Fail source-crop remove_bg when rembg is unavailable or invalid.",
    )
    parser.add_argument(
        "--source-crop-max-remove-bg-area-ratio", type=float, default=0.16,
        help="Disable RGBA remove-bg for large placeholders whose area ratio exceeds this value (default: 0.16)."
    )
    parser.add_argument(
        "--remove-bg-require-rembg",
        action="store_true",
        help="Fail generation remove_bg when rembg is unavailable or invalid.",
    )
    parser.add_argument(
        "--prefer-generate-remove-bg",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Prefer model generation before source-crop for placeholders marked remove_bg=true (default: on).",
    )
    parser.add_argument(
        "--allow-nested-group-shapes", action="store_true",
        help="Allow semantic-group then visual-group nested grouping (not recommended, may drift)."
    )
    parser.add_argument(
        "--hybrid-svgblip-icons",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Render complex icon sub-groups as embedded svgBlip pictures for higher visual fidelity (default: on).",
    )
    
    args = parser.parse_args()
    
    input_path = Path(args.input)
    
    # Collect SVG files
    if input_path.is_dir():
        all_svgs = list(input_path.glob("*.SVG")) + list(input_path.glob("*.svg"))
        svg_paths = sorted(
            [p for p in all_svgs if not p.name.startswith("._")],
            key=natural_sort_key
        )
    else:
        svg_paths = [input_path]
    
    if not svg_paths:
        print(f"No SVG files found in: {input_path}")
        return
    
    print(f"Found {len(svg_paths)} SVG file(s)")
    
    # Load placeholders and config
    placeholders_map = None
    chart_config = None
    
    if args.placeholders:
        placeholders_path = Path(args.placeholders)
        placeholders_map = load_placeholders(placeholders_path)
        if placeholders_map:
            print(f"Loaded {len(placeholders_map)} placeholder entries")
    
    config_path: Optional[Path] = None
    if args.config:
        config_path = Path(args.config)
    elif DEFAULT_RUNTIME_CONFIG_PATH.exists():
        config_path = DEFAULT_RUNTIME_CONFIG_PATH

    if config_path is not None:
        chart_config = load_config(config_path)
        if chart_config:
            print(f"Loaded config from {config_path}")

    chart_config = chart_config or {}
    profile_dir = Path(args.profile_dir).expanduser().resolve() if args.profile_dir else None
    profile_spec: Dict[str, Any] = {}
    try:
        profile_spec = load_profile_spec(args.profile, profile_dir=profile_dir)
    except Exception as exc:  # noqa: BLE001
        print(f"Warning: failed to load profile '{args.profile}': {exc}")
        profile_spec = {}
    if profile_spec:
        chart_config = apply_profile_overrides(chart_config, profile_spec, section="compiler")
        print(f"Loaded profile '{args.profile}'")

    profile_cli = profile_spec.get("compiler_cli") if isinstance(profile_spec.get("compiler_cli"), dict) else {}
    if profile_cli:
        if not _cli_flag_present("--semantic-mode"):
            prof_semantic = str(profile_cli.get("semantic_mode") or "").strip()
            if prof_semantic in ("absolute", "textbox"):
                args.semantic_mode = prof_semantic
        if not _cli_flag_present("--refill-mode"):
            prof_refill = str(profile_cli.get("refill_mode") or "").strip()
            if prof_refill in ("gemini", "source-crop", "auto"):
                args.refill_mode = prof_refill
        if not _cli_flag_present("--image-model"):
            prof_model = str(profile_cli.get("image_model") or "").strip()
            if prof_model:
                args.image_model = prof_model

    if not str(chart_config.get("base_url") or "").strip():
        chart_config["base_url"] = "https://cdn.12ai.org/v1"
    if not str(chart_config.get("chart_model") or "").strip():
        chart_config["chart_model"] = "gemini-3.1-pro-preview"
    if args.chart_fallback_image is not None:
        chart_config["chart_fallback_image"] = bool(args.chart_fallback_image)

    image_refill_config: Optional[Dict[str, Any]] = {
        "hybrid_svg_blip_icons": bool(args.hybrid_svgblip_icons),
    }
    if args.refill_placeholders:
        image_model = args.image_model
        if (
            str(image_model).strip() == "gemini-3.1-flash-image-preview"
            and str(chart_config.get("image_model") or "").strip()
        ):
            image_model = str(chart_config.get("image_model") or "").strip()
        cache_dir = Path(args.image_cache_dir) if args.image_cache_dir else (
            Path(args.output).resolve().parent / f"{Path(args.output).stem}_generated_images"
        )
        image_refill_config.update({
            "enabled": True,
            "mode": args.refill_mode,
            "api_key": (
                args.image_api_key
                or chart_config.get("image_api_key")
                or chart_config.get("api_key")
                or os.environ.get("OPENAI_API_KEY", "")
                or os.environ.get("GEMINI_API_KEY", "")
                or ""
            ).strip(),
            "api_base": (
                args.image_api_base
                or chart_config.get("image_api_base")
                or chart_config.get("base_url")
                or os.environ.get("OPENAI_API_BASE", "")
                or ""
            ).strip(),
            "model": image_model,
            "image_size": args.image_size,
            "aspect_ratio": args.image_aspect_ratio.strip(),
            "cache_dir": str(cache_dir),
            "max_concurrent_requests": max(1, int(args.image_max_concurrent or 8)),
            "style_prompt": args.image_style_prompt.strip(),
            "source_image": args.source_image.strip(),
            "source_image_dir": args.source_image_dir.strip(),
            "source_crop_expand": float(args.source_crop_expand or 0.0),
            "source_crop_fit": args.source_crop_fit,
            "source_crop_bg_threshold": float(args.source_crop_bg_threshold or 14.0),
            "source_crop_feather": float(args.source_crop_feather or 18.0),
            "source_crop_remove_bg": (not bool(args.source_crop_no_remove_bg)),
            "source_crop_skip_overlapped_text": (not bool(args.source_crop_no_suppress_overlap_text)),
            "source_crop_text_overlap_threshold": float(args.source_crop_text_overlap_threshold or 0.55),
            "source_crop_erase_overlapped_text_in_image": (not bool(args.source_crop_no_erase_overlap_text_in_image)),
            "source_crop_text_erase_overlap_threshold": float(args.source_crop_text_erase_overlap_threshold or 0.2),
            "source_crop_text_erase_overlap_threshold_remove_bg": float(
                args.source_crop_text_erase_overlap_threshold_remove_bg or 0.01
            ),
            "source_crop_text_erase_pad": float(args.source_crop_text_erase_pad or 2.0),
            "source_crop_text_erase_feather": float(args.source_crop_text_erase_feather or 1.0),
            "source_crop_text_erase_mode": args.source_crop_text_erase_mode,
            "source_crop_redraw_no_text": bool(args.source_crop_redraw_no_text),
            "source_crop_redraw_model": args.source_crop_redraw_model.strip(),
            "source_crop_redraw_prompt": args.source_crop_redraw_prompt.strip(),
            "source_crop_use_rembg": bool(args.source_crop_use_rembg),
            "source_crop_rembg_min_nonopaque_ratio": float(
                args.source_crop_rembg_min_nonopaque_ratio or 0.02
            ),
            "source_crop_rembg_max_nonopaque_ratio": float(
                args.source_crop_rembg_max_nonopaque_ratio or 0.995
            ),
            "source_crop_require_rembg": bool(args.source_crop_require_rembg),
            "source_crop_max_remove_bg_area_ratio": float(
                args.source_crop_max_remove_bg_area_ratio or 0.16
            ),
            "remove_bg_require_rembg": bool(args.remove_bg_require_rembg),
            "remove_bg_mode_default": str(
                chart_config.get("remove_bg_mode_default") or "auto"
            ).strip().lower(),
            "flat_remove_bg_threshold": float(
                chart_config.get("flat_remove_bg_threshold", 16.0) or 16.0
            ),
            "flat_remove_bg_feather": float(
                chart_config.get("flat_remove_bg_feather", 2.0) or 2.0
            ),
            "flat_chroma_enabled": bool(
                chart_config.get("flat_chroma_enabled", True)
            ),
            "flat_chroma_hue_low": float(
                chart_config.get("flat_chroma_hue_low", 70.0) or 70.0
            ),
            "flat_chroma_hue_high": float(
                chart_config.get("flat_chroma_hue_high", 170.0) or 170.0
            ),
            "flat_chroma_sat_min": float(
                chart_config.get("flat_chroma_sat_min", 0.18) or 0.18
            ),
            "flat_chroma_val_min": float(
                chart_config.get("flat_chroma_val_min", 0.10) or 0.10
            ),
            "flat_chroma_g_min_weak": int(
                chart_config.get("flat_chroma_g_min_weak", 90) or 90
            ),
            "flat_chroma_g_min_strong": int(
                chart_config.get("flat_chroma_g_min_strong", 120) or 120
            ),
            "flat_chroma_dom_weak": int(
                chart_config.get("flat_chroma_dom_weak", 18) or 18
            ),
            "flat_chroma_dom_strong": int(
                chart_config.get("flat_chroma_dom_strong", 42) or 42
            ),
            "flat_chroma_alpha_blur": float(
                chart_config.get("flat_chroma_alpha_blur", 0.6) or 0.6
            ),
            "flat_remove_small_green_islands": bool(
                chart_config.get("flat_remove_small_green_islands", True)
            ),
            "flat_green_island_min_area": int(
                chart_config.get("flat_green_island_min_area", 20) or 20
            ),
            "flat_green_island_dom": int(
                chart_config.get("flat_green_island_dom", 18) or 18
            ),
            "flat_green_island_ratio_min": float(
                chart_config.get("flat_green_island_ratio_min", 0.55) or 0.55
            ),
            "flat_despill_enabled": bool(
                chart_config.get("flat_despill_enabled", True)
            ),
            "flat_despill_strength": float(
                chart_config.get("flat_despill_strength", 0.8) or 0.8
            ),
            "flat_despill_alpha_max": int(
                chart_config.get("flat_despill_alpha_max", 252) or 252
            ),
            "flat_despill_dom_threshold": int(
                chart_config.get("flat_despill_dom_threshold", 8) or 8
            ),
            "flat_remove_bg_hard_threshold": int(
                chart_config.get("flat_remove_bg_hard_threshold", 242) or 242
            ),
            "flat_remove_bg_hard_blur": float(
                chart_config.get("flat_remove_bg_hard_blur", 0.8) or 0.8
            ),
            "flat_remove_bg_min_nonopaque_ratio": float(
                chart_config.get("flat_remove_bg_min_nonopaque_ratio", 0.01) or 0.01
            ),
            "flat_remove_bg_max_transparent_ratio": float(
                chart_config.get("flat_remove_bg_max_transparent_ratio", 0.999) or 0.999
            ),
            "flat_mode_fallback_rembg": bool(
                chart_config.get("flat_mode_fallback_rembg", True)
            ),
            "rembg_post_hard_enabled": bool(
                chart_config.get("rembg_post_hard_enabled", True)
            ),
            "rembg_post_hard_trigger_semi_ratio": float(
                chart_config.get("rembg_post_hard_trigger_semi_ratio", 0.2) or 0.2
            ),
            "rembg_post_hard_threshold": int(
                chart_config.get("rembg_post_hard_threshold", 238) or 238
            ),
            "rembg_post_hard_blur": float(
                chart_config.get("rembg_post_hard_blur", 0.7) or 0.7
            ),
            "remove_bg_style_prompt": str(
                chart_config.get("remove_bg_style_prompt") or ""
            ).strip(),
            "remove_bg_style_prompt_flat": str(
                chart_config.get("remove_bg_style_prompt_flat") or ""
            ).strip(),
            "remove_bg_style_prompt_photo": str(
                chart_config.get("remove_bg_style_prompt_photo") or ""
            ).strip(),
            "prefer_generate_remove_bg": bool(args.prefer_generate_remove_bg),
            "flatten_semantic_groups_for_visual": (not bool(args.allow_nested_group_shapes)),
            "_memo": {},
        })
        if (
            args.refill_mode in ("gemini", "auto")
            and not str(image_refill_config["api_key"]).strip()
        ):
            print("Warning: --refill-placeholders is enabled but image API key is empty.")
        if (
            args.refill_mode == "source-crop"
            and bool(args.source_crop_redraw_no_text)
            and not str(image_refill_config["api_key"]).strip()
        ):
            print("Warning: source-crop redraw is enabled but image API key is empty.")
        image_api_base = str(image_refill_config.get("api_base") or "").strip().lower()
        uses_google_refill = (
            "googleapis.com" in image_api_base
            or "ai.google.dev" in image_api_base
            or "generativelanguage.googleapis.com" in image_api_base
            or "aiplatform.googleapis.com" in image_api_base
            or (not image_api_base)
        )
        if args.refill_mode in ("gemini", "auto") and uses_google_refill and not HAS_GOOGLE_GENAI:
            print("Warning: google-genai is not installed; image refill will be skipped.")
        if (
            args.refill_mode == "source-crop"
            and bool(args.source_crop_redraw_no_text)
            and uses_google_refill
            and not HAS_GOOGLE_GENAI
        ):
            print("Warning: google-genai is not installed; source-crop redraw will be skipped.")
        print(
            f"Image refill enabled: mode={image_refill_config['mode']}, "
            f"model={image_refill_config['model']}, size={image_refill_config['image_size']}, "
            f"api_base={image_refill_config.get('api_base') or '(google-sdk-default)'}, "
            f"cache={image_refill_config['cache_dir']}"
        )
    
    # Skip elements - if skip_textboxes is set, we skip the semantic layer entirely
    skip_elements = set()
    if args.skip_textboxes:
        skip_elements.add("semantic-layer")
    
    # Build PPTX
    build_pptx_pro(
        svg_paths,
        Path(args.output),
        dpi=args.dpi,
        cjk_font=args.cjk_font,
        skip_elements=skip_elements if skip_elements else None,
        line_tol=args.line_tol,
        box_pad=args.box_pad,
        semantic_mode=args.semantic_mode,
        placeholders_map=placeholders_map,
        chart_config=chart_config,
        skip_charts=args.skip_charts,
        image_refill_config=image_refill_config,
    )


if __name__ == "__main__":
    main()
