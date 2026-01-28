#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Render SVG background to PNG and rebuild semantic textboxes in a PPTX slide.
Supports native chart reconstruction via code model when is_chart=true.
"""

import argparse
import json
import math
import os
import re
import tempfile
import threading
import time
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    import cairosvg  # type: ignore
except Exception:  # pragma: no cover - optional dependency
    cairosvg = None

try:
    from openai import OpenAI  # type: ignore
except Exception:
    OpenAI = None

from pptx import Presentation
try:
    from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
except Exception:
    from pptx.chart.data import CategoryChartData
    XyChartData = None
    BubbleChartData = None
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SVG_NS = "http://www.w3.org/2000/svg"
XML_DECL_RE = re.compile(r"<\?xml[^>]*\?>", re.I)
MATRIX_RE = re.compile(r"matrix\(([^)]+)\)")
TRANSLATE_RE = re.compile(r"translate\(([^)]+)\)")
CJK_RE = re.compile(r"[\u3400-\u9fff\u3000-\u303f\u3040-\u30ff\u31f0-\u31ff\uac00-\ud7af]")
NATURAL_SORT_RE = re.compile(r"(\d+)")
GENERIC_FONTS = {
    "sans-serif",
    "serif",
    "monospace",
    "system-ui",
    "ui-sans-serif",
    "ui-serif",
    "ui-monospace",
}
MERGE_GROUP_ROLES = {"body", "bullet", "numbered"}


class RateLimiter:
    def __init__(self, qps: float):
        self.min_interval = 1.0 / qps if qps > 0 else 0.0
        self.lock = threading.Lock()
        self.next_time = time.monotonic()

    def acquire(self) -> None:
        if self.min_interval <= 0:
            return
        with self.lock:
            now = time.monotonic()
            if now < self.next_time:
                time.sleep(self.next_time - now)
            self.next_time = max(now, self.next_time) + self.min_interval


def tag_name(elem: ET.Element) -> str:
    return elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag


def natural_sort_key(path: Path) -> List[Any]:
    """Generate a sort key for natural sorting (e.g., 幻灯片1, 幻灯片2, ..., 幻灯片10)."""
    parts = NATURAL_SORT_RE.split(path.name)
    return [int(p) if p.isdigit() else p.lower() for p in parts]


def parse_length(val: Optional[str]) -> float:
    if not val:
        return 0.0
    raw = val.strip()
    for suf in ("px", "pt", "mm", "cm", "in"):
        if raw.endswith(suf):
            raw = raw[: -len(suf)]
            break
    try:
        return float(raw)
    except Exception:
        return 0.0


def parse_color(val: Optional[str]) -> Optional[RGBColor]:
    if not val:
        return None
    v = val.strip()
    if v.startswith("#") and len(v) in (4, 7):
        if len(v) == 4:
            r = int(v[1] * 2, 16)
            g = int(v[2] * 2, 16)
            b = int(v[3] * 2, 16)
        else:
            r = int(v[1:3], 16)
            g = int(v[3:5], 16)
            b = int(v[5:7], 16)
        return RGBColor(r, g, b)
    if v.startswith("rgb(") and v.endswith(")"):
        parts = v[4:-1].split(",")
        try:
            r, g, b = [int(float(p.strip())) for p in parts[:3]]
            return RGBColor(r, g, b)
        except Exception:
            return None
    return None


def parse_transform_xy(transform: Optional[str]) -> Tuple[float, float]:
    if not transform:
        return (0.0, 0.0)
    match = MATRIX_RE.search(transform)
    if match:
        parts = [p for p in re.split(r"[ ,]+", match.group(1).strip()) if p]
        if len(parts) == 6:
            try:
                return (float(parts[4]), float(parts[5]))
            except Exception:
                return (0.0, 0.0)
    match = TRANSLATE_RE.search(transform)
    if match:
        parts = [p for p in re.split(r"[ ,]+", match.group(1).strip()) if p]
        if parts:
            try:
                x = float(parts[0])
                y = float(parts[1]) if len(parts) > 1 else 0.0
                return (x, y)
            except Exception:
                return (0.0, 0.0)
    return (0.0, 0.0)


def normalize_rotation(angle: float) -> float:
    while angle <= -180:
        angle += 360
    while angle > 180:
        angle -= 360
    return angle


def parse_transform_rotation(transform: Optional[str]) -> Optional[float]:
    if not transform:
        return None
    if "rotate(" in transform:
        match = re.search(r"rotate\(([^)]+)\)", transform)
        if match:
            parts = [p for p in re.split(r"[ ,]+", match.group(1).strip()) if p]
            if parts:
                try:
                    return normalize_rotation(float(parts[0]))
                except Exception:
                    return None
    match = MATRIX_RE.search(transform)
    if match:
        parts = [p for p in re.split(r"[ ,]+", match.group(1).strip()) if p]
        if len(parts) >= 4:
            try:
                a = float(parts[0])
                b = float(parts[1])
                if abs(a) < 1e-12 and abs(b) < 1e-12:
                    return None
                angle = math.degrees(math.atan2(b, a))
                return normalize_rotation(angle)
            except Exception:
                return None
    return None


def read_text(elem: ET.Element) -> str:
    parts: List[str] = []
    if elem.text:
        parts.append(elem.text)
    for child in elem.iter():
        if child is elem:
            continue
        if tag_name(child) == "tspan" and child.text:
            parts.append(child.text)
    return "".join(parts).replace("\u00a0", " ").strip()


def has_cjk(text: str) -> bool:
    return bool(CJK_RE.search(text))

def font_family_is_theme(font_family: Optional[str]) -> bool:
    tokens = [t.strip().strip("\"'") for t in (font_family or "").split(",") if t.strip()]
    meaningful = []
    for token in tokens:
        token_lower = token.lower()
        if token_lower in GENERIC_FONTS:
            continue
        if token.isdigit():
            continue
        meaningful.append(token)
    if not meaningful:
        return False
    return all("msfontservice" in t.lower() for t in meaningful)


def pick_font_name(font_family: Optional[str], text: str, cjk_font: str) -> Optional[str]:
    tokens = [t.strip().strip("\"'") for t in (font_family or "").split(",") if t.strip()]
    for token in tokens:
        token_lower = token.lower()
        if token_lower in GENERIC_FONTS:
            continue
        if "msfontservice" in token_lower:
            continue
        if any(c.isalpha() for c in token):
            return token
    if has_cjk(text) and cjk_font:
        return cjk_font
    return None


def set_run_ea_font(run, font_name: str) -> None:
    if not font_name:
        return
    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        r_pr.append(ea)
    ea.set("typeface", font_name)


def group_text_lines(items: List[Dict[str, object]], line_tol: float) -> List[Dict[str, object]]:
    if not items:
        return []
    items = sorted(items, key=lambda it: (it["y"], it["x"]))
    lines: List[Dict[str, object]] = []
    current: List[Dict[str, object]] = []
    current_y = items[0]["y"]
    for item in items:
        if abs(item["y"] - current_y) <= line_tol:
            current.append(item)
        else:
            lines.append({"y": current_y, "items": current})
            current = [item]
            current_y = item["y"]
    if current:
        lines.append({"y": current_y, "items": current})
    return lines


PUNCT_NO_SPACE_BEFORE = set(",.;:!?%)]}»。，、：；？！％）】》")
PUNCT_NO_SPACE_AFTER = set("([{\"'“‘")


def should_insert_space(prev_item: Dict[str, object], curr_item: Dict[str, object]) -> bool:
    prev_text = str(prev_item.get("text", ""))
    curr_text = str(curr_item.get("text", ""))
    if not prev_text or not curr_text:
        return False
    if prev_text[-1].isspace() or curr_text[0].isspace():
        return False
    if has_cjk(prev_text) or has_cjk(curr_text):
        return False
    if curr_text[0] in PUNCT_NO_SPACE_BEFORE:
        return False
    if prev_text[-1] in PUNCT_NO_SPACE_AFTER:
        return False
    if not re.search(r"[A-Za-z0-9]", prev_text) or not re.search(r"[A-Za-z0-9]", curr_text):
        return False
    prev_x = float(prev_item.get("x") or 0.0)
    curr_x = float(curr_item.get("x") or 0.0)
    size = max(float(prev_item.get("font_size") or 0.0), float(curr_item.get("font_size") or 0.0))
    gap = curr_x - prev_x
    if size <= 0.0:
        return gap > 1.0
    return gap > size * 0.6


def assemble_line_text(items: List[Dict[str, object]]) -> str:
    parts: List[str] = []
    prev_item: Optional[Dict[str, object]] = None
    for item in items:
        text = str(item.get("text", ""))
        if not text:
            continue
        if prev_item and should_insert_space(prev_item, item):
            parts.append(" ")
        parts.append(text)
        prev_item = item
    return "".join(parts).strip()


def remove_semantic_layer(tree: ET.ElementTree) -> None:
    root = tree.getroot()
    parent_map: Dict[ET.Element, ET.Element] = {}
    for parent in root.iter():
        for child in list(parent):
            parent_map[child] = parent
    for elem in list(root.iter()):
        if tag_name(elem) == "g" and (elem.get("id") == "semantic-layer" or elem.get("data-type") == "semantic-layer"):
            parent = parent_map.get(elem)
            if parent is not None:
                parent.remove(elem)


def remove_chart_placeholders(tree: ET.ElementTree, placeholder_ids: List[str]) -> None:
    """Remove image placeholders that are charts from the SVG tree."""
    if not placeholder_ids:
        return
    
    root = tree.getroot()
    parent_map: Dict[ET.Element, ET.Element] = {}
    for parent in root.iter():
        for child in list(parent):
            parent_map[child] = parent
    
    # Find and remove placeholders by ID
    for elem in list(root.iter()):
        if tag_name(elem) == "g" and elem.get("id") in placeholder_ids:
            parent = parent_map.get(elem)
            if parent is not None:
                parent.remove(elem)


def render_background(svg_path: Path, out_png: Path, remove_placeholder_ids: Optional[List[str]] = None) -> None:
    if cairosvg is None:
        raise SystemExit("cairosvg is required. Install: pip install cairosvg")
    tree = ET.parse(svg_path)
    remove_semantic_layer(tree)
    
    # Remove chart placeholders so they don't appear in background
    if remove_placeholder_ids:
        remove_chart_placeholders(tree, remove_placeholder_ids)
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".svg") as tmp:
        tmp_path = Path(tmp.name)
        tree.write(tmp_path, encoding="utf-8", xml_declaration=True)
    try:
        cairosvg.svg2png(url=str(tmp_path), write_to=str(out_png))
    finally:
        tmp_path.unlink(missing_ok=True)


def extract_textboxes(svg_path: Path) -> Tuple[Dict[str, float], List[Dict[str, object]]]:
    tree = ET.parse(svg_path)
    root = tree.getroot()
    width = parse_length(root.get("width"))
    height = parse_length(root.get("height"))
    canvas = {"w": width, "h": height}

    parent_map: Dict[ET.Element, ET.Element] = {}
    for parent in root.iter():
        for child in list(parent):
            parent_map[child] = parent

    merge_groups: set = set()
    for g in root.iter():
        if tag_name(g) != "g":
            continue
        if g.get("data-type") != "text-group":
            continue
        if g.get("data-role") not in MERGE_GROUP_ROLES:
            continue
        child_groups = [c for c in list(g) if tag_name(c) == "g"]
        if child_groups and all(c.get("data-type") == "textbox" for c in child_groups):
            merge_groups.add(g)

    def has_merge_ancestor(elem: ET.Element) -> bool:
        parent = parent_map.get(elem)
        while parent is not None:
            if parent in merge_groups:
                return True
            parent = parent_map.get(parent)
        return False

    def collect_text_items(container: ET.Element) -> List[Dict[str, object]]:
        items = []
        for text in container.iter():
            if tag_name(text) != "text":
                continue
            content = read_text(text)
            if not content:
                continue
            tx, ty = parse_transform_xy(text.get("transform"))
            rotation = parse_transform_rotation(text.get("transform"))
            if tx == 0.0 and ty == 0.0:
                tx = parse_length(text.get("x"))
                ty = parse_length(text.get("y"))
            items.append(
                {
                    "text": content,
                    "x": tx,
                    "y": ty,
                    "font_size": parse_length(text.get("font-size")),
                    "font_family": text.get("font-family"),
                    "font_theme": font_family_is_theme(text.get("font-family")),
                    "fill": text.get("fill"),
                    "text_anchor": text.get("text-anchor"),
                    "rotation": rotation,
                }
            )
        return items

    def infer_box_rotation(items: List[Dict[str, object]]) -> float:
        rotations = []
        for item in items:
            rot = item.get("rotation")
            if isinstance(rot, (int, float)):
                rot_val = normalize_rotation(float(rot))
                if abs(rot_val) >= 1.0:
                    rotations.append(rot_val)
        if not rotations:
            return 0.0
        rotations.sort()
        return rotations[len(rotations) // 2]

    textboxes = []
    for g in root.iter():
        if tag_name(g) != "g":
            continue
        if g in merge_groups:
            texts = collect_text_items(g)
            textboxes.append(
                {
                    "id": g.get("id") or "",
                    "x": parse_length(g.get("data-x")),
                    "y": parse_length(g.get("data-y")),
                    "w": parse_length(g.get("data-w")),
                    "h": parse_length(g.get("data-h")),
                    "role": g.get("data-role") or "",
                    "texts": texts,
                    "rotation": infer_box_rotation(texts),
                }
            )
            continue
        if g.get("data-type") != "textbox":
            continue
        if has_merge_ancestor(g):
            continue
        texts = collect_text_items(g)
        textboxes.append(
            {
                "id": g.get("id") or "",
                "x": parse_length(g.get("data-x")),
                "y": parse_length(g.get("data-y")),
                "w": parse_length(g.get("data-w")),
                "h": parse_length(g.get("data-h")),
                "role": g.get("data-role") or "",
                "texts": texts,
                "rotation": infer_box_rotation(texts),
            }
        )
    return canvas, textboxes


def add_textbox(
    slide,
    tb: Dict[str, object],
    dpi: float,
    line_tol: float,
    box_pad: float,
    cjk_font: str,
) -> None:
    pad = max(0.0, float(box_pad))
    base_x = float(tb["x"]) - pad
    base_y = float(tb["y"]) - pad
    base_w = float(tb["w"]) + 2 * pad
    base_h = float(tb["h"]) + 2 * pad
    rotation = normalize_rotation(float(tb.get("rotation", 0.0) or 0.0))
    swap_dims = abs(abs(rotation) - 90.0) <= 2.0
    if swap_dims:
        width_px = base_h
        height_px = base_w
        center_x = base_x + base_w / 2.0
        center_y = base_y + base_h / 2.0
        left = Inches((center_x - width_px / 2.0) / dpi)
        top = Inches((center_y - height_px / 2.0) / dpi)
    else:
        width_px = base_w
        height_px = base_h
        left = Inches(base_x / dpi)
        top = Inches(base_y / dpi)
    width = Inches(width_px / dpi)
    height = Inches(height_px / dpi)
    shape = slide.shapes.add_textbox(left, top, width, height)
    if abs(rotation) >= 1.0:
        shape.rotation = rotation
    if tb.get("id"):
        shape.name = str(tb["id"])
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    text_items = tb.get("texts", [])
    if abs(abs(rotation) - 90.0) <= 2.0:
        tf.word_wrap = False
        ordered = sorted(text_items, key=lambda it: it["y"])
        lines = [{"items": ordered}]
    else:
        lines = group_text_lines(text_items, line_tol=line_tol)
    for idx, line in enumerate(lines):
        line_items = sorted(line["items"], key=lambda it: it["x"])
        line_text = assemble_line_text(line_items)
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        if not line_items:
            continue
        run = p.runs[0] if p.runs else p.add_run()
        style = line_items[0]
        font_size = style.get("font_size") or 0.0
        if font_size:
            run.font.size = Pt(font_size * 0.75)
        use_theme_font = bool(style.get("font_theme"))
        font_name = None
        if not use_theme_font:
            font_name = pick_font_name(style.get("font_family"), line_text, cjk_font)
            if font_name:
                run.font.name = font_name
        if has_cjk(line_text) and not use_theme_font:
            ea_font = font_name or cjk_font
            if ea_font:
                set_run_ea_font(run, ea_font)
        color = parse_color(style.get("fill"))
        if color:
            run.font.color.rgb = color
        anchor = style.get("text_anchor")
        if anchor == "middle":
            p.alignment = PP_ALIGN.CENTER
        elif anchor == "end":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT


def load_config(config_path: Optional[Path]) -> Dict[str, Any]:
    """Load config.json for API settings."""
    if not config_path or not config_path.exists():
        return {}
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)


def load_placeholders(json_path: Optional[Path]) -> Dict[Tuple[str, str], Dict[str, Any]]:
    """Load image_placeholders.json and build {(svg_file, placeholder_id): entry} mapping."""
    if not json_path or not json_path.exists():
        return {}
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    mapping: Dict[Tuple[str, str], Dict[str, Any]] = {}
    for entry in data:
        key = (entry.get("svg_file", ""), entry.get("placeholder_id", ""))
        mapping[key] = entry
    return mapping


def parse_matrix(transform: str) -> Tuple[float, float, float, float, float, float]:
    """Parse matrix(a b c d e f) transform, returns (a, b, c, d, e, f)."""
    match = re.search(r"matrix\(([^)]+)\)", transform)
    if match:
        parts = [p for p in re.split(r"[ ,]+", match.group(1).strip()) if p]
        if len(parts) == 6:
            try:
                return tuple(float(p) for p in parts)
            except Exception:
                pass
    return (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)


def parse_scale(transform: str) -> Tuple[float, float]:
    """Parse scale(sx sy) transform, returns (sx, sy)."""
    match = re.search(r"scale\(([^)]+)\)", transform)
    if match:
        parts = [p for p in re.split(r"[ ,]+", match.group(1).strip()) if p]
        if parts:
            try:
                sx = float(parts[0])
                sy = float(parts[1]) if len(parts) > 1 else sx
                return (sx, sy)
            except Exception:
                pass
    return (1.0, 1.0)


def extract_chart_placeholders(svg_path: Path) -> List[Dict[str, Any]]:
    """Extract image placeholders with their bounding boxes from SVG.
    
    Correctly computes final position by applying transforms from element to root.
    """
    tree = ET.parse(svg_path)
    root = tree.getroot()
    
    # Build parent map
    parent_map: Dict[ET.Element, ET.Element] = {}
    for parent in root.iter():
        for child in parent:
            parent_map[child] = parent
    
    placeholders = []
    
    for elem in root.iter():
        if tag_name(elem) != "g":
            continue
        if elem.get("data-role") != "image-placeholder":
            continue
        
        placeholder_id = elem.get("id", "")
        caption = elem.get("data-caption", "")
        
        # Find dimensions from rect or image child
        rect_x, rect_y, rect_w, rect_h = 0.0, 0.0, 0.0, 0.0
        for child in elem:
            child_tag = tag_name(child)
            if child_tag in ("rect", "image"):
                rect_x = parse_length(child.get("x"))
                rect_y = parse_length(child.get("y"))
                rect_w = parse_length(child.get("width"))
                rect_h = parse_length(child.get("height"))
                break
        
        # Collect all transforms from element up to root
        transforms: List[str] = []
        current = elem
        while current is not None:
            t = current.get("transform", "")
            if t:
                transforms.append(t)
            current = parent_map.get(current)
        
        # Apply transforms from innermost (element) to outermost (root)
        # Start with rect position and dimensions
        x, y = rect_x, rect_y
        w, h = rect_w, rect_h
        
        for t in transforms:
            # Apply scale first
            sx, sy = parse_scale(t)
            x *= sx
            y *= sy
            w *= sx
            h *= sy
            
            # Apply matrix (includes scale and translate)
            a, b, c, d, e, f = parse_matrix(t)
            if not (a == 1.0 and b == 0.0 and c == 0.0 and d == 1.0 and e == 0.0 and f == 0.0):
                # Transform point: new_x = a*x + c*y + e, new_y = b*x + d*y + f
                # For non-rotated matrices (b=0, c=0): new_x = a*x + e, new_y = d*y + f
                new_x = a * x + e
                new_y = d * y + f
                x, y = new_x, new_y
                w *= a
                h *= d
            
            # Apply translate (if not already in matrix)
            if "translate(" in t and "matrix(" not in t:
                tx, ty = parse_transform_xy(t)
                x += tx
                y += ty
        
        placeholders.append({
            "id": placeholder_id,
            "caption": caption,
            "x": x,
            "y": y,
            "w": w,
            "h": h,
        })
    
    return placeholders


CHART_CODE_PROMPT = """You are a Python code generator. Generate python-pptx code to create a native PowerPoint chart based on this description:

Chart Description: {caption}

Requirements:
1. Generate ONLY a Python function named `add_chart_to_slide(slide, left, top, width, height)`
2. The function should use python-pptx to add a chart to the slide
3. Use the provided left, top, width, height parameters (already in Inches)
4. Import statements are NOT needed - CategoryChartData, XyChartData, BubbleChartData, XL_CHART_TYPE, RGBColor, Pt are available
5. Extract data values and labels from the description
6. Return ONLY the function code, no explanations, no markdown code blocks
7. Make sure the chart matches the description (bar chart, line chart, pie chart, etc.)

IMPORTANT RESTRICTIONS - DO NOT use these unsupported attributes:
- chart.fill (Chart has no fill attribute)
- chart.plot_area.format.fill (not supported)
- chart.chart_area (not supported)
- Any background color settings on the chart itself
- Do not reference undefined classes or missing imports

Only use these SAFE attributes:
- chart.series[i].format.fill.solid() and .fore_color.rgb for bar/column colors
- chart.value_axis / chart.category_axis for axis settings
- chart.has_legend for legend toggle
- axis.tick_labels.font for font settings

Example output format:
def add_chart_to_slide(slide, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = ['A', 'B', 'C']
    chart_data.add_series('Series 1', (1, 2, 3))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    ).chart
    # Set bar color
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0, 128, 128)
"""


def generate_chart_code(
    caption: str,
    api_key: Optional[str] = None,
    base_url: Optional[str] = None,
    model: Optional[str] = None,
    max_tokens: int = 800,
    temperature: float = 0.2,
    error_hint: Optional[str] = None,
) -> Optional[str]:
    """Call code model to generate python-pptx chart code from caption."""
    if OpenAI is None:
        print("Warning: openai package not installed. Skipping chart generation.")
        return None
    
    key = api_key or os.environ.get("OPENAI_API_KEY")
    if not key:
        print("Warning: API key not set. Skipping chart generation.")
        return None
    
    # Use provided model or default to gpt-4o
    use_model = model or "gpt-4o"
    
    try:
        # Create client with optional base_url for compatible APIs (e.g., SiliconFlow)
        client_kwargs = {"api_key": key}
        if base_url:
            client_kwargs["base_url"] = base_url + "/v1"
        client = OpenAI(**client_kwargs)
        user_prompt = CHART_CODE_PROMPT.format(caption=caption)
        if error_hint:
            user_prompt = (
                user_prompt
                + "\n\nPrevious attempt failed with error:\n"
                + str(error_hint)
                + "\nFix the code and regenerate."
            )
        response = client.chat.completions.create(
            model=use_model,
            messages=[
                {"role": "system", "content": "You are a Python code generator specializing in python-pptx charts."},
                {"role": "user", "content": user_prompt},
            ],
            temperature=temperature,
            max_tokens=max_tokens,
        )
        code = response.choices[0].message.content.strip()
        # Remove markdown code blocks if present
        if code.startswith("```"):
            lines = code.split("\n")
            code = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
        return code
    except Exception as e:
        print(f"Warning: Failed to generate chart code: {e}")
        return None


def generate_chart_codes(
    chart_requests: List[Dict[str, Any]],
    chart_config: Dict[str, Any],
) -> Dict[Tuple[str, str], Optional[str]]:
    if not chart_requests:
        return {}
    if OpenAI is None:
        print("Warning: openai package not installed. Skipping chart generation.")
        return {}
    api_key = chart_config.get("api_key") or os.environ.get("OPENAI_API_KEY")
    if not api_key:
        print("Warning: API key not set. Skipping chart generation.")
        return {}

    workers = int(chart_config.get("chart_workers", 1))
    qps = float(chart_config.get("chart_qps", 0))
    retries = int(chart_config.get("chart_retries", 0))
    limiter = RateLimiter(qps) if qps > 0 else None
    results: Dict[Tuple[str, str], Optional[str]] = {}

    def worker(req: Dict[str, Any]) -> Optional[str]:
        caption = req["caption"]
        for _ in range(retries + 1):
            if limiter:
                limiter.acquire()
            code = generate_chart_code(
                caption,
                api_key=api_key,
                base_url=chart_config.get("base_url"),
                model=chart_config.get("chart_model"),
                max_tokens=chart_config.get("chart_max_tokens", 800),
                temperature=chart_config.get("chart_temperature", 0.2),
            )
            if code:
                return code
        return None

    if workers <= 1:
        for req in chart_requests:
            key = (req["svg_name"], req["placeholder_id"])
            results[key] = worker(req)
    else:
        from concurrent.futures import ThreadPoolExecutor, as_completed

        with ThreadPoolExecutor(max_workers=workers) as executor:
            future_map = {
                executor.submit(worker, req): (req["svg_name"], req["placeholder_id"])
                for req in chart_requests
            }
            for future in as_completed(future_map):
                key = future_map[future]
                try:
                    results[key] = future.result()
                except Exception:
                    results[key] = None

    return results


def execute_chart_code(
    code: str,
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
) -> Tuple[bool, Optional[str]]:
    """Safely execute generated chart code."""
    # Prepare execution environment with necessary imports
    exec_globals = {
        "CategoryChartData": CategoryChartData,
        "XL_CHART_TYPE": XL_CHART_TYPE,
        "RGBColor": RGBColor,
        "Pt": Pt,
        "Inches": Inches,
    }
    if XyChartData is not None:
        exec_globals["XyChartData"] = XyChartData
    if BubbleChartData is not None:
        exec_globals["BubbleChartData"] = BubbleChartData
    exec_locals: Dict[str, Any] = {}
    
    try:
        # Execute the function definition
        exec(code, exec_globals, exec_locals)
        
        # Call the generated function
        if "add_chart_to_slide" in exec_locals:
            exec_locals["add_chart_to_slide"](slide, left, top, width, height)
            return True, None
        else:
            return False, "Generated code does not contain add_chart_to_slide function."
    except Exception as e:
        return False, str(e)


def build_pptx(
    svg_paths: List[Path],
    out_pptx: Path,
    dpi: float,
    line_tol: float,
    box_pad: float,
    cjk_font: str,
    placeholders_map: Optional[Dict[Tuple[str, str], Dict[str, Any]]] = None,
    skip_charts: bool = False,
    chart_config: Optional[Dict[str, Any]] = None,
) -> None:
    if not svg_paths:
        print("No SVG files to convert.")
        return

    placeholders_map = placeholders_map or {}
    chart_code_map: Dict[Tuple[str, str], Optional[str]] = {}
    if not skip_charts and placeholders_map:
        chart_requests: List[Dict[str, Any]] = []
        for svg_path in svg_paths:
            svg_name = svg_path.name
            chart_placeholders = extract_chart_placeholders(svg_path)
            for ph in chart_placeholders:
                placeholder_id = ph["id"]
                entry = placeholders_map.get((svg_name, placeholder_id))
                if not entry or not entry.get("is_chart", False):
                    continue
                caption = entry.get("caption", "")
                if not caption:
                    continue
                chart_requests.append(
                    {
                        "svg_name": svg_name,
                        "placeholder_id": placeholder_id,
                        "caption": caption,
                    }
                )
        if chart_requests:
            cfg = chart_config or {}
            workers = int(cfg.get("chart_workers", 1))
            print(f"Generating chart code ({len(chart_requests)} charts, workers={workers})...")
            chart_code_map = generate_chart_codes(chart_requests, cfg)
    first_canvas, _ = extract_textboxes(svg_paths[0])
    if not first_canvas["w"] or not first_canvas["h"]:
        first_canvas["w"], first_canvas["h"] = 1920.0, 1080.0

    prs = Presentation()
    prs.slide_width = Inches(first_canvas["w"] / dpi)
    prs.slide_height = Inches(first_canvas["h"] / dpi)

    for svg_path in svg_paths:
        canvas, textboxes = extract_textboxes(svg_path)
        if canvas["w"] and canvas["h"]:
            if abs(canvas["w"] - first_canvas["w"]) > 1 or abs(canvas["h"] - first_canvas["h"]) > 1:
                print(f"Warning: {svg_path.name} size differs from first slide.")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        svg_name = svg_path.name

        # Collect chart placeholder IDs to remove from background
        chart_placeholder_ids: List[str] = []
        chart_infos: List[Dict[str, Any]] = []
        
        if not skip_charts and placeholders_map:
            chart_placeholders = extract_chart_placeholders(svg_path)
            
            for ph in chart_placeholders:
                placeholder_id = ph["id"]
                key = (svg_name, placeholder_id)
                entry = placeholders_map.get(key)
                
                if entry and entry.get("is_chart", False):
                    chart_placeholder_ids.append(placeholder_id)
                    chart_infos.append({
                        "ph": ph,
                        "entry": entry,
                    })

        # Render background (with chart placeholders removed)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_png:
            tmp_png_path = Path(tmp_png.name)
        try:
            render_background(svg_path, tmp_png_path, remove_placeholder_ids=chart_placeholder_ids)
            slide.shapes.add_picture(
                str(tmp_png_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
        finally:
            tmp_png_path.unlink(missing_ok=True)

        # Add native charts
        for info in chart_infos:
            ph = info["ph"]
            entry = info["entry"]
            placeholder_id = ph["id"]
            caption = entry.get("caption", "")
            
            if not caption:
                continue
            
            # Calculate chart position in inches
            left = Inches(ph["x"] / dpi)
            top = Inches(ph["y"] / dpi)
            width = Inches(ph["w"] / dpi) if ph["w"] > 0 else Inches(5)
            height = Inches(ph["h"] / dpi) if ph["h"] > 0 else Inches(3)
            
            print(f"Generating native chart for {svg_name}/{placeholder_id}...")
            print(f"  Position: x={ph['x']:.1f}, y={ph['y']:.1f}, w={ph['w']:.1f}, h={ph['h']:.1f}")
            
            cfg = chart_config or {}
            exec_retries = cfg.get("chart_exec_retries", cfg.get("chart_retries", 2))
            try:
                exec_retries = int(exec_retries)
            except Exception:
                exec_retries = 2
            max_attempts = None if exec_retries < 0 else exec_retries + 1
            code = chart_code_map.get((svg_name, placeholder_id))
            last_error: Optional[str] = None
            success = False
            attempt = 0
            while True:
                if max_attempts is not None and attempt >= max_attempts:
                    break
                if not code:
                    code = generate_chart_code(
                        caption,
                        api_key=cfg.get("api_key"),
                        base_url=cfg.get("base_url"),
                        model=cfg.get("chart_model"),
                        max_tokens=cfg.get("chart_max_tokens", 800),
                        temperature=cfg.get("chart_temperature", 0.2),
                        error_hint=last_error,
                    )
                if not code:
                    last_error = "code generation failed"
                else:
                    success, last_error = execute_chart_code(code, slide, left, top, width, height)
                    if success:
                        break
                code = None
                attempt += 1
            if success:
                print("  ✓ Chart created successfully")
            else:
                raise SystemExit(
                    f"Chart creation failed for {svg_name}/{placeholder_id}: {last_error or 'unknown error'}"
                )

        # Add textboxes
        for tb in textboxes:
            if tb.get("w", 0) <= 0 or tb.get("h", 0) <= 0:
                continue
            add_textbox(slide, tb, dpi=dpi, line_tol=line_tol, box_pad=box_pad, cjk_font=cjk_font)

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


def main() -> None:
    parser = argparse.ArgumentParser(description="Convert semantic SVG to PPTX (slide only).")
    parser.add_argument("--input", required=True, help="Semantic SVG file or directory.")
    parser.add_argument("--output", required=True, help="Output PPTX path.")
    parser.add_argument("--dpi", type=float, default=96.0, help="SVG pixel DPI.")
    parser.add_argument("--line-tol", type=float, default=2.0, help="Line grouping tolerance (px).")
    parser.add_argument("--box-pad", type=float, default=12.0, help="Textbox padding in px.")
    parser.add_argument("--cjk-font", default="PingFang SC", help="Font for CJK text.")
    parser.add_argument("--placeholders", help="Path to image_placeholders.json for chart detection.")
    parser.add_argument("--skip-charts", action="store_true", help="Skip native chart reconstruction.")
    parser.add_argument("--config", help="Path to config.json for API settings.")
    parser.add_argument("--api-key", help="API key (overrides config.json).")
    parser.add_argument("--chart-workers", type=int, default=None, help="Parallel chart code workers.")
    parser.add_argument("--chart-qps", type=float, default=None, help="Chart code QPS limit.")
    parser.add_argument("--chart-retries", type=int, default=None, help="Chart code retries.")
    args = parser.parse_args()

    input_path = Path(args.input)
    if input_path.is_dir():
        # Collect all SVG files and sort naturally (幻灯片1, 幻灯片2, ..., 幻灯片10)
        all_svgs = list(input_path.glob("*.SVG")) + list(input_path.glob("*.svg"))
        svg_paths = sorted(
            [p for p in all_svgs if not p.name.startswith("._")],
            key=natural_sort_key
        )
    else:
        svg_paths = [input_path]

    # Load config for API settings
    chart_config: Dict[str, Any] = {}
    if args.config:
        chart_config = load_config(Path(args.config))
        if chart_config:
            print(f"Loaded config: model={chart_config.get('chart_model')}, base_url={chart_config.get('base_url')}")
    
    # Command-line api-key overrides config
    if args.api_key:
        chart_config["api_key"] = args.api_key
    if args.chart_workers is not None:
        chart_config["chart_workers"] = args.chart_workers
    if args.chart_qps is not None:
        chart_config["chart_qps"] = args.chart_qps
    if args.chart_retries is not None:
        chart_config["chart_retries"] = args.chart_retries

    # Load placeholders mapping for chart detection
    placeholders_map = {}
    if args.placeholders:
        placeholders_map = load_placeholders(Path(args.placeholders))
        if placeholders_map:
            chart_count = sum(1 for v in placeholders_map.values() if v.get("is_chart"))
            print(f"Loaded {len(placeholders_map)} placeholders ({chart_count} charts)")

    build_pptx(
        svg_paths,
        Path(args.output),
        dpi=args.dpi,
        line_tol=args.line_tol,
        box_pad=args.box_pad,
        cjk_font=args.cjk_font,
        placeholders_map=placeholders_map,
        skip_charts=args.skip_charts,
        chart_config=chart_config,
    )
    print(f"Wrote: {args.output}")


if __name__ == "__main__":
    main()
