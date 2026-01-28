#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SVG to PPTX Pro: Convert SVG elements to editable PPTX shapes.

Unlike svg_to_pptx_slide.py which renders SVG as a background image,
this script converts each SVG element (rect, circle, path, image, text, etc.)
to native PPTX shapes, preserving editability and layer order.
"""

import argparse
import base64
import io
import json
import math
import os
import re
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

# Optional: openai for chart generation
try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

# Chart data types
try:
    from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
except ImportError:
    CategoryChartData = None
    XyChartData = None
    BubbleChartData = None

try:
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    XL_CHART_TYPE = None

# Optional: svgpathtools for complex path parsing
try:
    from svgpathtools import parse_path, Line, CubicBezier, QuadraticBezier, Arc
    HAS_SVGPATHTOOLS = True
except ImportError:
    HAS_SVGPATHTOOLS = False
    parse_path = None

# Optional: PIL for image cropping
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# Optional: CairoSVG for rasterizing non-rect clip paths
try:
    import cairosvg
    HAS_CAIROSVG = True
except ImportError:
    HAS_CAIROSVG = False

# ============================================================================
# Constants and Regex Patterns
# ============================================================================

SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"

# EMU conversion: 1 inch = 914400 EMU
EMU_PER_INCH = 914400

# Default DPI for SVG (pixels per inch)
DEFAULT_DPI = 96.0

# Regex patterns
NATURAL_SORT_RE = re.compile(r"(\d+)")
MATRIX_RE = re.compile(r"matrix\s*\(\s*([^)]+)\s*\)")
TRANSLATE_RE = re.compile(r"translate\s*\(\s*([^)]+)\s*\)")
SCALE_RE = re.compile(r"scale\s*\(\s*([^)]+)\s*\)")
ROTATE_RE = re.compile(r"rotate\s*\(\s*([^)]+)\s*\)")
CJK_RE = re.compile(r"[\u3400-\u9fff\u3000-\u303f\u3040-\u30ff\u31f0-\u31ff\uac00-\ud7af]")

# Named colors mapping (common SVG/CSS colors)
NAMED_COLORS = {
    "black": (0, 0, 0),
    "white": (255, 255, 255),
    "red": (255, 0, 0),
    "green": (0, 128, 0),
    "blue": (0, 0, 255),
    "yellow": (255, 255, 0),
    "cyan": (0, 255, 255),
    "magenta": (255, 0, 255),
    "gray": (128, 128, 128),
    "grey": (128, 128, 128),
    "silver": (192, 192, 192),
    "maroon": (128, 0, 0),
    "olive": (128, 128, 0),
    "lime": (0, 255, 0),
    "aqua": (0, 255, 255),
    "teal": (0, 128, 128),
    "navy": (0, 0, 128),
    "fuchsia": (255, 0, 255),
    "purple": (128, 0, 128),
    "orange": (255, 165, 0),
    "pink": (255, 192, 203),
    "brown": (165, 42, 42),
    "transparent": None,
    "none": None,
}

GENERIC_FONTS = {
    "sans-serif", "serif", "monospace", "system-ui",
    "ui-sans-serif", "ui-serif", "ui-monospace",
}


# ============================================================================
# Utility Functions
# ============================================================================

def tag_name(elem: ET.Element) -> str:
    """Extract local tag name without namespace."""
    return elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag


def natural_sort_key(path: Path) -> List[Any]:
    """Generate sort key for natural sorting (e.g., 幻灯片1, 幻灯片2, ..., 幻灯片10)."""
    parts = NATURAL_SORT_RE.split(path.name)
    return [int(p) if p.isdigit() else p.lower() for p in parts]


def parse_length(val: Optional[str], default: float = 0.0) -> float:
    """Parse SVG length value (px, pt, mm, cm, in, %)."""
    if not val:
        return default
    raw = val.strip()
    
    # Handle percentage (will need context for actual conversion)
    if raw.endswith("%"):
        try:
            return float(raw[:-1]) / 100.0 * default if default else 0.0
        except ValueError:
            return default
    
    # Remove unit suffixes
    for suffix in ("px", "pt", "mm", "cm", "in", "em", "ex"):
        if raw.endswith(suffix):
            raw = raw[:-len(suffix)]
            break
    
    try:
        return float(raw)
    except ValueError:
        return default


def parse_color(val: Optional[str]) -> Optional[RGBColor]:
    """Parse SVG color value to RGBColor."""
    if not val:
        return None
    
    v = val.strip().lower()
    
    # Handle 'none' and 'transparent'
    if v in ("none", "transparent", "inherit", "currentcolor"):
        return None
    
    # Named colors
    if v in NAMED_COLORS:
        rgb = NAMED_COLORS[v]
        return RGBColor(*rgb) if rgb else None
    
    # Hex colors
    if v.startswith("#"):
        hex_color = v[1:]
        try:
            if len(hex_color) == 3:
                r = int(hex_color[0] * 2, 16)
                g = int(hex_color[1] * 2, 16)
                b = int(hex_color[2] * 2, 16)
            elif len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
            elif len(hex_color) == 8:
                # RGBA - ignore alpha for now
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
            else:
                return None
            return RGBColor(r, g, b)
        except ValueError:
            return None
    
    # rgb() and rgba() functions
    rgb_match = re.match(r"rgba?\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", v)
    if rgb_match:
        try:
            r = int(rgb_match.group(1))
            g = int(rgb_match.group(2))
            b = int(rgb_match.group(3))
            return RGBColor(
                min(255, max(0, r)),
                min(255, max(0, g)),
                min(255, max(0, b))
            )
        except ValueError:
            return None
    
    return None


def parse_opacity(val: Optional[str]) -> float:
    """Parse opacity value (0.0 to 1.0)."""
    if not val:
        return 1.0
    try:
        opacity = float(val.strip())
        return max(0.0, min(1.0, opacity))
    except ValueError:
        return 1.0


def normalize_rotation(angle: float) -> float:
    """Normalize rotation angle to -180 to 180 degrees."""
    while angle <= -180:
        angle += 360
    while angle > 180:
        angle -= 360
    return angle


def parse_transform_rotation(transform: Optional[str]) -> Optional[float]:
    """Extract rotation angle from SVG transform attribute."""
    if not transform:
        return None
    if "rotate(" in transform:
        match = re.search(r"rotate\s*\(\s*([^)]+)\s*\)", transform)
        if match:
            parts = [p for p in re.split(r"[\s,]+", match.group(1).strip()) if p]
            if parts:
                try:
                    return normalize_rotation(float(parts[0]))
                except Exception:
                    return None
    match = MATRIX_RE.search(transform)
    if match:
        parts = [p for p in re.split(r"[\s,]+", match.group(1).strip()) if p]
        if len(parts) >= 4:
            try:
                a = float(parts[0])
                b = float(parts[1])
                if abs(a) < 1e-12 and abs(b) < 1e-12:
                    return None
                angle = math.degrees(math.atan2(b, a))
                return normalize_rotation(angle)
            except Exception:
                return None
    return None


def parse_transform_xy(transform: Optional[str]) -> Tuple[float, float]:
    """Extract translation (x, y) from SVG transform attribute."""
    if not transform:
        return (0.0, 0.0)
    match = MATRIX_RE.search(transform)
    if match:
        parts = [p for p in re.split(r"[\s,]+", match.group(1).strip()) if p]
        if len(parts) == 6:
            try:
                return (float(parts[4]), float(parts[5]))
            except Exception:
                return (0.0, 0.0)
    match = TRANSLATE_RE.search(transform)
    if match:
        parts = [p for p in re.split(r"[\s,]+", match.group(1).strip()) if p]
        if parts:
            try:
                x = float(parts[0])
                y = float(parts[1]) if len(parts) > 1 else 0.0
                return (x, y)
            except Exception:
                return (0.0, 0.0)
    return (0.0, 0.0)


def read_text_content(elem: ET.Element) -> str:
    """Read text content from SVG text element including tspans."""
    parts: List[str] = []
    if elem.text:
        parts.append(elem.text)
    for child in elem.iter():
        if child is elem:
            continue
        if tag_name(child) == "tspan" and child.text:
            parts.append(child.text)
    return "".join(parts).replace("\u00a0", " ").strip()


def has_cjk(text: str) -> bool:
    """Check if text contains CJK characters."""
    return bool(CJK_RE.search(text))

def font_family_is_theme(font_family: Optional[str]) -> bool:
    tokens = [t.strip().strip("\"'") for t in (font_family or "").split(",") if t.strip()]
    meaningful = []
    for token in tokens:
        token_lower = token.lower()
        if token_lower in GENERIC_FONTS:
            continue
        if token.isdigit():
            continue
        meaningful.append(token)
    if not meaningful:
        return False
    return all("msfontservice" in t.lower() for t in meaningful)


def pick_font_name(font_family: Optional[str], text: str, cjk_font: str) -> Optional[str]:
    """Pick appropriate font name from font-family string."""
    tokens = [t.strip().strip("\"'") for t in (font_family or "").split(",") if t.strip()]
    for token in tokens:
        token_lower = token.lower()
        if token_lower in GENERIC_FONTS:
            continue
        if "msfontservice" in token_lower:
            continue
        if any(c.isalpha() for c in token):
            return token
    if has_cjk(text) and cjk_font:
        return cjk_font
    return None


def set_run_ea_font(run, font_name: str) -> None:
    """Set East Asian font for a text run."""
    if not font_name:
        return
    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        r_pr.append(ea)
    ea.set("typeface", font_name)


def group_text_lines(items: List[Dict[str, object]], line_tol: float) -> List[Dict[str, object]]:
    """Group text items into lines based on y-coordinate proximity."""
    if not items:
        return []
    items = sorted(items, key=lambda it: (it["y"], it["x"]))
    lines: List[Dict[str, object]] = []
    current: List[Dict[str, object]] = []
    current_y = items[0]["y"]
    for item in items:
        if abs(item["y"] - current_y) <= line_tol:
            current.append(item)
        else:
            lines.append({"y": current_y, "items": current})
            current = [item]
            current_y = item["y"]
    if current:
        lines.append({"y": current_y, "items": current})
    return lines


PUNCT_NO_SPACE_BEFORE = set(",.;:!?%)]}»。，、：；？！％）】》")
PUNCT_NO_SPACE_AFTER = set("([{\"'“‘")


def should_insert_space(prev_item: Dict[str, object], curr_item: Dict[str, object]) -> bool:
    prev_text = str(prev_item.get("text", ""))
    curr_text = str(curr_item.get("text", ""))
    if not prev_text or not curr_text:
        return False
    if prev_text[-1].isspace() or curr_text[0].isspace():
        return False
    if has_cjk(prev_text) or has_cjk(curr_text):
        return False
    if curr_text[0] in PUNCT_NO_SPACE_BEFORE:
        return False
    if prev_text[-1] in PUNCT_NO_SPACE_AFTER:
        return False
    if not re.search(r"[A-Za-z0-9]", prev_text) or not re.search(r"[A-Za-z0-9]", curr_text):
        return False
    prev_x = float(prev_item.get("x") or 0.0)
    curr_x = float(curr_item.get("x") or 0.0)
    size = max(float(prev_item.get("font_size") or 0.0), float(curr_item.get("font_size") or 0.0))
    gap = curr_x - prev_x
    if size <= 0.0:
        return gap > 1.0
    return gap > size * 0.6


def assemble_line_text(items: List[Dict[str, object]]) -> str:
    parts: List[str] = []
    prev_item: Optional[Dict[str, object]] = None
    for item in items:
        text = str(item.get("text", ""))
        if not text:
            continue
        if prev_item and should_insert_space(prev_item, item):
            parts.append(" ")
        parts.append(text)
        prev_item = item
    return "".join(parts).strip()


# ============================================================================
# Image Placeholder Extraction (Correct Transform Chain Handling)
# ============================================================================

def parse_matrix_simple(transform: str) -> Tuple[float, float, float, float, float, float]:
    """Parse matrix(a b c d e f) transform, returns (a, b, c, d, e, f)."""
    match = re.search(r"matrix\(([^)]+)\)", transform)
    if match:
        parts = [p for p in re.split(r"[ ,]+", match.group(1).strip()) if p]
        if len(parts) == 6:
            try:
                return tuple(float(p) for p in parts)
            except Exception:
                pass
    return (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)

def apply_transform_chain(
    x: float,
    y: float,
    w: float,
    h: float,
    transforms: List[str],
) -> Tuple[float, float, float, float]:
    """Apply a list of SVG transforms (element -> root) to a rect."""
    for t in transforms:
        sx, sy = parse_scale_simple(t)
        x *= sx
        y *= sy
        w *= sx
        h *= sy

        a, b, c, d, e, f = parse_matrix_simple(t)
        if not (a == 1.0 and b == 0.0 and c == 0.0 and d == 1.0 and e == 0.0 and f == 0.0):
            new_x = a * x + c * y + e
            new_y = b * x + d * y + f
            x, y = new_x, new_y
            w *= a
            h *= d

        if "translate(" in t and "matrix(" not in t:
            tx, ty = parse_transform_xy(t)
            x += tx
            y += ty
    return x, y, w, h


def clip_path_is_rect(clip_elem: ET.Element) -> bool:
    """Return True if clipPath contains a single rect child and nothing else."""
    children = [child for child in list(clip_elem) if isinstance(child.tag, str)]
    if len(children) != 1:
        return False
    return tag_name(children[0]) == "rect"


def rasterize_clipped_placeholder(
    placeholder: Dict[str, Any],
    converter: "CoordinateConverter",
) -> Optional[Tuple[str, Tuple[int, int, int, int]]]:
    """Rasterize a non-rect clipped placeholder to a temporary PNG.

    Returns (png_path, bbox) where bbox is (left, top, right, bottom) in SVG px.
    """
    if not HAS_CAIROSVG:
        print("Warning: cairosvg not available; skipping non-rect clip rasterization.")
        return None
    if not HAS_PIL:
        print("Warning: PIL not available; skipping non-rect clip rasterization.")
        return None

    svg_width = int(round(converter.svg_width))
    svg_height = int(round(converter.svg_height))
    if svg_width <= 0 or svg_height <= 0:
        return None

    svg_root = ET.Element(
        "svg",
        {
            "xmlns": SVG_NS,
            "xmlns:xlink": XLINK_NS,
            "width": str(svg_width),
            "height": str(svg_height),
            "viewBox": f"0 0 {svg_width} {svg_height}",
        },
    )
    defs = ET.SubElement(svg_root, "defs")
    for clip_xml in placeholder.get("clip_defs", {}).values():
        try:
            defs.append(ET.fromstring(clip_xml))
        except ET.ParseError:
            continue

    parent = svg_root
    for wrapper in placeholder.get("clip_chain", []):
        group = ET.SubElement(parent, "g")
        if wrapper.get("transform"):
            group.set("transform", wrapper["transform"])
        if wrapper.get("clip_path"):
            group.set("clip-path", wrapper["clip_path"])
        if wrapper.get("opacity"):
            group.set("opacity", wrapper["opacity"])
        parent = group

    elem_xml = placeholder.get("elem_xml")
    if not elem_xml:
        return None
    try:
        parent.append(ET.fromstring(elem_xml))
    except ET.ParseError:
        return None

    svg_bytes = ET.tostring(svg_root, encoding="utf-8", xml_declaration=True)
    try:
        png_bytes = cairosvg.svg2png(
            bytestring=svg_bytes,
            output_width=svg_width,
            output_height=svg_height,
        )
    except Exception:
        return None

    img = Image.open(io.BytesIO(png_bytes)).convert("RGBA")
    alpha = img.split()[-1]
    bbox = alpha.getbbox()
    if not bbox:
        return None

    cropped = img.crop(bbox)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
        cropped.save(tmp, format="PNG")
        tmp_path = tmp.name

    return tmp_path, bbox


def parse_scale_simple(transform: str) -> Tuple[float, float]:
    """Parse scale(sx sy) transform, returns (sx, sy)."""
    match = re.search(r"scale\(([^)]+)\)", transform)
    if match:
        parts = [p for p in re.split(r"[ ,]+", match.group(1).strip()) if p]
        if parts:
            try:
                sx = float(parts[0])
                sy = float(parts[1]) if len(parts) > 1 else sx
                return (sx, sy)
            except Exception:
                pass
    return (1.0, 1.0)


def extract_image_placeholders(svg_path: Path) -> List[Dict[str, Any]]:
    """Extract image placeholders with their bounding boxes from SVG.
    
    Correctly computes final position by applying transforms from element to root.
    This handles complex nested transforms like scale(8000) inside matrix(0.0001 ...).
    """
    tree = ET.parse(svg_path)
    root = tree.getroot()
    
    # Build parent map
    parent_map: Dict[ET.Element, ET.Element] = {}
    for parent in root.iter():
        for child in parent:
            parent_map[child] = parent
    
    placeholders = []
    clip_rects_map: Dict[str, Dict[str, float]] = {}
    clip_units: Dict[str, str] = {}
    clip_is_rect: Dict[str, bool] = {}
    clip_xml_map: Dict[str, str] = {}

    for elem in root.iter():
        if tag_name(elem) != "clipPath":
            continue
        clip_id = elem.get("id")
        if not clip_id:
            continue
        clip_units[clip_id] = elem.get("clipPathUnits", "userSpaceOnUse")
        clip_is_rect[clip_id] = clip_path_is_rect(elem)
        clip_xml_map[clip_id] = ET.tostring(elem, encoding="utf-8").decode("utf-8")
        if clip_is_rect[clip_id]:
            rect = None
            for child in list(elem):
                if tag_name(child) == "rect":
                    rect = child
                    break
            if rect is None:
                continue
            clip_rects_map[clip_id] = {
                "x": parse_length(rect.get("x")),
                "y": parse_length(rect.get("y")),
                "w": parse_length(rect.get("width")),
                "h": parse_length(rect.get("height")),
            }
    
    for elem in root.iter():
        if tag_name(elem) != "g":
            continue
        if elem.get("data-role") != "image-placeholder":
            continue
        
        placeholder_id = elem.get("id", "")
        caption = elem.get("data-caption", "")
        
        # Find dimensions from rect or image child
        rect_x, rect_y, rect_w, rect_h = 0.0, 0.0, 0.0, 0.0
        image_elem = None
        for child in elem:
            child_tag = tag_name(child)
            if child_tag in ("rect", "image"):
                rect_x = parse_length(child.get("x"))
                rect_y = parse_length(child.get("y"))
                rect_w = parse_length(child.get("width"))
                rect_h = parse_length(child.get("height"))
                if child_tag == "image":
                    image_elem = child
                break
        
        # Collect all transforms from element up to root
        transforms: List[str] = []
        current = elem
        while current is not None:
            t = current.get("transform", "")
            if t:
                transforms.append(t)
            current = parent_map.get(current)
        
        # Apply transforms from innermost (element) to outermost (root)
        x, y, w, h = apply_transform_chain(rect_x, rect_y, rect_w, rect_h, transforms)

        # Resolve clip-paths on ancestors (userSpaceOnUse only)
        clip_chain: List[Dict[str, Optional[str]]] = []
        clip_rects: List[Tuple[float, float, float, float]] = []
        current = parent_map.get(elem)
        while current is not None:
            clip_ref = current.get("clip-path", "")
            transform = current.get("transform", "")
            opacity = current.get("opacity", "")
            if clip_ref or transform or opacity:
                clip_chain.append(
                    {
                        "clip_path": clip_ref or None,
                        "transform": transform or None,
                        "opacity": opacity or None,
                    }
                )
            if clip_ref.startswith("url(#") and clip_ref.endswith(")"):
                clip_id = clip_ref[5:-1]
                units = clip_units.get(clip_id, "userSpaceOnUse")
                if clip_is_rect.get(clip_id) and units == "userSpaceOnUse":
                    rect = clip_rects_map.get(clip_id)
                    if not rect:
                        current = parent_map.get(current)
                        continue
                    # Apply transforms from this ancestor up to root
                    t_chain: List[str] = []
                    cur2 = current
                    while cur2 is not None:
                        t_val = cur2.get("transform", "")
                        if t_val:
                            t_chain.append(t_val)
                        cur2 = parent_map.get(cur2)
                    cx, cy, cw, ch = apply_transform_chain(
                        rect["x"], rect["y"], rect["w"], rect["h"], t_chain
                    )
                    clip_rects.append((cx, cy, cw, ch))
            current = parent_map.get(current)

        clip_chain = list(reversed(clip_chain))
        clip_box = None
        if clip_rects:
            left = max(r[0] for r in clip_rects)
            top = max(r[1] for r in clip_rects)
            right = min(r[0] + r[2] for r in clip_rects)
            bottom = min(r[1] + r[3] for r in clip_rects)
            if right > left and bottom > top:
                clip_box = {"x": left, "y": top, "w": right - left, "h": bottom - top}

        clip_ids = []
        needs_raster = False
        for wrapper in clip_chain:
            clip_path = wrapper.get("clip_path") or ""
            if clip_path.startswith("url(#") and clip_path.endswith(")"):
                clip_id = clip_path[5:-1]
                clip_ids.append(clip_id)
                if not clip_is_rect.get(clip_id, False):
                    needs_raster = True
        clip_defs = {clip_id: clip_xml_map[clip_id] for clip_id in clip_ids if clip_id in clip_xml_map}

        placeholders.append({
            "id": placeholder_id,
            "caption": caption,
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "clip": clip_box,
            "clip_chain": clip_chain,
            "clip_defs": clip_defs,
            "clip_non_rect": needs_raster,
            "elem_xml": ET.tostring(elem, encoding="utf-8").decode("utf-8"),
            "image_elem": image_elem,  # Store reference to the image element
        })
    
    return placeholders


# ============================================================================
# Config and Placeholders Loading
# ============================================================================

def load_config(config_path: Optional[Path]) -> Dict[str, Any]:
    """Load config.json for API settings."""
    if not config_path or not config_path.exists():
        return {}
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)


def load_placeholders(json_path: Optional[Path]) -> Dict[Tuple[str, str], Dict[str, Any]]:
    """Load image_placeholders.json and build {(svg_file, placeholder_id): entry} mapping."""
    if not json_path or not json_path.exists():
        return {}
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    mapping: Dict[Tuple[str, str], Dict[str, Any]] = {}
    for entry in data:
        key = (entry.get("svg_file", ""), entry.get("placeholder_id", ""))
        mapping[key] = entry
    return mapping


# ============================================================================
# Chart Generation (from LLM)
# ============================================================================

CHART_CODE_PROMPT = """You are a Python code generator. Generate python-pptx code to create a native PowerPoint chart based on this description:

Chart Description: {caption}

Requirements:
1. Generate ONLY a Python function named `add_chart_to_slide(slide, left, top, width, height)`
2. The function should use python-pptx to add a chart to the slide
3. Use the provided left, top, width, height parameters (already in Inches)
4. Import statements are NOT needed - CategoryChartData, XyChartData, BubbleChartData, XL_CHART_TYPE, RGBColor, Pt are available
5. Extract data values and labels from the description
6. Return ONLY the function code, no explanations, no markdown code blocks
7. Make sure the chart matches the description (bar chart, line chart, pie chart, etc.)

IMPORTANT RESTRICTIONS - DO NOT use these unsupported attributes:
- chart.fill (Chart has no fill attribute)
- chart.plot_area.format.fill (not supported)
- chart.chart_area (not supported)
- Any background color settings on the chart itself
- Do not reference undefined classes or missing imports

Only use these SAFE attributes:
- chart.series[i].format.fill.solid() and .fore_color.rgb for bar/column colors
- chart.value_axis / chart.category_axis for axis settings
- chart.has_legend for legend toggle
- axis.tick_labels.font for font settings

Example output format:
def add_chart_to_slide(slide, left, top, width, height):
    chart_data = CategoryChartData()
    chart_data.categories = ['A', 'B', 'C']
    chart_data.add_series('Series 1', (1, 2, 3))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    ).chart
    # Set bar color
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0, 128, 128)
"""


def generate_chart_code(
    caption: str,
    api_key: Optional[str] = None,
    base_url: Optional[str] = None,
    model: Optional[str] = None,
    max_tokens: int = 800,
    temperature: float = 0.2,
    error_hint: Optional[str] = None,
) -> Optional[str]:
    """Call code model to generate python-pptx chart code from caption."""
    if OpenAI is None:
        print("Warning: openai package not installed. Skipping chart generation.")
        return None
    
    key = api_key or os.environ.get("OPENAI_API_KEY")
    if not key:
        print("Warning: API key not set. Skipping chart generation.")
        return None
    
    # Use provided model or default to gpt-4o
    use_model = model or "gpt-4o"
    
    try:
        # Create client with optional base_url for compatible APIs (e.g., SiliconFlow)
        client_kwargs = {"api_key": key}
        if base_url:
            client_kwargs["base_url"] = base_url + "/v1"
        client = OpenAI(**client_kwargs)
        user_prompt = CHART_CODE_PROMPT.format(caption=caption)
        if error_hint:
            user_prompt = (
                user_prompt
                + "\n\nPrevious attempt failed with error:\n"
                + str(error_hint)
                + "\nFix the code and regenerate."
            )
        response = client.chat.completions.create(
            model=use_model,
            messages=[
                {"role": "system", "content": "You are a Python code generator specializing in python-pptx charts."},
                {"role": "user", "content": user_prompt},
            ],
            temperature=temperature,
            max_tokens=max_tokens,
        )
        code = response.choices[0].message.content.strip()
        # Remove markdown code blocks if present
        if code.startswith("```"):
            lines = code.split("\n")
            code = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
        return code
    except Exception as e:
        print(f"Warning: Failed to generate chart code: {e}")
        return None


def execute_chart_code(
    code: str,
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
) -> Tuple[bool, Optional[str]]:
    """Safely execute generated chart code."""
    # Prepare execution environment with necessary imports
    exec_globals = {
        "CategoryChartData": CategoryChartData,
        "XL_CHART_TYPE": XL_CHART_TYPE,
        "RGBColor": RGBColor,
        "Pt": Pt,
        "Inches": Inches,
    }
    if XyChartData is not None:
        exec_globals["XyChartData"] = XyChartData
    if BubbleChartData is not None:
        exec_globals["BubbleChartData"] = BubbleChartData
    exec_locals: Dict[str, Any] = {}
    
    try:
        # Execute the function definition
        exec(code, exec_globals, exec_locals)
        
        # Call the generated function
        if "add_chart_to_slide" in exec_locals:
            exec_locals["add_chart_to_slide"](slide, left, top, width, height)
            return True, None
        else:
            return False, "Generated code does not contain add_chart_to_slide function."
    except Exception as e:
        return False, str(e)


# ============================================================================
# Transform Matrix Handling
# ============================================================================

class TransformMatrix:
    """2D affine transform matrix [a, b, c, d, e, f]
    
    | a  c  e |
    | b  d  f |
    | 0  0  1 |
    
    Transform: x' = a*x + c*y + e
               y' = b*x + d*y + f
    """
    
    def __init__(self, a: float = 1.0, b: float = 0.0, c: float = 0.0,
                 d: float = 1.0, e: float = 0.0, f: float = 0.0):
        self.a = a
        self.b = b
        self.c = c
        self.d = d
        self.e = e
        self.f = f
    
    @classmethod
    def identity(cls) -> "TransformMatrix":
        return cls()
    
    @classmethod
    def translate(cls, tx: float, ty: float = 0.0) -> "TransformMatrix":
        return cls(1.0, 0.0, 0.0, 1.0, tx, ty)
    
    @classmethod
    def scale(cls, sx: float, sy: Optional[float] = None) -> "TransformMatrix":
        if sy is None:
            sy = sx
        return cls(sx, 0.0, 0.0, sy, 0.0, 0.0)
    
    @classmethod
    def rotate(cls, angle_deg: float, cx: float = 0.0, cy: float = 0.0) -> "TransformMatrix":
        rad = math.radians(angle_deg)
        cos_a = math.cos(rad)
        sin_a = math.sin(rad)
        # Rotation around (cx, cy)
        # translate(-cx, -cy) -> rotate -> translate(cx, cy)
        return cls(
            cos_a, sin_a, -sin_a, cos_a,
            cx - cos_a * cx + sin_a * cy,
            cy - sin_a * cx - cos_a * cy
        )
    
    def multiply(self, other: "TransformMatrix") -> "TransformMatrix":
        """Return self * other (apply other first, then self)."""
        return TransformMatrix(
            a=self.a * other.a + self.c * other.b,
            b=self.b * other.a + self.d * other.b,
            c=self.a * other.c + self.c * other.d,
            d=self.b * other.c + self.d * other.d,
            e=self.a * other.e + self.c * other.f + self.e,
            f=self.b * other.e + self.d * other.f + self.f,
        )
    
    def transform_point(self, x: float, y: float) -> Tuple[float, float]:
        """Apply transform to a point."""
        return (
            self.a * x + self.c * y + self.e,
            self.b * x + self.d * y + self.f
        )
    
    def transform_vector(self, dx: float, dy: float) -> Tuple[float, float]:
        """Apply transform to a vector (without translation)."""
        return (
            self.a * dx + self.c * dy,
            self.b * dx + self.d * dy
        )
    
    def get_scale(self) -> Tuple[float, float]:
        """Extract approximate scale factors."""
        sx = math.sqrt(self.a * self.a + self.b * self.b)
        sy = math.sqrt(self.c * self.c + self.d * self.d)
        return (sx, sy)
    
    def get_rotation_degrees(self) -> float:
        """Extract rotation angle in degrees."""
        return math.degrees(math.atan2(self.b, self.a))
    
    def get_translation(self) -> Tuple[float, float]:
        """Get translation components."""
        return (self.e, self.f)


def parse_transform(transform_str: Optional[str]) -> TransformMatrix:
    """Parse SVG transform attribute into a TransformMatrix."""
    if not transform_str:
        return TransformMatrix.identity()
    
    result = TransformMatrix.identity()
    
    # Parse matrix()
    for match in MATRIX_RE.finditer(transform_str):
        parts = [p.strip() for p in re.split(r"[\s,]+", match.group(1).strip()) if p]
        if len(parts) == 6:
            try:
                a, b, c, d, e, f = [float(p) for p in parts]
                result = result.multiply(TransformMatrix(a, b, c, d, e, f))
            except ValueError:
                pass
    
    # Parse translate()
    for match in TRANSLATE_RE.finditer(transform_str):
        if "matrix" in transform_str[:match.start()]:
            continue  # Skip if already in matrix
        parts = [p.strip() for p in re.split(r"[\s,]+", match.group(1).strip()) if p]
        if parts:
            try:
                tx = float(parts[0])
                ty = float(parts[1]) if len(parts) > 1 else 0.0
                result = result.multiply(TransformMatrix.translate(tx, ty))
            except ValueError:
                pass
    
    # Parse scale()
    for match in SCALE_RE.finditer(transform_str):
        parts = [p.strip() for p in re.split(r"[\s,]+", match.group(1).strip()) if p]
        if parts:
            try:
                sx = float(parts[0])
                sy = float(parts[1]) if len(parts) > 1 else sx
                result = result.multiply(TransformMatrix.scale(sx, sy))
            except ValueError:
                pass
    
    # Parse rotate()
    for match in ROTATE_RE.finditer(transform_str):
        parts = [p.strip() for p in re.split(r"[\s,]+", match.group(1).strip()) if p]
        if parts:
            try:
                angle = float(parts[0])
                cx = float(parts[1]) if len(parts) > 1 else 0.0
                cy = float(parts[2]) if len(parts) > 2 else 0.0
                result = result.multiply(TransformMatrix.rotate(angle, cx, cy))
            except ValueError:
                pass
    
    return result


# ============================================================================
# SVG Style Inheritance
# ============================================================================

class StyleContext:
    """Track inherited styles through SVG hierarchy."""
    
    def __init__(self):
        self.fill: Optional[str] = None
        self.stroke: Optional[str] = None
        self.stroke_width: Optional[str] = None
        self.opacity: float = 1.0
        self.fill_opacity: float = 1.0
        self.stroke_opacity: float = 1.0
        self.font_family: Optional[str] = None
        self.font_size: Optional[str] = None
        self.font_weight: Optional[str] = None
        self.text_anchor: Optional[str] = None
    
    def copy(self) -> "StyleContext":
        ctx = StyleContext()
        ctx.fill = self.fill
        ctx.stroke = self.stroke
        ctx.stroke_width = self.stroke_width
        ctx.opacity = self.opacity
        ctx.fill_opacity = self.fill_opacity
        ctx.stroke_opacity = self.stroke_opacity
        ctx.font_family = self.font_family
        ctx.font_size = self.font_size
        ctx.font_weight = self.font_weight
        ctx.text_anchor = self.text_anchor
        return ctx
    
    def update_from_element(self, elem: ET.Element) -> "StyleContext":
        """Create a new context inheriting from this one with element's styles."""
        ctx = self.copy()
        
        # Parse style attribute
        style_attr = elem.get("style", "")
        style_dict = {}
        for part in style_attr.split(";"):
            if ":" in part:
                key, val = part.split(":", 1)
                style_dict[key.strip().lower()] = val.strip()
        
        # Direct attributes override style attribute
        def get_style(name: str, attr_name: Optional[str] = None) -> Optional[str]:
            attr_name = attr_name or name
            return elem.get(attr_name) or style_dict.get(name)
        
        if get_style("fill"):
            ctx.fill = get_style("fill")
        if get_style("stroke"):
            ctx.stroke = get_style("stroke")
        if get_style("stroke-width"):
            ctx.stroke_width = get_style("stroke-width")
        if get_style("opacity"):
            ctx.opacity = parse_opacity(get_style("opacity")) * ctx.opacity
        if get_style("fill-opacity"):
            ctx.fill_opacity = parse_opacity(get_style("fill-opacity"))
        if get_style("stroke-opacity"):
            ctx.stroke_opacity = parse_opacity(get_style("stroke-opacity"))
        if get_style("font-family"):
            ctx.font_family = get_style("font-family")
        if get_style("font-size"):
            ctx.font_size = get_style("font-size")
        if get_style("font-weight"):
            ctx.font_weight = get_style("font-weight")
        if get_style("text-anchor"):
            ctx.text_anchor = get_style("text-anchor")
        
        return ctx


# ============================================================================
# Coordinate Conversion
# ============================================================================

class CoordinateConverter:
    """Convert SVG coordinates to PPTX EMU."""
    
    def __init__(self, svg_width: float, svg_height: float,
                 slide_width_emu: int, slide_height_emu: int, dpi: float = 96.0):
        self.svg_width = svg_width
        self.svg_height = svg_height
        self.slide_width_emu = slide_width_emu
        self.slide_height_emu = slide_height_emu
        self.dpi = dpi
        
        # Calculate scale factors
        self.scale_x = slide_width_emu / svg_width if svg_width else 1.0
        self.scale_y = slide_height_emu / svg_height if svg_height else 1.0
    
    def to_emu_x(self, svg_x: float) -> int:
        """Convert SVG x coordinate to EMU."""
        return int(svg_x * self.scale_x)
    
    def to_emu_y(self, svg_y: float) -> int:
        """Convert SVG y coordinate to EMU."""
        return int(svg_y * self.scale_y)
    
    def to_emu_width(self, svg_width: float) -> int:
        """Convert SVG width to EMU."""
        return int(svg_width * self.scale_x)
    
    def to_emu_height(self, svg_height: float) -> int:
        """Convert SVG height to EMU."""
        return int(svg_height * self.scale_y)
    
    def to_emu_length(self, svg_length: float) -> int:
        """Convert SVG length (using average scale) to EMU."""
        avg_scale = (self.scale_x + self.scale_y) / 2
        return int(svg_length * avg_scale)


# ============================================================================
# Shape Addition Functions
# ============================================================================

def apply_fill_to_shape(shape, fill_color: Optional[str], opacity: float = 1.0) -> None:
    """Apply fill color to a PPTX shape."""
    color = parse_color(fill_color)
    if color is None or fill_color in ("none", "transparent"):
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        # Note: Opacity requires OXML manipulation for full support


def apply_stroke_to_shape(shape, stroke_color: Optional[str], stroke_width: Optional[str],
                          converter: CoordinateConverter, opacity: float = 1.0) -> None:
    """Apply stroke to a PPTX shape."""
    color = parse_color(stroke_color)
    if color is None or stroke_color in ("none", "transparent"):
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        if stroke_width:
            width_px = parse_length(stroke_width, 1.0)
            # Convert to EMU, then to Pt for line width
            shape.line.width = Pt(width_px * 0.75)  # Approximate px to pt


def add_svg_rect(slide, elem: ET.Element, transform: TransformMatrix,
                 style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG rect element as PPTX rectangle shape."""
    x = parse_length(elem.get("x"), 0.0)
    y = parse_length(elem.get("y"), 0.0)
    width = parse_length(elem.get("width"), 0.0)
    height = parse_length(elem.get("height"), 0.0)
    rx = parse_length(elem.get("rx"), 0.0)
    ry = parse_length(elem.get("ry"), 0.0)
    
    if width <= 0 or height <= 0:
        return None
    
    # Apply transform to corners
    x1, y1 = transform.transform_point(x, y)
    x2, y2 = transform.transform_point(x + width, y + height)
    
    # Get final bounding box
    left = min(x1, x2)
    top = min(y1, y2)
    final_width = abs(x2 - x1)
    final_height = abs(y2 - y1)
    
    # Convert to EMU
    left_emu = converter.to_emu_x(left)
    top_emu = converter.to_emu_y(top)
    width_emu = converter.to_emu_width(final_width)
    height_emu = converter.to_emu_height(final_height)
    
    # Choose shape type based on rounded corners
    if rx > 0 or ry > 0:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE
    
    shape = slide.shapes.add_shape(
        shape_type, left_emu, top_emu, width_emu, height_emu
    )
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity)
    
    # Apply rotation if present
    rotation = transform.get_rotation_degrees()
    if abs(rotation) > 0.1:
        shape.rotation = rotation
    
    return shape


def add_svg_circle(slide, elem: ET.Element, transform: TransformMatrix,
                   style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG circle element as PPTX oval shape."""
    cx = parse_length(elem.get("cx"), 0.0)
    cy = parse_length(elem.get("cy"), 0.0)
    r = parse_length(elem.get("r"), 0.0)
    
    if r <= 0:
        return None
    
    # Transform center and radius
    center_x, center_y = transform.transform_point(cx, cy)
    sx, sy = transform.get_scale()
    final_rx = r * sx
    final_ry = r * sy
    
    # Calculate bounding box
    left = center_x - final_rx
    top = center_y - final_ry
    width = final_rx * 2
    height = final_ry * 2
    
    # Convert to EMU
    left_emu = converter.to_emu_x(left)
    top_emu = converter.to_emu_y(top)
    width_emu = converter.to_emu_width(width)
    height_emu = converter.to_emu_height(height)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left_emu, top_emu, width_emu, height_emu
    )
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity)
    
    return shape


def add_svg_ellipse(slide, elem: ET.Element, transform: TransformMatrix,
                    style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG ellipse element as PPTX oval shape."""
    cx = parse_length(elem.get("cx"), 0.0)
    cy = parse_length(elem.get("cy"), 0.0)
    rx = parse_length(elem.get("rx"), 0.0)
    ry = parse_length(elem.get("ry"), 0.0)
    
    if rx <= 0 or ry <= 0:
        return None
    
    # Transform center and radii
    center_x, center_y = transform.transform_point(cx, cy)
    sx, sy = transform.get_scale()
    final_rx = rx * sx
    final_ry = ry * sy
    
    # Calculate bounding box
    left = center_x - final_rx
    top = center_y - final_ry
    width = final_rx * 2
    height = final_ry * 2
    
    # Convert to EMU
    left_emu = converter.to_emu_x(left)
    top_emu = converter.to_emu_y(top)
    width_emu = converter.to_emu_width(width)
    height_emu = converter.to_emu_height(height)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left_emu, top_emu, width_emu, height_emu
    )
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity)
    
    # Apply rotation if present
    rotation = transform.get_rotation_degrees()
    if abs(rotation) > 0.1:
        shape.rotation = rotation
    
    return shape


def add_svg_line(slide, elem: ET.Element, transform: TransformMatrix,
                 style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG line element as PPTX connector."""
    x1 = parse_length(elem.get("x1"), 0.0)
    y1 = parse_length(elem.get("y1"), 0.0)
    x2 = parse_length(elem.get("x2"), 0.0)
    y2 = parse_length(elem.get("y2"), 0.0)
    
    # Transform endpoints
    tx1, ty1 = transform.transform_point(x1, y1)
    tx2, ty2 = transform.transform_point(x2, y2)
    
    # Convert to EMU
    x1_emu = converter.to_emu_x(tx1)
    y1_emu = converter.to_emu_y(ty1)
    x2_emu = converter.to_emu_x(tx2)
    y2_emu = converter.to_emu_y(ty2)
    
    # Add as connector shape
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        x1_emu, y1_emu, x2_emu, y2_emu
    )
    
    # Apply stroke style
    stroke = style.stroke or elem.get("stroke") or "#000000"
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    color = parse_color(stroke)
    if color:
        connector.line.color.rgb = color
    if stroke_width:
        width_px = parse_length(stroke_width, 1.0)
        connector.line.width = Pt(width_px * 0.75)
    
    return connector


def add_svg_polygon(slide, elem: ET.Element, transform: TransformMatrix,
                    style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG polygon element as PPTX freeform shape."""
    points_str = elem.get("points", "")
    if not points_str:
        return None
    
    # Parse points
    points = []
    coords = re.findall(r"[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?", points_str)
    for i in range(0, len(coords) - 1, 2):
        try:
            x = float(coords[i])
            y = float(coords[i + 1])
            points.append((x, y))
        except ValueError:
            continue
    
    if len(points) < 3:
        return None
    
    # Transform all points
    transformed = [transform.transform_point(x, y) for x, y in points]
    
    # Find bounding box
    xs = [p[0] for p in transformed]
    ys = [p[1] for p in transformed]
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    width = max_x - min_x
    height = max_y - min_y
    
    if width <= 0 or height <= 0:
        return None
    
    # Convert bounding box to EMU
    left_emu = converter.to_emu_x(min_x)
    top_emu = converter.to_emu_y(min_y)
    
    # Calculate scale: local units to EMU
    # We use 1000 local units for the shape
    local_units = 1000
    scale_x = converter.to_emu_width(width) / local_units
    scale_y = converter.to_emu_height(height) / local_units
    
    # Create freeform shape using FreeformBuilder
    # Normalize points to local coordinate system (0 to local_units)
    start_local_x = int((transformed[0][0] - min_x) / width * local_units) if width else 0
    start_local_y = int((transformed[0][1] - min_y) / height * local_units) if height else 0
    
    builder = slide.shapes.build_freeform(start_local_x, start_local_y, scale=(scale_x, scale_y))
    
    # Draw lines to other points
    line_segments = []
    for px, py in transformed[1:]:
        local_x = int((px - min_x) / width * local_units) if width else 0
        local_y = int((py - min_y) / height * local_units) if height else 0
        line_segments.append((local_x, local_y))
    
    if line_segments:
        builder.add_line_segments(line_segments, close=True)
    
    shape = builder.convert_to_shape(left_emu, top_emu)
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity)
    
    return shape


def add_svg_polyline(slide, elem: ET.Element, transform: TransformMatrix,
                     style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG polyline element as PPTX freeform shape."""
    points_str = elem.get("points", "")
    if not points_str:
        return None
    
    # Parse points
    points = []
    coords = re.findall(r"[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?", points_str)
    for i in range(0, len(coords) - 1, 2):
        try:
            x = float(coords[i])
            y = float(coords[i + 1])
            points.append((x, y))
        except ValueError:
            continue
    
    if len(points) < 2:
        return None
    
    # Transform all points
    transformed = [transform.transform_point(x, y) for x, y in points]
    
    # Find bounding box
    xs = [p[0] for p in transformed]
    ys = [p[1] for p in transformed]
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    width = max_x - min_x or 1
    height = max_y - min_y or 1
    
    # Convert bounding box to EMU
    left_emu = converter.to_emu_x(min_x)
    top_emu = converter.to_emu_y(min_y)
    
    # Calculate scale: local units to EMU
    local_units = 1000
    scale_x = converter.to_emu_width(width) / local_units
    scale_y = converter.to_emu_height(height) / local_units
    
    # Normalize points to local coordinate system
    start_local_x = int((transformed[0][0] - min_x) / width * local_units)
    start_local_y = int((transformed[0][1] - min_y) / height * local_units)
    
    builder = slide.shapes.build_freeform(start_local_x, start_local_y, scale=(scale_x, scale_y))
    
    # Draw lines to other points (not closed)
    line_segments = []
    for px, py in transformed[1:]:
        local_x = int((px - min_x) / width * local_units)
        local_y = int((py - min_y) / height * local_units)
        line_segments.append((local_x, local_y))
    
    if line_segments:
        builder.add_line_segments(line_segments, close=False)
    
    shape = builder.convert_to_shape(left_emu, top_emu)
    
    # Polyline typically has no fill
    fill = style.fill or elem.get("fill") or "none"
    stroke = style.stroke or elem.get("stroke") or "#000000"
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity)
    
    return shape


def add_svg_path(slide, elem: ET.Element, transform: TransformMatrix,
                 style: StyleContext, converter: CoordinateConverter) -> Optional[Any]:
    """Add SVG path element as PPTX freeform shape.
    
    This is the most complex conversion - SVG paths can contain:
    - M/m: moveto
    - L/l: lineto
    - H/h: horizontal lineto
    - V/v: vertical lineto
    - C/c: cubic bezier
    - S/s: smooth cubic bezier
    - Q/q: quadratic bezier
    - T/t: smooth quadratic bezier
    - A/a: arc
    - Z/z: close path
    """
    d = elem.get("d", "")
    if not d:
        return None
    
    # If svgpathtools is available, use it for accurate parsing
    if HAS_SVGPATHTOOLS and parse_path:
        return _add_svg_path_with_svgpathtools(
            slide, elem, d, transform, style, converter
        )
    else:
        # Fallback: simple path parsing
        return _add_svg_path_simple(
            slide, elem, d, transform, style, converter
        )


def _add_svg_path_with_svgpathtools(slide, elem: ET.Element, d: str,
                                     transform: TransformMatrix,
                                     style: StyleContext,
                                     converter: CoordinateConverter) -> Optional[Any]:
    """Parse path with svgpathtools and convert to PPTX freeform."""
    try:
        path = parse_path(d)
    except Exception:
        return _add_svg_path_simple(slide, elem, d, transform, style, converter)
    
    if len(path) == 0:
        return None
    
    # Sample points from the path
    num_samples = 20  # Points per path segment (reduced for performance)
    points = []
    
    for segment in path:
        # Sample each segment
        for t in [i / num_samples for i in range(num_samples + 1)]:
            try:
                pt = segment.point(t)
                x, y = pt.real, pt.imag
                tx, ty = transform.transform_point(x, y)
                points.append((tx, ty))
            except Exception:
                continue
    
    if len(points) < 2:
        return None
    
    # Remove duplicate consecutive points
    unique_points = [points[0]]
    for pt in points[1:]:
        if abs(pt[0] - unique_points[-1][0]) > 0.1 or abs(pt[1] - unique_points[-1][1]) > 0.1:
            unique_points.append(pt)
    
    if len(unique_points) < 2:
        return None
    
    # Find bounding box
    xs = [p[0] for p in unique_points]
    ys = [p[1] for p in unique_points]
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    width = max_x - min_x or 1
    height = max_y - min_y or 1
    
    # Convert to EMU
    left_emu = converter.to_emu_x(min_x)
    top_emu = converter.to_emu_y(min_y)
    
    # Calculate scale: local units to EMU
    local_units = 1000
    scale_x = converter.to_emu_width(width) / local_units if width > 0 else 1
    scale_y = converter.to_emu_height(height) / local_units if height > 0 else 1
    
    # Normalize points to local coordinate system
    start_local_x = int((unique_points[0][0] - min_x) / width * local_units) if width else 0
    start_local_y = int((unique_points[0][1] - min_y) / height * local_units) if height else 0
    
    builder = slide.shapes.build_freeform(start_local_x, start_local_y, scale=(scale_x, scale_y))
    
    # Build line segments
    line_segments = []
    for px, py in unique_points[1:]:
        local_x = int((px - min_x) / width * local_units) if width else 0
        local_y = int((py - min_y) / height * local_units) if height else 0
        line_segments.append((local_x, local_y))
    
    # Check if path is closed
    is_closed = d.strip().upper().endswith("Z")
    
    if line_segments:
        builder.add_line_segments(line_segments, close=is_closed)
    
    shape = builder.convert_to_shape(left_emu, top_emu)
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity)
    
    return shape


def _add_svg_path_simple(slide, elem: ET.Element, d: str,
                         transform: TransformMatrix,
                         style: StyleContext,
                         converter: CoordinateConverter) -> Optional[Any]:
    """Simple path parsing fallback - handles basic M, L, H, V, Z commands."""
    points = []
    current_x, current_y = 0.0, 0.0
    start_x, start_y = 0.0, 0.0
    is_closed = False
    
    # Tokenize the path
    tokens = re.findall(r"[MLHVCSQTAZmlhvcsqtaz]|[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?", d)
    
    i = 0
    cmd = "M"
    
    while i < len(tokens):
        token = tokens[i]
        
        if token.isalpha():
            cmd = token
            i += 1
            if cmd in ("Z", "z"):
                is_closed = True
                current_x, current_y = start_x, start_y
            continue
        
        try:
            if cmd in ("M", "m"):
                x = float(tokens[i])
                y = float(tokens[i + 1])
                i += 2
                if cmd == "m":
                    x += current_x
                    y += current_y
                current_x, current_y = x, y
                start_x, start_y = x, y
                points.append((x, y))
                cmd = "L" if cmd == "M" else "l"
            
            elif cmd in ("L", "l"):
                x = float(tokens[i])
                y = float(tokens[i + 1])
                i += 2
                if cmd == "l":
                    x += current_x
                    y += current_y
                current_x, current_y = x, y
                points.append((x, y))
            
            elif cmd in ("H", "h"):
                x = float(tokens[i])
                i += 1
                if cmd == "h":
                    x += current_x
                current_x = x
                points.append((x, current_y))
            
            elif cmd in ("V", "v"):
                y = float(tokens[i])
                i += 1
                if cmd == "v":
                    y += current_y
                current_y = y
                points.append((current_x, y))
            
            elif cmd in ("C", "c"):
                # Cubic bezier - sample points
                x1 = float(tokens[i])
                y1 = float(tokens[i + 1])
                x2 = float(tokens[i + 2])
                y2 = float(tokens[i + 3])
                x = float(tokens[i + 4])
                y = float(tokens[i + 5])
                i += 6
                if cmd == "c":
                    x1 += current_x
                    y1 += current_y
                    x2 += current_x
                    y2 += current_y
                    x += current_x
                    y += current_y
                # Sample bezier curve
                for t in [0.25, 0.5, 0.75, 1.0]:
                    bx = (1-t)**3 * current_x + 3*(1-t)**2*t * x1 + 3*(1-t)*t**2 * x2 + t**3 * x
                    by = (1-t)**3 * current_y + 3*(1-t)**2*t * y1 + 3*(1-t)*t**2 * y2 + t**3 * y
                    points.append((bx, by))
                current_x, current_y = x, y
            
            elif cmd in ("Q", "q"):
                # Quadratic bezier
                x1 = float(tokens[i])
                y1 = float(tokens[i + 1])
                x = float(tokens[i + 2])
                y = float(tokens[i + 3])
                i += 4
                if cmd == "q":
                    x1 += current_x
                    y1 += current_y
                    x += current_x
                    y += current_y
                # Sample bezier
                for t in [0.25, 0.5, 0.75, 1.0]:
                    bx = (1-t)**2 * current_x + 2*(1-t)*t * x1 + t**2 * x
                    by = (1-t)**2 * current_y + 2*(1-t)*t * y1 + t**2 * y
                    points.append((bx, by))
                current_x, current_y = x, y
            
            else:
                # Skip unsupported commands
                i += 1
        
        except (ValueError, IndexError):
            i += 1
            continue
    
    if len(points) < 2:
        return None
    
    # Transform points
    transformed = [transform.transform_point(x, y) for x, y in points]
    
    # Find bounding box
    xs = [p[0] for p in transformed]
    ys = [p[1] for p in transformed]
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    width = max_x - min_x or 1
    height = max_y - min_y or 1
    
    # Convert to EMU
    left_emu = converter.to_emu_x(min_x)
    top_emu = converter.to_emu_y(min_y)
    
    # Calculate scale: local units to EMU
    local_units = 1000
    scale_x = converter.to_emu_width(width) / local_units if width > 0 else 1
    scale_y = converter.to_emu_height(height) / local_units if height > 0 else 1
    
    # Normalize points to local coordinate system
    start_local_x = int((transformed[0][0] - min_x) / width * local_units)
    start_local_y = int((transformed[0][1] - min_y) / height * local_units)
    
    builder = slide.shapes.build_freeform(start_local_x, start_local_y, scale=(scale_x, scale_y))
    
    # Build line segments
    line_segments = []
    for px, py in transformed[1:]:
        local_x = int((px - min_x) / width * local_units)
        local_y = int((py - min_y) / height * local_units)
        line_segments.append((local_x, local_y))
    
    if line_segments:
        builder.add_line_segments(line_segments, close=is_closed)
    
    shape = builder.convert_to_shape(left_emu, top_emu)
    
    # Apply styles
    fill = style.fill or elem.get("fill")
    stroke = style.stroke or elem.get("stroke")
    stroke_width = style.stroke_width or elem.get("stroke-width")
    
    apply_fill_to_shape(shape, fill, style.fill_opacity * style.opacity)
    apply_stroke_to_shape(shape, stroke, stroke_width, converter, style.stroke_opacity * style.opacity)
    
    return shape


def add_svg_image(slide, elem: ET.Element, transform: TransformMatrix,
                  style: StyleContext, converter: CoordinateConverter,
                  svg_path: Optional[Path] = None) -> Optional[Any]:
    """Add SVG image element as PPTX picture."""
    x = parse_length(elem.get("x"), 0.0)
    y = parse_length(elem.get("y"), 0.0)
    width = parse_length(elem.get("width"), 0.0)
    height = parse_length(elem.get("height"), 0.0)
    
    # Get href (try both with and without namespace)
    href = elem.get(f"{{{XLINK_NS}}}href") or elem.get("href", "")
    
    if not href or width <= 0 or height <= 0:
        return None
    
    # Transform position
    tx, ty = transform.transform_point(x, y)
    sx, sy = transform.get_scale()
    final_width = width * sx
    final_height = height * sy
    
    # Convert to EMU
    left_emu = converter.to_emu_x(tx)
    top_emu = converter.to_emu_y(ty)
    width_emu = converter.to_emu_width(final_width)
    height_emu = converter.to_emu_height(final_height)
    
    # Handle data URI
    if href.startswith("data:"):
        # Parse data URI: data:[<mediatype>][;base64],<data>
        match = re.match(r"data:([^;,]+)?(?:;base64)?,(.+)", href, re.DOTALL)
        if not match:
            return None
        
        media_type = match.group(1) or "image/png"
        data = match.group(2)
        
        try:
            image_data = base64.b64decode(data)
        except Exception:
            return None
        
        # Create temp file with correct extension
        ext = ".png"
        if "jpeg" in media_type or "jpg" in media_type:
            ext = ".jpg"
        elif "gif" in media_type:
            ext = ".gif"
        elif "webp" in media_type:
            ext = ".webp"
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(image_data)
            tmp_path = tmp.name
        
        try:
            shape = slide.shapes.add_picture(
                tmp_path, left_emu, top_emu, width_emu, height_emu
            )
        finally:
            Path(tmp_path).unlink(missing_ok=True)
        
        return shape
    
    # Handle external file reference
    if svg_path:
        image_path = svg_path.parent / href
        if image_path.exists():
            shape = slide.shapes.add_picture(
                str(image_path), left_emu, top_emu, width_emu, height_emu
            )
            return shape
    
    return None


def add_svg_text(slide, elem: ET.Element, transform: TransformMatrix,
                 style: StyleContext, converter: CoordinateConverter,
                 cjk_font: str = "PingFang SC") -> Optional[Any]:
    """Add SVG text element as PPTX textbox."""
    # Get text content
    text_content = ""
    if elem.text:
        text_content += elem.text
    for child in elem:
        if tag_name(child) == "tspan":
            if child.text:
                text_content += child.text
            if child.tail:
                text_content += child.tail
    
    text_content = text_content.replace("\u00a0", " ").strip()
    if not text_content:
        return None
    
    # Get position
    x = parse_length(elem.get("x"), 0.0)
    y = parse_length(elem.get("y"), 0.0)
    
    # Transform position
    tx, ty = transform.transform_point(x, y)
    
    # Get font size
    font_size_str = style.font_size or elem.get("font-size", "16")
    font_size = parse_length(font_size_str, 16.0)
    sx, sy = transform.get_scale()
    final_font_size = font_size * (sx + sy) / 2
    
    # Estimate text dimensions (rough approximation)
    char_width = final_font_size * 0.6
    text_width = len(text_content) * char_width
    text_height = final_font_size * 1.5
    
    # Adjust position based on text-anchor
    anchor = style.text_anchor or elem.get("text-anchor", "start")
    if anchor == "middle":
        tx -= text_width / 2
    elif anchor == "end":
        tx -= text_width
    
    # Adjust y position (SVG text y is baseline, PPTX is top)
    ty -= final_font_size * 0.8
    
    # Convert to EMU
    left_emu = converter.to_emu_x(tx)
    top_emu = converter.to_emu_y(ty)
    width_emu = converter.to_emu_width(text_width * 1.2)  # Add some padding
    height_emu = converter.to_emu_height(text_height)
    
    # Create textbox
    shape = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text_content
    
    # Set alignment
    if anchor == "middle":
        p.alignment = PP_ALIGN.CENTER
    elif anchor == "end":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    
    # Apply font style
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(final_font_size * 0.75)  # Approximate px to pt
        
        # Set font family
        font_family = style.font_family or elem.get("font-family")
        use_theme_font = font_family_is_theme(font_family)
        if font_family:
            # Clean up font family string
            fonts = [f.strip().strip("\"'") for f in font_family.split(",")]
            for font in fonts:
                if font.lower() not in GENERIC_FONTS and "msfontservice" not in font.lower():
                    run.font.name = font
                    break
            else:
                # Use CJK font if text contains CJK characters
                if CJK_RE.search(text_content) and cjk_font and not use_theme_font:
                    run.font.name = cjk_font
        if CJK_RE.search(text_content) and not use_theme_font:
            ea_font = run.font.name or cjk_font
            if ea_font:
                set_run_ea_font(run, ea_font)
        
        # Set color
        fill = style.fill or elem.get("fill")
        color = parse_color(fill)
        if color:
            run.font.color.rgb = color
        
        # Set bold
        font_weight = style.font_weight or elem.get("font-weight", "")
        if font_weight in ("bold", "700", "800", "900"):
            run.font.bold = True
    
    return shape


# ============================================================================
# Semantic Layer Textbox Extraction and Addition
# ============================================================================

def extract_semantic_textboxes(svg_path: Path) -> Tuple[Dict[str, float], List[Dict[str, object]]]:
    """Extract textboxes from semantic layer in SVG file.
    
    This handles the data-type="textbox" elements in the semantic-layer group.
    Each textbox is handled individually - no merging of body/bullet/numbered roles.
    """
    tree = ET.parse(svg_path)
    root = tree.getroot()
    width = parse_length(root.get("width"))
    height = parse_length(root.get("height"))
    canvas = {"w": width, "h": height}

    def collect_text_items(container: ET.Element) -> List[Dict[str, object]]:
        items = []
        for text in container.iter():
            if tag_name(text) != "text":
                continue
            content = read_text_content(text)
            if not content:
                continue
            tx, ty = parse_transform_xy(text.get("transform"))
            rotation = parse_transform_rotation(text.get("transform"))
            if tx == 0.0 and ty == 0.0:
                tx = parse_length(text.get("x"))
                ty = parse_length(text.get("y"))
            items.append(
                {
                    "text": content,
                    "x": tx,
                    "y": ty,
                    "font_size": parse_length(text.get("font-size")),
                    "font_family": text.get("font-family"),
                    "font_theme": font_family_is_theme(text.get("font-family")),
                    "fill": text.get("fill"),
                    "text_anchor": text.get("text-anchor"),
                    "rotation": rotation,
                }
            )
        return items

    def infer_box_rotation(items: List[Dict[str, object]]) -> float:
        rotations = []
        for item in items:
            rot = item.get("rotation")
            if isinstance(rot, (int, float)):
                rot_val = normalize_rotation(float(rot))
                if abs(rot_val) >= 1.0:
                    rotations.append(rot_val)
        if not rotations:
            return 0.0
        rotations.sort()
        return rotations[len(rotations) // 2]

    # Simply collect all data-type="textbox" elements - no merging
    textboxes = []
    for g in root.iter():
        if tag_name(g) != "g":
            continue
        if g.get("data-type") != "textbox":
            continue
        texts = collect_text_items(g)
        textboxes.append(
            {
                "id": g.get("id") or "",
                "x": parse_length(g.get("data-x")),
                "y": parse_length(g.get("data-y")),
                "w": parse_length(g.get("data-w")),
                "h": parse_length(g.get("data-h")),
                "role": g.get("data-role") or "",
                "texts": texts,
                "rotation": infer_box_rotation(texts),
            }
        )
    return canvas, textboxes


def add_semantic_textbox(
    slide,
    tb: Dict[str, object],
    dpi: float,
    line_tol: float,
    box_pad: float,
    cjk_font: str,
    width_expand: float = 1.15,  # Expand width by 15% to avoid forced line breaks
) -> Optional[Any]:
    """Add a semantic textbox to the slide.
    
    Args:
        width_expand: Factor to expand the textbox width (e.g., 1.15 = 15% wider).
                      This helps avoid forced line breaks due to font rendering differences.
    """
    pad = max(0.0, float(box_pad))
    base_x = float(tb["x"]) - pad
    base_y = float(tb["y"]) - pad
    # Expand width to avoid forced line breaks
    base_w = (float(tb["w"]) + 2 * pad) * width_expand
    base_h = float(tb["h"]) + 2 * pad
    
    if base_w <= 0 or base_h <= 0:
        return None
    
    rotation = normalize_rotation(float(tb.get("rotation", 0.0) or 0.0))
    swap_dims = abs(abs(rotation) - 90.0) <= 2.0
    
    if swap_dims:
        width_px = base_h
        height_px = base_w
        center_x = base_x + base_w / 2.0
        center_y = base_y + base_h / 2.0
        left = Inches((center_x - width_px / 2.0) / dpi)
        top = Inches((center_y - height_px / 2.0) / dpi)
    else:
        width_px = base_w
        height_px = base_h
        left = Inches(base_x / dpi)
        top = Inches(base_y / dpi)
    
    width = Inches(width_px / dpi)
    height = Inches(height_px / dpi)
    
    shape = slide.shapes.add_textbox(left, top, width, height)
    if abs(rotation) >= 1.0:
        shape.rotation = rotation
    if tb.get("id"):
        shape.name = str(tb["id"])
    
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    text_items = tb.get("texts", [])
    lines = group_text_lines(text_items, line_tol=line_tol)
    
    for idx, line in enumerate(lines):
        line_items = sorted(line["items"], key=lambda it: it["x"])
        line_text = assemble_line_text(line_items)
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        if not line_items:
            continue
        run = p.runs[0] if p.runs else p.add_run()
        style = line_items[0]
        font_size = style.get("font_size") or 0.0
        if font_size:
            run.font.size = Pt(font_size * 0.75)
        use_theme_font = bool(style.get("font_theme"))
        font_name = None
        if not use_theme_font:
            font_name = pick_font_name(style.get("font_family"), line_text, cjk_font)
            if font_name:
                run.font.name = font_name
        if has_cjk(line_text) and not use_theme_font:
            ea_font = font_name or cjk_font
            if ea_font:
                set_run_ea_font(run, ea_font)
        color = parse_color(style.get("fill"))
        if color:
            run.font.color.rgb = color
        anchor = style.get("text_anchor")
        if anchor == "middle":
            p.alignment = PP_ALIGN.CENTER
        elif anchor == "end":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
    
    return shape


# ============================================================================
# Image Placeholder Processing (with Correct Position Calculation)
# ============================================================================

def add_image_placeholder(
    slide,
    placeholder: Dict[str, Any],
    converter: CoordinateConverter,
    svg_path: Optional[Path] = None,
    clip_to_canvas: bool = True,
) -> Optional[Any]:
    """Add an image placeholder to the slide using pre-calculated position.
    
    The position is calculated by extract_image_placeholders which correctly
    handles complex nested transforms.
    
    Uses PPTX's built-in crop properties instead of PIL cropping to preserve
    the correct aspect ratio when the image is stretched in the SVG.
    """
    orig_x = placeholder.get("x", 0)
    orig_y = placeholder.get("y", 0)
    orig_w = placeholder.get("w", 0)
    orig_h = placeholder.get("h", 0)
    image_elem = placeholder.get("image_elem")
    
    if orig_w <= 0 or orig_h <= 0:
        return None

    # Rasterize non-rect clip paths to preserve masks
    if placeholder.get("clip_non_rect"):
        raster = rasterize_clipped_placeholder(placeholder, converter)
        if raster:
            tmp_path, bbox = raster
            left_emu = converter.to_emu_x(bbox[0])
            top_emu = converter.to_emu_y(bbox[1])
            width_emu = converter.to_emu_width(bbox[2] - bbox[0])
            height_emu = converter.to_emu_height(bbox[3] - bbox[1])
            try:
                shape = slide.shapes.add_picture(tmp_path, left_emu, top_emu, width_emu, height_emu)
                return shape
            finally:
                Path(tmp_path).unlink(missing_ok=True)
    
    # Get the image data from the image element
    if image_elem is None:
        return None
    
    href = image_elem.get(f"{{{XLINK_NS}}}href") or image_elem.get("href", "")
    if not href:
        return None
    
    # Calculate crop region if clipping to canvas or clip-path
    x, y, w, h = orig_x, orig_y, orig_w, orig_h
    crop_fractions = None  # (left, top, right, bottom) as fractions
    clip_box = placeholder.get("clip")
    
    if clip_to_canvas:
        canvas_w = converter.svg_width
        canvas_h = converter.svg_height
        
        # Calculate visible region (intersection with clip-path + canvas)
        visible_left = orig_x
        visible_top = orig_y
        visible_right = orig_x + orig_w
        visible_bottom = orig_y + orig_h
        if isinstance(clip_box, dict):
            try:
                clip_left = float(clip_box.get("x", visible_left))
                clip_top = float(clip_box.get("y", visible_top))
                clip_right = clip_left + float(clip_box.get("w", 0.0))
                clip_bottom = clip_top + float(clip_box.get("h", 0.0))
                visible_left = max(visible_left, clip_left)
                visible_top = max(visible_top, clip_top)
                visible_right = min(visible_right, clip_right)
                visible_bottom = min(visible_bottom, clip_bottom)
            except Exception:
                pass
        visible_left = max(0, visible_left)
        visible_top = max(0, visible_top)
        visible_right = min(canvas_w, visible_right)
        visible_bottom = min(canvas_h, visible_bottom)
        
        # Check if image is completely outside canvas
        if visible_right <= visible_left or visible_bottom <= visible_top:
            return None
        
        # Calculate crop fractions (relative to the image's displayed bounds)
        crop_left = (visible_left - orig_x) / orig_w if orig_w > 0 else 0
        crop_top = (visible_top - orig_y) / orig_h if orig_h > 0 else 0
        crop_right = (orig_x + orig_w - visible_right) / orig_w if orig_w > 0 else 0
        crop_bottom = (orig_y + orig_h - visible_bottom) / orig_h if orig_h > 0 else 0
        
        # Only set crop if we actually need to crop
        if crop_left > 0.001 or crop_top > 0.001 or crop_right > 0.001 or crop_bottom > 0.001:
            crop_fractions = (crop_left, crop_top, crop_right, crop_bottom)
        
        # Use the visible region for positioning
        x = visible_left
        y = visible_top
        w = visible_right - visible_left
        h = visible_bottom - visible_top
    
    # Convert to EMU using the correct coordinates
    left_emu = converter.to_emu_x(x)
    top_emu = converter.to_emu_y(y)
    width_emu = converter.to_emu_width(w)
    height_emu = converter.to_emu_height(h)
    
    # Get SVG image element dimensions for preserveAspectRatio calculation
    svg_img_w = parse_length(image_elem.get("width"), 0.0)
    svg_img_h = parse_length(image_elem.get("height"), 0.0)
    preserve_aspect = image_elem.get("preserveAspectRatio", "xMidYMid")
    
    # Helper function to add picture and apply cropping via PIL
    def add_picture_with_crop(image_path_or_stream, left, top, width, height, crop=None):
        """Add picture and apply cropping using PIL.
        
        Correctly handles preserveAspectRatio: xMidYMid slice which:
        1. Scales the image to COVER the display area (using larger scale)
        2. Centers the image
        3. Clips overflow
        """
        # Load image
        if isinstance(image_path_or_stream, bytes):
            img = Image.open(io.BytesIO(image_path_or_stream))
        else:
            img = Image.open(image_path_or_stream)
        
        img_w, img_h = img.size
        
        if crop and svg_img_w > 0 and svg_img_h > 0:
            crop_l, crop_t, crop_r, crop_b = crop
            
            # Handle preserveAspectRatio: xMidYMid slice
            # This scales the image to COVER the viewbox, then centers it
            
            # Calculate scale factors
            scale_by_width = svg_img_w / img_w
            scale_by_height = svg_img_h / img_h
            
            if "slice" in preserve_aspect:
                # slice: use larger scale to cover
                par_scale = max(scale_by_width, scale_by_height)
            else:
                # meet: use smaller scale to fit
                par_scale = min(scale_by_width, scale_by_height)
            
            # Scaled image dimensions
            scaled_w = img_w * par_scale
            scaled_h = img_h * par_scale
            
            # Calculate offset for centering (xMidYMid)
            offset_x = (scaled_w - svg_img_w) / 2  # Pixels hidden on left
            offset_y = (scaled_h - svg_img_h) / 2  # Pixels hidden on top
            
            # The crop fractions are relative to the SVG viewbox (svg_img_w x svg_img_h)
            # We need to map them back to the original image pixels
            
            # Visible region in SVG viewbox coordinates
            vis_x_in_view = crop_l * svg_img_w
            vis_y_in_view = crop_t * svg_img_h
            vis_w_in_view = svg_img_w * (1 - crop_l - crop_r)
            vis_h_in_view = svg_img_h * (1 - crop_t - crop_b)
            
            # Map to scaled image coordinates (add offset for centering)
            vis_x_in_scaled = offset_x + vis_x_in_view
            vis_y_in_scaled = offset_y + vis_y_in_view
            vis_w_in_scaled = vis_w_in_view
            vis_h_in_scaled = vis_h_in_view
            
            # Map to original image pixels
            left_px = int(vis_x_in_scaled / par_scale)
            top_px = int(vis_y_in_scaled / par_scale)
            right_px = int((vis_x_in_scaled + vis_w_in_scaled) / par_scale)
            bottom_px = int((vis_y_in_scaled + vis_h_in_scaled) / par_scale)
            
            # Ensure valid crop box
            left_px = max(0, min(left_px, img_w - 1))
            top_px = max(0, min(top_px, img_h - 1))
            right_px = max(left_px + 1, min(right_px, img_w))
            bottom_px = max(top_px + 1, min(bottom_px, img_h))
            
            # Crop the image
            img = img.crop((left_px, top_px, right_px, bottom_px))
        elif crop:
            # Fallback for cases without SVG dimensions
            crop_l, crop_t, crop_r, crop_b = crop
            left_px = int(crop_l * img_w)
            top_px = int(crop_t * img_h)
            right_px = int(img_w * (1 - crop_r))
            bottom_px = int(img_h * (1 - crop_b))
            
            left_px = max(0, min(left_px, img_w - 1))
            top_px = max(0, min(top_px, img_h - 1))
            right_px = max(left_px + 1, min(right_px, img_w))
            bottom_px = max(top_px + 1, min(bottom_px, img_h))
            
            img = img.crop((left_px, top_px, right_px, bottom_px))
        
        # Save to temp file
        ext = ".png"
        img_format = "PNG"
        if hasattr(img, 'format') and img.format:
            if img.format.lower() in ("jpeg", "jpg"):
                ext = ".jpg"
                img_format = "JPEG"
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            img.save(tmp, format=img_format)
            tmp_path = tmp.name
        
        try:
            shape = slide.shapes.add_picture(tmp_path, left, top, width, height)
            return shape
        finally:
            Path(tmp_path).unlink(missing_ok=True)
    
    # Handle data URI
    if href.startswith("data:"):
        match = re.match(r"data:([^;,]+)?(?:;base64)?,(.+)", href, re.DOTALL)
        if not match:
            return None
        
        data = match.group(2)
        
        try:
            image_data = base64.b64decode(data)
        except Exception:
            return None
        
        return add_picture_with_crop(image_data, left_emu, top_emu, width_emu, height_emu, crop_fractions)
    
    # Handle external file reference
    if svg_path:
        image_path = svg_path.parent / href
        if image_path.exists():
            return add_picture_with_crop(image_path, left_emu, top_emu, width_emu, height_emu, crop_fractions)
    
    return None


# ============================================================================
# Main SVG Processing
# ============================================================================

def process_svg_element(slide, elem: ET.Element, parent_transform: TransformMatrix,
                        parent_style: StyleContext, converter: CoordinateConverter,
                        svg_path: Optional[Path] = None, cjk_font: str = "PingFang SC",
                        skip_elements: Optional[set] = None,
                        skip_image_placeholders: bool = True) -> List[Any]:
    """Process a single SVG element and its children, return list of shapes added.
    
    Args:
        skip_image_placeholders: If True, skip elements with data-role="image-placeholder"
                                 (they are processed separately with correct transform handling)
    """
    shapes = []
    tag = tag_name(elem)
    
    # Skip certain elements
    skip_elements = skip_elements or set()
    elem_id = elem.get("id", "")
    if elem_id in skip_elements:
        return shapes
    
    # Skip defs, clipPath, mask, etc.
    if tag in ("defs", "clipPath", "mask", "symbol", "use", "metadata", "title", "desc"):
        return shapes
    
    # Skip semantic layer (handled separately if needed)
    if elem.get("data-type") == "semantic-layer" or elem_id == "semantic-layer":
        return shapes
    
    # Skip image placeholders if requested (they are processed separately)
    if skip_image_placeholders and elem.get("data-role") == "image-placeholder":
        return shapes
    
    # Update transform
    elem_transform_str = elem.get("transform", "")
    elem_transform = parse_transform(elem_transform_str)
    current_transform = parent_transform.multiply(elem_transform)
    
    # Update style
    current_style = parent_style.update_from_element(elem)
    
    # Process based on element type
    shape = None
    
    if tag == "rect":
        shape = add_svg_rect(slide, elem, current_transform, current_style, converter)
    
    elif tag == "circle":
        shape = add_svg_circle(slide, elem, current_transform, current_style, converter)
    
    elif tag == "ellipse":
        shape = add_svg_ellipse(slide, elem, current_transform, current_style, converter)
    
    elif tag == "line":
        shape = add_svg_line(slide, elem, current_transform, current_style, converter)
    
    elif tag == "polygon":
        shape = add_svg_polygon(slide, elem, current_transform, current_style, converter)
    
    elif tag == "polyline":
        shape = add_svg_polyline(slide, elem, current_transform, current_style, converter)
    
    elif tag == "path":
        shape = add_svg_path(slide, elem, current_transform, current_style, converter)
    
    elif tag == "image":
        shape = add_svg_image(slide, elem, current_transform, current_style, converter, svg_path)
    
    elif tag == "text":
        shape = add_svg_text(slide, elem, current_transform, current_style, converter, cjk_font)
    
    elif tag == "g":
        # Process group children recursively
        for child in elem:
            child_shapes = process_svg_element(
                slide, child, current_transform, current_style, converter,
                svg_path, cjk_font, skip_elements, skip_image_placeholders
            )
            shapes.extend(child_shapes)
    
    elif tag == "svg":
        # Nested SVG - process children
        for child in elem:
            child_shapes = process_svg_element(
                slide, child, current_transform, current_style, converter,
                svg_path, cjk_font, skip_elements, skip_image_placeholders
            )
            shapes.extend(child_shapes)
    
    if shape:
        shapes.append(shape)
    
    return shapes


def convert_svg_to_slide(
    prs: Presentation,
    svg_path: Path,
    dpi: float = 96.0,
    cjk_font: str = "PingFang SC",
    skip_elements: Optional[set] = None,
    line_tol: float = 2.0,
    box_pad: float = 12.0,
    placeholders_map: Optional[Dict[Tuple[str, str], Dict[str, Any]]] = None,
    chart_config: Optional[Dict[str, Any]] = None,
    skip_charts: bool = False,
) -> None:
    """Convert a single SVG file to a PPTX slide.
    
    Args:
        placeholders_map: Mapping from (svg_file, placeholder_id) to placeholder info
        chart_config: Configuration for chart generation (API key, model, etc.)
        skip_charts: If True, skip chart generation
    """
    # Parse SVG
    tree = ET.parse(svg_path)
    root = tree.getroot()
    
    # Get SVG dimensions
    svg_width = parse_length(root.get("width"), 1920.0)
    svg_height = parse_length(root.get("height"), 1080.0)
    
    # Handle viewBox
    viewbox = root.get("viewBox", "")
    if viewbox:
        vb_parts = viewbox.split()
        if len(vb_parts) == 4:
            try:
                svg_width = float(vb_parts[2])
                svg_height = float(vb_parts[3])
            except ValueError:
                pass
    
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Create coordinate converter
    converter = CoordinateConverter(
        svg_width, svg_height,
        prs.slide_width, prs.slide_height,
        dpi
    )
    
    # Initialize transform and style
    base_transform = TransformMatrix.identity()
    base_style = StyleContext()
    
    # Build skip set - always skip semantic-layer for regular SVG element processing
    # (textboxes from semantic-layer are handled separately)
    actual_skip = set(skip_elements) if skip_elements else set()
    actual_skip.add("semantic-layer")
    
    # Check if we should skip textboxes entirely
    skip_textboxes = skip_elements and "semantic-layer" in skip_elements
    
    # Extract image placeholders with correct position calculation
    image_placeholders = extract_image_placeholders(svg_path)
    placeholder_ids = {ph["id"] for ph in image_placeholders}
    
    # Process all children of root (except semantic layer and image placeholders)
    for child in root:
        process_svg_element(
            slide, child, base_transform, base_style, converter,
            svg_path, cjk_font, actual_skip, skip_image_placeholders=True
        )
    
    # Add image placeholders with correct positioning
    svg_name = svg_path.name
    placeholders_map = placeholders_map or {}
    chart_config = chart_config or {}
    image_count = 0
    chart_count = 0
    
    for ph in image_placeholders:
        placeholder_id = ph["id"]
        entry = placeholders_map.get((svg_name, placeholder_id))
        
        # Check if this is a chart
        is_chart = entry and entry.get("is_chart", False)
        
        if is_chart and not skip_charts:
            # This is a chart placeholder - must generate chart, not image
            caption = entry.get("caption", "")
            if caption and chart_config:
                # Generate chart code
                code = generate_chart_code(
                    caption,
                    api_key=chart_config.get("api_key"),
                    base_url=chart_config.get("base_url"),
                    model=chart_config.get("chart_model"),
                    max_tokens=chart_config.get("chart_max_tokens", 800),
                    temperature=chart_config.get("chart_temperature", 0.2),
                )
                if code:
                    # Execute chart code
                    left_inches = Inches(ph["x"] / dpi)
                    top_inches = Inches(ph["y"] / dpi)
                    width_inches = Inches(ph["w"] / dpi)
                    height_inches = Inches(ph["h"] / dpi)
                    
                    success, error = execute_chart_code(
                        code, slide, left_inches, top_inches, width_inches, height_inches
                    )
                    if success:
                        chart_count += 1
                    else:
                        print(f"    Warning: Failed to create chart '{placeholder_id}': {error}")
                else:
                    print(f"    Warning: Failed to generate chart code for '{placeholder_id}'")
            else:
                if not caption:
                    print(f"    Warning: Chart '{placeholder_id}' has no caption")
                if not chart_config:
                    print(f"    Warning: No chart config provided for '{placeholder_id}'")
            # Do NOT fall back to image for chart placeholders
            continue
        
        # Add as image (only for non-chart placeholders)
        shape = add_image_placeholder(slide, ph, converter, svg_path)
        if shape:
            image_count += 1
    
    # Now extract and add semantic textboxes (unless skipped)
    textbox_count = 0
    if not skip_textboxes:
        _, textboxes = extract_semantic_textboxes(svg_path)
        for tb in textboxes:
            if tb.get("w", 0) <= 0 or tb.get("h", 0) <= 0:
                continue
            shape = add_semantic_textbox(
                slide, tb, dpi=dpi, line_tol=line_tol, box_pad=box_pad, cjk_font=cjk_font
            )
            if shape:
                textbox_count += 1
    
    total_shapes = len(slide.shapes)
    parts = [f"{total_shapes} shapes"]
    if image_count:
        parts.append(f"{image_count} images")
    if chart_count:
        parts.append(f"{chart_count} charts")
    if textbox_count:
        parts.append(f"{textbox_count} textboxes")
    print(f"  Converted: {svg_path.name} ({', '.join(parts)})")


def build_pptx_pro(
    svg_paths: List[Path],
    out_pptx: Path,
    dpi: float = 96.0,
    cjk_font: str = "PingFang SC",
    skip_elements: Optional[set] = None,
    line_tol: float = 2.0,
    box_pad: float = 12.0,
    placeholders_map: Optional[Dict[Tuple[str, str], Dict[str, Any]]] = None,
    chart_config: Optional[Dict[str, Any]] = None,
    skip_charts: bool = False,
) -> None:
    """Build PPTX from multiple SVG files with editable shapes.
    
    Args:
        placeholders_map: Mapping from (svg_file, placeholder_id) to placeholder info
        chart_config: Configuration for chart generation (API key, model, etc.)
        skip_charts: If True, skip chart generation
    """
    if not svg_paths:
        print("No SVG files to convert.")
        return
    
    # Parse first SVG to get dimensions
    first_tree = ET.parse(svg_paths[0])
    first_root = first_tree.getroot()
    svg_width = parse_length(first_root.get("width"), 1920.0)
    svg_height = parse_length(first_root.get("height"), 1080.0)
    
    # Handle viewBox
    viewbox = first_root.get("viewBox", "")
    if viewbox:
        vb_parts = viewbox.split()
        if len(vb_parts) == 4:
            try:
                svg_width = float(vb_parts[2])
                svg_height = float(vb_parts[3])
            except ValueError:
                pass
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(svg_width / dpi)
    prs.slide_height = Inches(svg_height / dpi)
    
    print(f"Creating PPTX: {svg_width}x{svg_height} px @ {dpi} DPI")
    print(f"Slide size: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    
    # Convert each SVG
    for svg_path in svg_paths:
        convert_svg_to_slide(
            prs, svg_path, dpi, cjk_font, skip_elements, line_tol, box_pad,
            placeholders_map, chart_config, skip_charts
        )
    
    # Save
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    print(f"\nSaved: {out_pptx}")


# ============================================================================
# Command Line Interface
# ============================================================================

def main() -> None:
    parser = argparse.ArgumentParser(
        description="SVG to PPTX Pro: Convert SVG to editable PPTX shapes."
    )
    parser.add_argument(
        "--input", "-i", required=True,
        help="Input SVG file or directory containing SVG files."
    )
    parser.add_argument(
        "--output", "-o", required=True,
        help="Output PPTX file path."
    )
    parser.add_argument(
        "--dpi", type=float, default=96.0,
        help="SVG pixel DPI (default: 96.0)."
    )
    parser.add_argument(
        "--cjk-font", default="PingFang SC",
        help="Font for CJK text (default: PingFang SC)."
    )
    parser.add_argument(
        "--line-tol", type=float, default=2.0,
        help="Line grouping tolerance in px for textboxes (default: 2.0)."
    )
    parser.add_argument(
        "--box-pad", type=float, default=20.0,
        help="Textbox padding in px (default: 20.0)."
    )
    parser.add_argument(
        "--skip-textboxes", action="store_true",
        help="Skip semantic layer textboxes (don't add text)."
    )
    parser.add_argument(
        "--placeholders", type=str,
        help="Path to image_placeholders.json for chart detection."
    )
    parser.add_argument(
        "--config", type=str,
        help="Path to config.json for API settings (chart generation)."
    )
    parser.add_argument(
        "--skip-charts", action="store_true",
        help="Skip chart generation (add images instead)."
    )
    
    args = parser.parse_args()
    
    input_path = Path(args.input)
    
    # Collect SVG files
    if input_path.is_dir():
        all_svgs = list(input_path.glob("*.SVG")) + list(input_path.glob("*.svg"))
        svg_paths = sorted(
            [p for p in all_svgs if not p.name.startswith("._")],
            key=natural_sort_key
        )
    else:
        svg_paths = [input_path]
    
    if not svg_paths:
        print(f"No SVG files found in: {input_path}")
        return
    
    print(f"Found {len(svg_paths)} SVG file(s)")
    
    # Load placeholders and config
    placeholders_map = None
    chart_config = None
    
    if args.placeholders:
        placeholders_path = Path(args.placeholders)
        placeholders_map = load_placeholders(placeholders_path)
        if placeholders_map:
            print(f"Loaded {len(placeholders_map)} placeholder entries")
    
    if args.config:
        config_path = Path(args.config)
        chart_config = load_config(config_path)
        if chart_config:
            print(f"Loaded config from {config_path}")
    
    # Skip elements - if skip_textboxes is set, we skip the semantic layer entirely
    skip_elements = set()
    if args.skip_textboxes:
        skip_elements.add("semantic-layer")
    
    # Build PPTX
    build_pptx_pro(
        svg_paths,
        Path(args.output),
        dpi=args.dpi,
        cjk_font=args.cjk_font,
        skip_elements=skip_elements if skip_elements else None,
        line_tol=args.line_tol,
        box_pad=args.box_pad,
        placeholders_map=placeholders_map,
        chart_config=chart_config,
        skip_charts=args.skip_charts,
    )


if __name__ == "__main__":
    main()
